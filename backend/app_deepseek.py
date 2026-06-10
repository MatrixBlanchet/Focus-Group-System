from flask import Flask, request, jsonify, session, make_response, Response, stream_with_context, has_request_context, redirect, g, copy_current_request_context
from flask_sqlalchemy import SQLAlchemy
from flask_cors import CORS
from sqlalchemy import or_
import requests
import json
import time
import random
import hashlib
import hmac
import uuid
import re
import secrets
import string
import smtplib
import threading
import base64
from queue import Queue
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime, timedelta
from email.mime.text import MIMEText
from email.header import Header
from email.utils import formataddr

import os
basedir = os.path.abspath(os.path.dirname(__file__))
REPORT_EVENT_LOG_PATH = os.path.join(basedir, 'logs', 'report_generation.jsonl')
_report_event_log_lock = threading.Lock()

from io import BytesIO
from urllib.parse import quote
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import platform
from docx import Document

# ---------------------------
# Flask应用初始化配置
# ---------------------------
app = Flask(__name__, static_folder='static')
# 数据库连接配置：SQLite数据库
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///' + os.path.join(basedir, 'focus_group.db')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False  # 禁用修改跟踪，提升性能
app.config['JSON_AS_ASCII'] = False  # 支持中文JSON输出
app.config['JSONIFY_MIMETYPE'] = 'application/json; charset=utf-8'
app.secret_key = 'your-secret-key-here-make-it-very-long-and-secret-for-production'  # Session密钥，生产环境需修改
db = SQLAlchemy(app)  # 初始化SQLAlchemy实例
CORS(app)  # 启用跨域资源共享

# ---------------------------
# 配置文件加载（DeepSeek API密钥）
# ---------------------------
DEEPSEEK_API_KEY = None
DEEPSEEK_API_KEYS = []
DEEPSEEK_BASE_URL = None
MODEL = None
USER_AI_CONFIG_SECRET = None


def get_report_request_id():
    if not has_request_context():
        return None
    request_id = getattr(g, 'report_request_id', None)
    if not request_id:
        request_id = uuid.uuid4().hex[:12]
        g.report_request_id = request_id
    return request_id


def extract_error_message_from_response(response_obj):
    if response_obj is None:
        return ''
    try:
        payload = json.loads(response_obj.get_data(as_text=True) or '{}')
        if isinstance(payload, dict):
            return str(payload.get('error') or payload.get('message') or '').strip()
    except Exception:
        pass
    return getattr(response_obj, 'status', '') or ''


def make_log_safe(value):
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, datetime):
        return value.isoformat()
    if isinstance(value, dict):
        return {str(key): make_log_safe(item) for key, item in value.items()}
    if isinstance(value, (list, tuple, set)):
        return [make_log_safe(item) for item in value]
    return str(value)


def log_report_event(event, scenario_id=None, level='INFO', **fields):
    payload = {
        'timestamp': datetime.now().isoformat(),
        'level': level,
        'event': event,
        'scenario_id': scenario_id,
    }
    if has_request_context():
        payload.update({
            'request_id': get_report_request_id(),
            'path': request.path,
            'method': request.method,
            'remote_addr': request.headers.get('X-Forwarded-For', request.remote_addr),
            'user_id': session.get('user_id'),
        })
    for key, value in fields.items():
        if value is not None:
            payload[key] = make_log_safe(value)

    try:
        os.makedirs(os.path.dirname(REPORT_EVENT_LOG_PATH), exist_ok=True)
        line = json.dumps(payload, ensure_ascii=False)
        with _report_event_log_lock:
            with open(REPORT_EVENT_LOG_PATH, 'a', encoding='utf-8') as log_file:
                log_file.write(line + '\n')
    except Exception as exc:
        print(f"[ReportLog] Failed to write report event: {exc}")

# 验证码配置
VERIFICATION_CHANNEL = "test"   # test / email / sms / auto
VERIFICATION_TTL = 600
VERIFICATION_COOLDOWN = 60
VERIFICATION_MAX_ATTEMPTS = 3

EMAIL_SMTP_HOST = ""
EMAIL_SMTP_PORT = 465
EMAIL_SMTP_USE_SSL = True
EMAIL_SMTP_USER = ""
EMAIL_SMTP_PASSWORD = ""
EMAIL_FROM_NAME = "焦点小组系统"

SMS_PROVIDER = "aliyun"
SMS_ACCESS_KEY_ID = ""
SMS_ACCESS_KEY_SECRET = ""
SMS_SIGN_NAME = ""
SMS_TEMPLATE_CODE = ""

# 登录速率限制（每60秒最多5次）
LOGIN_RATE_LIMIT = 5
LOGIN_RATE_WINDOW = 60

with open(os.path.join(basedir, 'config.txt'), 'r', encoding='utf-8') as f:
    for line in f:
        if line.startswith('DEEPSEEK_API_KEY='):
            DEEPSEEK_API_KEY = line.split('=', 1)[1].strip()
        elif line.startswith('DEEPSEEK_API_KEYS='):
            DEEPSEEK_API_KEYS = [item.strip() for item in line.split('=', 1)[1].split(',') if item.strip()]
        elif line.startswith('DEEPSEEK_BASE_URL='):
            DEEPSEEK_BASE_URL = line.split('=', 1)[1].strip()
        elif line.startswith('MODEL='):
            MODEL = line.split('=', 1)[1].strip()
        elif line.startswith('USER_AI_CONFIG_SECRET='):
            USER_AI_CONFIG_SECRET = line.split('=', 1)[1].strip()
        elif line.startswith('VERIFICATION_CHANNEL='):
            VERIFICATION_CHANNEL = line.split('=', 1)[1].strip()
        elif line.startswith('VERIFICATION_TTL='):
            VERIFICATION_TTL = int(line.split('=', 1)[1].strip())
        elif line.startswith('VERIFICATION_COOLDOWN='):
            VERIFICATION_COOLDOWN = int(line.split('=', 1)[1].strip())
        elif line.startswith('VERIFICATION_MAX_ATTEMPTS='):
            VERIFICATION_MAX_ATTEMPTS = int(line.split('=', 1)[1].strip())
        elif line.startswith('EMAIL_SMTP_HOST='):
            EMAIL_SMTP_HOST = line.split('=', 1)[1].strip()
        elif line.startswith('EMAIL_SMTP_PORT='):
            EMAIL_SMTP_PORT = int(line.split('=', 1)[1].strip())
        elif line.startswith('EMAIL_SMTP_USE_SSL='):
            EMAIL_SMTP_USE_SSL = line.split('=', 1)[1].strip().lower() == 'true'
        elif line.startswith('EMAIL_SMTP_USER='):
            EMAIL_SMTP_USER = line.split('=', 1)[1].strip()
        elif line.startswith('EMAIL_SMTP_PASSWORD='):
            EMAIL_SMTP_PASSWORD = line.split('=', 1)[1].strip()
        elif line.startswith('EMAIL_FROM_NAME='):
            EMAIL_FROM_NAME = line.split('=', 1)[1].strip()
        elif line.startswith('SMS_ACCESS_KEY_ID='):
            SMS_ACCESS_KEY_ID = line.split('=', 1)[1].strip()
        elif line.startswith('SMS_ACCESS_KEY_SECRET='):
            SMS_ACCESS_KEY_SECRET = line.split('=', 1)[1].strip()
        elif line.startswith('SMS_SIGN_NAME='):
            SMS_SIGN_NAME = line.split('=', 1)[1].strip()
        elif line.startswith('SMS_TEMPLATE_CODE='):
            SMS_TEMPLATE_CODE = line.split('=', 1)[1].strip()

if not DEEPSEEK_API_KEYS and DEEPSEEK_API_KEY:
    DEEPSEEK_API_KEYS = [DEEPSEEK_API_KEY]

def safe_json_loads(value, fallback):
    if value in (None, ''):
        return fallback
    try:
        return json.loads(value)
    except Exception:
        return fallback

def json_dumps(value, default):
    source = value if value is not None else default
    return json.dumps(source, ensure_ascii=False)

def normalize_text_list(value):
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    if isinstance(value, str):
        return [item.strip() for item in re.split(r'[,\n，；;]', value) if item.strip()]
    return []

_ai_key_rotation_lock = threading.Lock()
_ai_key_rotation_state = {}


def normalize_api_keys_input(value):
    keys = []
    if isinstance(value, list):
        for item in value:
            keys.extend(normalize_text_list(item))
    else:
        keys = normalize_text_list(value)

    unique_keys = []
    seen = set()
    for key in keys:
        if key not in seen:
            unique_keys.append(key)
            seen.add(key)
    return unique_keys


def mask_api_key(key):
    text = str(key or '').strip()
    if not text:
        return ''
    if len(text) <= 8:
        return text[:2] + '*' * max(len(text) - 4, 1) + text[-2:]
    return text[:4] + '*' * (len(text) - 8) + text[-4:]


def build_chat_completions_url(base_url):
    cleaned = str(base_url or '').strip().rstrip('/')
    if not cleaned:
        return ''
    if cleaned.endswith('/chat/completions'):
        return cleaned
    return f"{cleaned}/chat/completions"


def _require_user_ai_secret(user=None, user_id=None):
    explicit_secret = (USER_AI_CONFIG_SECRET or '').strip()
    if explicit_secret:
        return explicit_secret.encode('utf-8')

    effective_user_id = user_id
    if effective_user_id is None and user is not None:
        effective_user_id = getattr(user, 'id', None)

    app_secret = str(getattr(app, 'secret_key', '') or '').strip()
    if not app_secret:
        raise ValueError('app.secret_key is not configured')

    # Derive a stable per-user secret from the app secret and user id.
    scope = f"user-ai:{effective_user_id or 'global'}".encode('utf-8')
    return hmac.new(app_secret.encode('utf-8'), scope, hashlib.sha256).digest()


def _xor_with_derived_stream(data_bytes, secret_bytes, nonce_bytes):
    stream = bytearray()
    counter = 0
    while len(stream) < len(data_bytes):
        block = hashlib.sha256(
            secret_bytes + nonce_bytes + counter.to_bytes(4, 'big')
        ).digest()
        stream.extend(block)
        counter += 1
    return bytes(a ^ b for a, b in zip(data_bytes, stream[:len(data_bytes)]))


def encrypt_user_api_keys(api_keys, user=None, user_id=None):
    normalized_keys = normalize_api_keys_input(api_keys)
    secret_bytes = _require_user_ai_secret(user=user, user_id=user_id)
    nonce_bytes = secrets.token_bytes(16)
    payload_bytes = json.dumps(normalized_keys, ensure_ascii=False).encode('utf-8')
    cipher_bytes = _xor_with_derived_stream(payload_bytes, secret_bytes, nonce_bytes)
    signature = hmac.new(secret_bytes, nonce_bytes + cipher_bytes, hashlib.sha256).digest()
    return base64.urlsafe_b64encode(nonce_bytes + signature + cipher_bytes).decode('ascii')


def decrypt_user_api_keys(encrypted_value, user=None, user_id=None):
    if not encrypted_value:
        return []

    secret_bytes = _require_user_ai_secret(user=user, user_id=user_id)
    raw_bytes = base64.urlsafe_b64decode(str(encrypted_value).encode('ascii'))
    nonce_bytes = raw_bytes[:16]
    signature = raw_bytes[16:48]
    cipher_bytes = raw_bytes[48:]
    expected_signature = hmac.new(secret_bytes, nonce_bytes + cipher_bytes, hashlib.sha256).digest()
    if not hmac.compare_digest(signature, expected_signature):
        raise ValueError('Invalid encrypted AI key payload')

    payload_bytes = _xor_with_derived_stream(cipher_bytes, secret_bytes, nonce_bytes)
    return normalize_api_keys_input(json.loads(payload_bytes.decode('utf-8')))


def _build_runtime_ai_request(endpoint_url, model_name, prompt, system_prompt, max_tokens):
    return {
        "model": model_name,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.7,
        "max_tokens": max_tokens
    }


def _pick_api_keys_in_order(pool_name, api_keys):
    cleaned_keys = normalize_api_keys_input(api_keys)
    if len(cleaned_keys) <= 1:
        return cleaned_keys

    with _ai_key_rotation_lock:
        start_index = _ai_key_rotation_state.get(pool_name, 0) % len(cleaned_keys)
        _ai_key_rotation_state[pool_name] = (start_index + 1) % len(cleaned_keys)

    return cleaned_keys[start_index:] + cleaned_keys[:start_index]


def _post_ai_completion(endpoint_url, api_key, model_name, prompt, system_prompt, max_tokens=3500, timeout=120):
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }
    data = _build_runtime_ai_request(endpoint_url, model_name, prompt, system_prompt, max_tokens)
    return requests.post(
        endpoint_url,
        headers=headers,
        json=data,
        timeout=timeout
    )


def get_session_user():
    if not has_request_context():
        return None
    user_id = session.get('user_id')
    if not user_id:
        return None
    return db.session.get(User, user_id)


def get_user_ai_config_summary(user):
    masked_keys = []
    try:
        masked_keys = [
            mask_api_key(key)
            for key in decrypt_user_api_keys(getattr(user, 'ai_api_keys_encrypted', ''), user=user)
        ]
    except Exception:
        masked_keys = []

    has_config = bool(
        getattr(user, 'ai_endpoint_url', '') and
        getattr(user, 'ai_model_name', '') and
        getattr(user, 'ai_api_keys_encrypted', '')
    )

    return {
        "endpoint_url": getattr(user, 'ai_endpoint_url', '') or '',
        "model_name": getattr(user, 'ai_model_name', '') or '',
        "has_config": has_config,
        "enabled": bool(getattr(user, 'ai_config_enabled', False) and has_config),
        "key_count": len(masked_keys),
        "masked_keys": masked_keys,
        "last_test_status": getattr(user, 'ai_last_test_status', '') or '',
        "last_test_message": getattr(user, 'ai_last_test_message', '') or '',
        "last_tested_at": user.ai_last_tested_at.isoformat() if getattr(user, 'ai_last_tested_at', None) else ''
    }


def build_system_ai_runtime_config():
    api_keys = normalize_api_keys_input(DEEPSEEK_API_KEYS or [DEEPSEEK_API_KEY])
    endpoint_url = build_chat_completions_url(DEEPSEEK_BASE_URL)
    if not endpoint_url or not MODEL or not api_keys:
        return None
    return {
        "name": "system",
        "endpoint_url": endpoint_url,
        "model_name": MODEL,
        "api_keys": api_keys
    }


def build_user_ai_runtime_config(user):
    if not user or not getattr(user, 'ai_config_enabled', False):
        return None

    endpoint_url = (getattr(user, 'ai_endpoint_url', '') or '').strip()
    model_name = (getattr(user, 'ai_model_name', '') or '').strip()
    encrypted_keys = getattr(user, 'ai_api_keys_encrypted', '') or ''
    if not endpoint_url or not model_name or not encrypted_keys:
        return None

    api_keys = decrypt_user_api_keys(encrypted_keys, user=user)
    if not api_keys:
        return None

    return {
        "name": f"user:{user.id}",
        "endpoint_url": endpoint_url,
        "model_name": model_name,
        "api_keys": api_keys
    }


def resolve_runtime_ai_configs():
    configs = []
    user = get_session_user()
    if user:
        try:
            user_config = build_user_ai_runtime_config(user)
        except Exception:
            user_config = None
        if user_config:
            configs.append(user_config)

    system_config = build_system_ai_runtime_config()
    if system_config:
        configs.append(system_config)
    return configs


def _call_ai_with_config(config, prompt, system_prompt, max_retries=3, max_tokens=3500):
    retry_delay = 3
    last_error = None

    for api_key in _pick_api_keys_in_order(config["name"], config["api_keys"]):
        for attempt in range(max_retries):
            try:
                response = _post_ai_completion(
                    config["endpoint_url"],
                    api_key,
                    config["model_name"],
                    prompt,
                    system_prompt,
                    max_tokens=max_tokens,
                    timeout=120
                )
            except requests.exceptions.RequestException as e:
                last_error = f"Request Error: {str(e)}"
                if attempt < max_retries - 1:
                    time.sleep(retry_delay * (attempt + 1))
                    continue
                break

            if response.status_code == 200:
                return response.json()["choices"][0]["message"]["content"], None

            if response.status_code in [429, 500, 502, 503, 504]:
                last_error = f"API Error {response.status_code}: 服务暂时不可用，正在重试..."
                if attempt < max_retries - 1:
                    time.sleep(retry_delay * (attempt + 1))
                    continue
            else:
                last_error = f"API Error: {response.status_code} - {response.text}"
            break

    return None, last_error or "请求失败，请稍后重试"


def validate_user_ai_config_payload(endpoint_url, model_name, api_keys):
    endpoint = str(endpoint_url or '').strip()
    model = str(model_name or '').strip()
    keys = normalize_api_keys_input(api_keys)

    if not endpoint:
        return None, None, None, "请填写完整接口地址"
    if not model:
        return None, None, None, "请填写模型名称"
    if not keys:
        return None, None, None, "请至少填写一个 API Key"
    if not (endpoint.startswith('http://') or endpoint.startswith('https://')):
        return None, None, None, "接口地址必须以 http:// 或 https:// 开头"

    return endpoint, model, keys, None


def test_user_ai_config(endpoint_url, model_name, api_keys):
    endpoint, model, keys, error_message = validate_user_ai_config_payload(
        endpoint_url,
        model_name,
        api_keys
    )
    if error_message:
        return False, {
            "status": "error",
            "message": error_message,
            "results": [],
            "valid_key_count": 0,
            "invalid_key_count": 0
        }, []

    def probe(single_key):
        try:
            response = _post_ai_completion(
                endpoint,
                single_key,
                model,
                "Reply with OK only.",
                "You are a connectivity check assistant.",
                max_tokens=16,
                timeout=30
            )
            if response.status_code == 200:
                return {
                    "masked_key": mask_api_key(single_key),
                    "is_valid": True,
                    "message": "连接成功"
                }
            return {
                "masked_key": mask_api_key(single_key),
                "is_valid": False,
                "message": f"HTTP {response.status_code}"
            }
        except requests.exceptions.RequestException as e:
            return {
                "masked_key": mask_api_key(single_key),
                "is_valid": False,
                "message": str(e)
            }

    max_workers = min(len(keys), 4) or 1
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        results = list(executor.map(probe, keys))

    valid_keys = [key for key, item in zip(keys, results) if item["is_valid"]]
    invalid_count = len(results) - len(valid_keys)
    success = bool(valid_keys)
    message = "至少 1 个 Key 可用" if success else "没有可用的 API Key"
    return success, {
        "status": "success" if success else "error",
        "message": message,
        "results": results,
        "valid_key_count": len(valid_keys),
        "invalid_key_count": invalid_count
    }, valid_keys


def normalize_custom_params(value):
    if isinstance(value, dict):
        return value
    if isinstance(value, str):
        text = value.strip()
        if not text:
            return {}
        parsed = safe_json_loads(text, None)
        if isinstance(parsed, dict):
            return parsed
        return {"notes": text}
    return {}

def sanitize_participant_seed(seed, default_name=None):
    raw = seed if isinstance(seed, dict) else {}
    return {
        'persona_name': str(raw.get('persona_name') or raw.get('name') or default_name or '').strip(),
        'persona_tags': normalize_text_list(raw.get('persona_tags', raw.get('tags', []))),
        'personality': str(raw.get('personality') or '').strip(),
        'background': str(raw.get('background') or '').strip(),
        'usage_goal': str(raw.get('usage_goal') or '').strip(),
        'budget_sensitivity': str(raw.get('budget_sensitivity') or '').strip(),
        'brand_preference': str(raw.get('brand_preference') or '').strip(),
        'risk_aversion': str(raw.get('risk_aversion') or '').strip(),
        'decision_style': str(raw.get('decision_style') or '').strip(),
        'deal_breakers': normalize_text_list(raw.get('deal_breakers', [])),
        'stance_summary': str(raw.get('stance_summary') or '').strip(),
        'custom_params': normalize_custom_params(raw.get('custom_params', {}))
    }

def build_participant_stance_state(stance_summary=''):
    return json_dumps({
        "current_position": stance_summary or '',
        "confidence": "中",
        "last_updated_round": 0
    }, {})

def parse_participant_payload(data, require_name=False, existing_custom_params=None):
    payload = {
        'persona_name': str(data.get('persona_name') or data.get('name') or '').strip(),
        'persona_tags': normalize_text_list(data.get('persona_tags', data.get('tags', []))),
        'personality': str(data.get('personality') or '').strip(),
        'background': str(data.get('background') or '').strip(),
        'usage_goal': str(data.get('usage_goal') or '').strip(),
        'budget_sensitivity': str(data.get('budget_sensitivity') or '').strip(),
        'brand_preference': str(data.get('brand_preference') or '').strip(),
        'risk_aversion': str(data.get('risk_aversion') or '').strip(),
        'decision_style': str(data.get('decision_style') or '').strip(),
        'deal_breakers': normalize_text_list(data.get('deal_breakers', [])),
        'stance_summary': str(data.get('stance_summary') or '').strip()
    }
    if require_name and not payload['persona_name']:
        return None, "参与者姓名不能为空"

    custom_params = existing_custom_params.copy() if isinstance(existing_custom_params, dict) else {}
    if 'custom_params' in data:
        custom_params = normalize_custom_params(data.get('custom_params'))
    payload['custom_params'] = custom_params
    return payload, None

def build_legacy_profile_constraints(profile):
    if not isinstance(profile, dict):
        return []

    labels = {
        'professional_level': '专业级别',
        'income_level': '收入水平',
        'industry': '所在行业',
        'age_range': '年龄段',
        'education': '教育背景',
        'city_tier': '所在城市',
        'budget_sensitivity': '预算敏感度',
        'risk_aversion': '风险偏好',
        'decision_style': '决策风格',
        'brand_preference': '品牌偏好'
    }
    constraints = []
    for key, label in labels.items():
        value = str(profile.get(key) or '').strip()
        if value:
            constraints.append(f"{label}：{value}")

    needs = normalize_text_list(profile.get('needs', []))
    if needs:
        constraints.append(f"核心需求：{'、'.join(needs)}")

    deal_breakers = normalize_text_list(profile.get('deal_breakers', []))
    if deal_breakers:
        constraints.append(f"反感点：{'、'.join(deal_breakers)}")

    if profile.get('diversity_mix'):
        constraints.append("不同参与者之间需要体现明显差异化")
    return constraints

def build_generated_participant_prompt(scenario, role_type, count, persona_instructions, seed_participants, legacy_profile):
    selling_points = safe_json_loads(scenario.core_selling_points, [])
    topics = safe_json_loads(scenario.discussion_topics, [])
    assumptions = normalize_text_list(safe_json_loads(scenario.validation_assumptions, scenario.validation_assumptions or []))
    seed_lines = []
    for index, seed in enumerate(seed_participants, start=1):
        seed_lines.append(
            f"{index}. 姓名={seed.get('persona_name') or '待AI命名'}；"
            f"标签={'/'.join(seed.get('persona_tags', [])) or '待补充'}；"
            f"背景={seed.get('background') or '待补充'}；"
            f"性格={seed.get('personality') or '待补充'}；"
            f"使用目标={seed.get('usage_goal') or '待补充'}；"
            f"预算敏感度={seed.get('budget_sensitivity') or '待补充'}；"
            f"品牌偏好={seed.get('brand_preference') or '待补充'}；"
            f"风险厌恶={seed.get('risk_aversion') or '待补充'}；"
            f"决策风格={seed.get('decision_style') or '待补充'}；"
            f"反感点={'、'.join(seed.get('deal_breakers', [])) or '待补充'}；"
            f"稳定立场={seed.get('stance_summary') or '待补充'}"
        )

    instructions = []
    if persona_instructions:
        instructions.append(persona_instructions)
    instructions.extend(build_legacy_profile_constraints(legacy_profile))
    instructions_text = "\n".join([f"- {item}" for item in instructions]) if instructions else "- 无额外约束，按场景生成差异化参与者"
    seed_text = "\n".join(seed_lines) if seed_lines else "无。请直接生成完整参与者列表。"

    prompt = f"""请为以下企业研究场景生成 {count} 个{role_type}参与者画像，用于模拟讨论。

产品名称：{scenario.product_name}
产品概念：{scenario.product_concept}
核心卖点：{'、'.join(selling_points) if selling_points else '暂无'}
研究目标：{scenario.research_goal or '判断该概念是否值得继续推进'}
核心决策问题：{scenario.decision_problem or '当前阶段是否值得继续验证'}
目标用户：{scenario.target_user_profile or '暂无'}
竞品/替代方案：{scenario.competitor_context or '暂无'}
讨论主题：{'、'.join(topics) if topics else '暂无'}
待验证假设：{'、'.join(assumptions) if assumptions else '暂无'}
场景类型：{scenario.occasion_description or scenario.occasion_type}

用户额外要求：
{instructions_text}

已提供的种子参与者：
{seed_text}

要求：
1. 输出必须是 JSON 数组，不要输出解释文字或 Markdown。
2. 数组长度固定为 {count}。
3. 如果已提供种子参与者，优先保留其核心设定，并补全缺失字段。
4. 不同参与者之间要有明显差异，避免只是措辞不同。
5. 每个对象都必须包含：name, tags, personality, background, usage_goal, budget_sensitivity, brand_preference, risk_aversion, decision_style, deal_breakers, stance_summary。
6. tags 和 deal_breakers 输出数组，其余字段输出字符串。
7. 内容要贴近真实用户或真实岗位，适合围绕该产品进行讨论。"""

    system_prompt = (
        "你是一名企业研究顾问，擅长为产品讨论会构造真实、差异化、可持续扮演的参与者画像。"
        "你必须严格输出合法 JSON。"
    )
    return prompt, system_prompt

def build_fallback_personas(seed_participants, count):
    personas = []
    for index in range(count):
        seed = seed_participants[index] if index < len(seed_participants) else {}
        personas.append({
            "name": seed.get('persona_name') or f"参与者{index + 1}",
            "tags": seed.get('persona_tags', []) or ["普通用户"],
            "personality": seed.get('personality') or "表达直接，愿意分享真实使用顾虑",
            "background": seed.get('background') or "具备相关使用场景，但仍在权衡是否值得尝试",
            "usage_goal": seed.get('usage_goal') or "确认产品是否能真正解决自己的问题",
            "budget_sensitivity": seed.get('budget_sensitivity') or "中等，愿意为明确价值付费",
            "brand_preference": seed.get('brand_preference') or "更看重体验可信度，不会只因概念新鲜买单",
            "risk_aversion": seed.get('risk_aversion') or "中等偏高，担心功能不成熟",
            "decision_style": seed.get('decision_style') or "会比较收益和风险后再决定",
            "deal_breakers": seed.get('deal_breakers', []) or ["体验复杂", "价值不清晰"],
            "stance_summary": seed.get('stance_summary') or "对产品有兴趣，但需要更多证据支持"
        })
    return personas

def merge_generated_personas(personas, seed_participants, count):
    merged = []
    for index in range(count):
        seed = seed_participants[index] if index < len(seed_participants) else {}
        generated = personas[index] if index < len(personas) and isinstance(personas[index], dict) else {}
        merged.append({
            'persona_name': str(generated.get('name') or generated.get('persona_name') or seed.get('persona_name') or f"参与者{index + 1}").strip(),
            'persona_tags': normalize_text_list(generated.get('tags', generated.get('persona_tags', seed.get('persona_tags', [])))),
            'personality': str(generated.get('personality') or seed.get('personality') or '').strip(),
            'background': str(generated.get('background') or seed.get('background') or '').strip(),
            'usage_goal': str(generated.get('usage_goal') or seed.get('usage_goal') or '').strip(),
            'budget_sensitivity': str(generated.get('budget_sensitivity') or seed.get('budget_sensitivity') or '').strip(),
            'brand_preference': str(generated.get('brand_preference') or seed.get('brand_preference') or '').strip(),
            'risk_aversion': str(generated.get('risk_aversion') or seed.get('risk_aversion') or '').strip(),
            'decision_style': str(generated.get('decision_style') or seed.get('decision_style') or '').strip(),
            'deal_breakers': normalize_text_list(generated.get('deal_breakers', seed.get('deal_breakers', []))),
            'stance_summary': str(generated.get('stance_summary') or seed.get('stance_summary') or '').strip(),
            'custom_params': seed.get('custom_params', {})
        })
    return merged

def normalize_strength_level(value):
    text = (value or '').strip()
    if text in {'高', '中', '低'}:
        return text
    if '高' in text:
        return '高'
    if '低' in text:
        return '低'
    return '中'

# ======================== 密码安全工具函数 ========================

EMAIL_RE = re.compile(r"^[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}$")

def hash_password(password: str) -> str:
    """使用 pbkdf2_hmac 加密密码，带随机 salt"""
    salt = secrets.token_hex(16)
    dk = hashlib.pbkdf2_hmac('sha256', password.encode('utf-8'), salt.encode('utf-8'), 120000)
    return f"pbkdf2_sha256$120000${salt}${dk.hex()}"

def verify_password(password: str, stored: str) -> bool:
    """验证密码，兼容旧的 SHA256 格式"""
    try:
        if stored.startswith("pbkdf2_sha256$"):
            _, iterations, salt, hex_hash = stored.split('$')
            dk = hashlib.pbkdf2_hmac('sha256', password.encode('utf-8'), salt.encode('utf-8'), int(iterations))
            return dk.hex() == hex_hash
        return hashlib.sha256(password.encode()).hexdigest() == stored
    except Exception:
        return False

def validate_password(pwd: str):
    """验证密码强度"""
    if len(pwd) < 8:
        return "密码至少 8 位"
    if not re.search(r"[A-Za-z]", pwd):
        return "密码必须包含字母"
    if not re.search(r"[0-9]", pwd):
        return "密码必须包含数字"
    return None

def validate_email(email: str):
    """验证邮箱格式"""
    return EMAIL_RE.match(email) is not None


def json_api_response(payload, status=200):
    return app.response_class(
        response=json.dumps(payload, ensure_ascii=False),
        status=status,
        mimetype='application/json; charset=utf-8'
    )


def json_api_error(message, status=400):
    return json_api_response({"error": message}, status=status)


def normalize_auth_target(target):
    return (target or '').strip().lower()


def is_email_target(target):
    return validate_email(normalize_auth_target(target))


def find_user_by_email(email):
    normalized = normalize_auth_target(email)
    if not normalized or not validate_email(normalized):
        return None
    return User.query.filter_by(email=normalized).first()


def apply_login_session(user, update_stats=True):
    if update_stats:
        user.last_login = datetime.now()
        user.login_count = (user.login_count or 0) + 1
        db.session.commit()

    session['user_id'] = user.id
    session['username'] = user.username


def register_user_account(username, email, password):
    normalized_email = normalize_auth_target(email)
    normalized_username = (username or '').strip()

    if not normalized_username or not normalized_email or not password:
        return None, "请填写完整信息", 400
    if not validate_email(normalized_email):
        return None, "邮箱格式不正确", 400

    pwd_err = validate_password(password)
    if pwd_err:
        return None, pwd_err, 400

    if User.query.filter_by(username=normalized_username).first():
        return None, "用户名已存在", 400
    if User.query.filter_by(email=normalized_email).first():
        return None, "邮箱已被注册", 400

    user = User(
        username=normalized_username,
        email=normalized_email,
        password_hash=hash_password(password),
        nickname=normalized_username
    )
    db.session.add(user)
    db.session.commit()
    return user, None, 201


def build_auto_username_from_email(email):
    local_part = normalize_auth_target(email).split('@', 1)[0]
    base = re.sub(r'[^a-z0-9_]+', '_', local_part.lower()).strip('_') or 'user'
    candidate = base[:32]
    suffix = 1
    while User.query.filter_by(username=candidate).first():
        candidate = f"{base[:24]}_{suffix}"
        suffix += 1
    return candidate


def create_code_login_user(email):
    username = build_auto_username_from_email(email)
    user = User(
        username=username,
        email=normalize_auth_target(email),
        password_hash=hash_password(secrets.token_urlsafe(24)),
        nickname=username
    )
    db.session.add(user)
    db.session.commit()
    return user


def generate_meeting_room_code(length=6):
    alphabet = "ABCDEFGHJKLMNPQRSTUVWXYZ23456789"
    while True:
        code = ''.join(secrets.choice(alphabet) for _ in range(length))
        if not MeetingRoom.query.filter_by(room_code=code).first():
            return code


def generate_meeting_invite_token():
    while True:
        token = secrets.token_urlsafe(18)
        if not MeetingRoom.query.filter_by(invite_token=token).first():
            return token


def build_room_invite_link(invite_token):
    token = (invite_token or '').strip()
    if not token:
        return ''
    base_url = request.host_url.rstrip('/') if has_request_context() else 'http://127.0.0.1:5000'
    return f"{base_url}/meeting-room?room_invite={token}"


def get_room_member_display_name(user):
    if not user:
        return '成员'
    return (user.nickname or user.username or user.email or '成员').strip()


def get_current_user_model():
    user_id = session.get('user_id')
    if not user_id:
        return None
    return db.session.get(User, user_id)


def get_meeting_room_or_error(room_id, require_membership=True, allow_owner=True):
    user = get_current_user_model()
    if not user:
        return None, None, json_api_error("请先登录", 401)

    room = db.session.get(MeetingRoom, room_id)
    if not room:
        return None, None, json_api_error("会议室不存在", 404)

    membership = MeetingRoomMember.query.filter_by(room_id=room.id, user_id=user.id).first()
    if require_membership and not membership and not (allow_owner and room.owner_user_id == user.id):
        return None, None, json_api_error("无权访问该会议室", 403)

    if not membership and allow_owner and room.owner_user_id == user.id:
        membership = MeetingRoomMember.query.filter_by(room_id=room.id, user_id=user.id).first()

    return room, membership, None


def serialize_meeting_room_bundle(room, scenario=None, include_messages=False):
    scenario = scenario or db.session.get(ProductScenario, room.scenario_id)
    members = MeetingRoomMember.query.filter_by(room_id=room.id).order_by(MeetingRoomMember.joined_at.asc()).all()
    payload = {
        'room': room.to_dict(),
        'scenario': scenario.to_dict() if scenario else None,
        'members': [member.to_dict() for member in members]
    }
    if include_messages:
        messages = MeetingRoomMessage.query.filter_by(room_id=room.id).order_by(MeetingRoomMessage.created_at.asc()).all()
        payload['messages'] = [message.to_dict() for message in messages]
    return payload


def create_meeting_room_system_message(room_id, content):
    message = MeetingRoomMessage(
        room_id=room_id,
        user_id=None,
        sender_name='系统通知',
        content=(content or '').strip(),
        message_type='system'
    )
    db.session.add(message)
    return message


def build_meeting_scenario_from_payload(user_id, data):
    discussion_topics = normalize_text_list(data.get('discussion_topics', []))
    selling_points = normalize_text_list(data.get('core_selling_points', []))
    validation_assumptions = normalize_text_list(data.get('validation_assumptions', []))
    topic_title = str(data.get('topic_title') or data.get('decision_problem') or data.get('room_name') or '').strip()

    scenario = ProductScenario(
        user_id=user_id,
        product_name=(data.get('product_name') or topic_title or '会议议题').strip(),
        product_concept=(data.get('product_concept') or data.get('topic_notes') or '待补充产品概念').strip(),
        core_selling_points=json_dumps(selling_points, []),
        discussion_topics=json_dumps(discussion_topics, []),
        occasion_type=data.get('occasion_type', 'focus_group'),
        occasion_description=(data.get('occasion_description') or '标准焦点小组讨论').strip(),
        research_goal=(data.get('research_goal') or '').strip(),
        decision_problem=(data.get('decision_problem') or topic_title).strip(),
        target_user_profile=(data.get('target_user_profile') or '').strip(),
        competitor_context=(data.get('competitor_context') or '').strip(),
        validation_assumptions=json_dumps(validation_assumptions, []),
        research_plan_status='pending',
        meeting_status='meeting_waiting'
    )
    scenario.research_plan = build_research_plan(scenario)
    return scenario


def normalize_meeting_status_value(value, fallback='standalone'):
    normalized = (value or '').strip().lower()
    mapping = {
        '': fallback,
        'standalone': 'standalone',
        'waiting': 'meeting_waiting',
        'active': 'meeting_active',
        'ended': 'meeting_ended',
        'meeting_waiting': 'meeting_waiting',
        'meeting_active': 'meeting_active',
        'meeting_ended': 'meeting_ended',
    }
    return mapping.get(normalized, fallback)


def sync_scenario_meeting_status(scenario, room=None):
    if not scenario:
        return None
    if room is None:
        room = MeetingRoom.query.filter_by(scenario_id=scenario.id).first()
    if room:
        room_status = normalize_meeting_status_value(room.status, 'meeting_waiting')
        scenario.meeting_status = room_status
    else:
        scenario.meeting_status = normalize_meeting_status_value(scenario.meeting_status, 'standalone')
    return scenario.meeting_status


def is_standalone_scenario(scenario):
    return normalize_meeting_status_value(getattr(scenario, 'meeting_status', None), 'standalone') == 'standalone'


def get_standalone_participants(scenario_id):
    return (
        VirtualParticipant.query
        .filter(VirtualParticipant.scenario_id == scenario_id)
        .filter(or_(VirtualParticipant.room_managed.is_(False), VirtualParticipant.room_managed.is_(None)))
        .order_by(VirtualParticipant.id.asc())
        .all()
    )


def reject_if_not_standalone(scenario, message):
    if is_standalone_scenario(scenario):
        return None
    return json_api_error(message, 409)


def upsert_room_member(room, user, member_role='member'):
    member = MeetingRoomMember.query.filter_by(room_id=room.id, user_id=user.id).first()
    if member:
        member.member_role = member.member_role or member_role
        member.display_name = get_room_member_display_name(user)
        member.status = 'arrived'
        member.last_seen_at = datetime.now()
        return member

    member = MeetingRoomMember(
        room_id=room.id,
        user_id=user.id,
        member_role=member_role,
        display_name=get_room_member_display_name(user),
        status='arrived',
        joined_at=datetime.now(),
        last_seen_at=datetime.now()
    )
    db.session.add(member)
    return member


def generate_participants_for_room_gap(scenario, gap_count):
    if gap_count <= 0:
        return []

    existing_count = VirtualParticipant.query.filter_by(scenario_id=scenario.id).count()
    role_mapping = {
        "product_team": "企业产品团队成员",
        "sales_conversation": "销售相关角色",
        "focus_group": "目标用户",
        "user_interview": "深度访谈用户",
        "brainstorming": "创意团队成员"
    }
    role_type = role_mapping.get(scenario.occasion_type, "目标用户")
    prompt, system_prompt = build_generated_participant_prompt(
        scenario,
        role_type,
        gap_count,
        "请补足会议人数缺口，保持角色多样性，并贴合当前议题。",
        [],
        {}
    )

    result = call_deepseek(prompt, system_prompt)
    try:
        personas = json.loads(result)
        if not isinstance(personas, list):
            personas = [personas]
    except Exception:
        personas = build_fallback_personas([], gap_count)

    merged_personas = merge_generated_personas(personas, [], gap_count)
    created_participants = []
    for index, persona in enumerate(merged_personas, start=1):
        participant = VirtualParticipant(
            scenario_id=scenario.id,
            persona_name=persona.get('persona_name') or f"AI参与者{existing_count + index}",
            persona_tags=json_dumps(persona.get('persona_tags'), []),
            personality=persona.get('personality') or '理性务实',
            background=persona.get('background') or '基于当前会议主题补位生成',
            usage_goal=persona.get('usage_goal') or '',
            budget_sensitivity=persona.get('budget_sensitivity') or '',
            brand_preference=persona.get('brand_preference') or '',
            risk_aversion=persona.get('risk_aversion') or '',
            decision_style=persona.get('decision_style') or '',
            deal_breakers=json_dumps(persona.get('deal_breakers'), []),
            stance_summary=persona.get('stance_summary') or '',
            stance_state=build_participant_stance_state(persona.get('stance_summary') or ''),
            is_custom=False,
            is_ai_generated=True,
            custom_params=json_dumps(persona.get('custom_params'), {})
        )
        db.session.add(participant)
        created_participants.append(participant)

    db.session.flush()
    return [participant.to_dict() for participant in created_participants]


def get_meeting_room_members(room):
    return MeetingRoomMember.query.filter_by(room_id=room.id).order_by(MeetingRoomMember.joined_at.asc()).all()


def normalize_meeting_turn_order(room, members=None):
    members = members or get_meeting_room_members(room)
    valid_ids = [member.id for member in members]
    stored_order = safe_json_loads(room.turn_order, [])
    if not isinstance(stored_order, list):
        stored_order = []

    ordered_ids = []
    seen = set()
    for member_id in stored_order:
        try:
            normalized_id = int(member_id)
        except (TypeError, ValueError):
            continue
        if normalized_id in valid_ids and normalized_id not in seen:
            ordered_ids.append(normalized_id)
            seen.add(normalized_id)

    for member_id in valid_ids:
        if member_id not in seen:
            ordered_ids.append(member_id)
            seen.add(member_id)
    return ordered_ids


def get_room_managed_participants(scenario_id):
    return (
        VirtualParticipant.query
        .filter_by(scenario_id=scenario_id, room_managed=True)
        .order_by(VirtualParticipant.seat_order.asc(), VirtualParticipant.id.asc())
        .all()
    )


def build_meeting_member_participant(member, scenario, seat_order):
    tags = ['会议成员']
    if (member.member_role or '') == 'owner':
        tags.append('房主')

    participant = VirtualParticipant(
        scenario_id=scenario.id,
        persona_name=member.display_name or f"成员 {member.user_id}",
        persona_tags=json_dumps(tags, []),
        personality='表达直接、愿意围绕主题给出真实判断',
        background='真实参会成员，将基于自己的立场手动发表观点',
        usage_goal='围绕会议主题给出自己的真实看法',
        budget_sensitivity='中等',
        brand_preference='更关注实际价值和落地效果',
        risk_aversion='中等',
        decision_style='先表达核心判断，再根据讨论调整观点',
        deal_breakers=json_dumps(["主题不清晰", "无法落地执行"], []),
        stance_summary=f"将围绕“{scenario.decision_problem or scenario.product_name}”给出真实看法",
        stance_state=build_participant_stance_state(''),
        is_custom=True,
        is_ai_generated=False,
        custom_params=json_dumps({'source': 'meeting_member'}, {}),
        linked_user_id=member.user_id,
        meeting_member_id=member.id,
        speaker_origin='human',
        seat_order=seat_order,
        room_managed=True
    )
    normalize_participant_profile(participant, scenario)
    return participant


def materialize_room_ai_participants(scenario, generated_personas, starting_seat_order):
    participants = []
    ordered_ids = []
    for persona in generated_personas or []:
        persona_id = persona.get('id') if isinstance(persona, dict) else None
        participant = db.session.get(VirtualParticipant, persona_id) if persona_id else None
        if not participant:
            participant = VirtualParticipant(
                scenario_id=scenario.id,
                persona_name=(persona or {}).get('persona_name') or f"AI 参与者 {starting_seat_order + len(participants) + 1}",
                persona_tags=json_dumps((persona or {}).get('persona_tags'), []),
                personality=(persona or {}).get('personality') or '理性务实',
                background=(persona or {}).get('background') or '基于会议主题自动补位',
                usage_goal=(persona or {}).get('usage_goal') or '',
                budget_sensitivity=(persona or {}).get('budget_sensitivity') or '',
                brand_preference=(persona or {}).get('brand_preference') or '',
                risk_aversion=(persona or {}).get('risk_aversion') or '',
                decision_style=(persona or {}).get('decision_style') or '',
                deal_breakers=json_dumps((persona or {}).get('deal_breakers'), []),
                stance_summary=(persona or {}).get('stance_summary') or '',
                stance_state=build_participant_stance_state((persona or {}).get('stance_summary') or ''),
                is_custom=False,
                is_ai_generated=True,
                custom_params=json_dumps((persona or {}).get('custom_params'), {}),
            )
            db.session.add(participant)
            db.session.flush()

        participant.linked_user_id = None
        participant.meeting_member_id = None
        participant.speaker_origin = 'ai'
        participant.room_managed = True
        participant.seat_order = starting_seat_order + len(participants)
        normalize_participant_profile(participant, scenario)
        participants.append(participant)
        ordered_ids.append(participant.id)
    return participants


def serialize_discussion_record(record, participant_map):
    participant = participant_map.get(record.participant_id)
    participant_name = '系统'
    linked_user_id = None
    speaker_origin = 'system'
    seat_order = None
    if record.participant_id == 0:
        participant_name = 'AI主持人'
        speaker_origin = 'host'
    elif record.participant_id == -1:
        participant_name = '我'
        speaker_origin = 'human'
    elif participant:
        participant_name = participant.persona_name
        linked_user_id = participant.linked_user_id
        speaker_origin = participant.speaker_origin or ('human' if participant.linked_user_id else 'ai')
        seat_order = participant.seat_order

    return {
        "id": record.id,
        "participant_id": record.participant_id,
        "participant_name": participant_name,
        "linked_user_id": linked_user_id,
        "speaker_origin": speaker_origin,
        "seat_order": seat_order,
        "content": record.content,
        "timestamp": record.timestamp.isoformat() if record.timestamp else '',
        "is_host": record.is_host,
        "message_type": record.message_type
    }


def build_meeting_discussion_payload(room, scenario=None, current_user=None):
    scenario = scenario or db.session.get(ProductScenario, room.scenario_id)
    current_user = current_user or get_current_user_model()
    participants = get_room_managed_participants(room.scenario_id)
    participant_map = {participant.id: participant for participant in participants}
    records = ConversationRecord.query.filter_by(scenario_id=room.scenario_id).order_by(ConversationRecord.timestamp.asc(), ConversationRecord.id.asc()).all()
    messages = [serialize_discussion_record(record, participant_map) for record in records]

    active_participant = participant_map.get(room.active_speaker_participant_id)
    active_speaker = active_participant.to_dict() if active_participant else None
    current_user_id = current_user.id if current_user else None
    can_speak_now = bool(
        active_participant and
        (room.discussion_phase or 'waiting') == 'live' and
        (active_participant.speaker_origin or 'ai') == 'human' and
        active_participant.linked_user_id == current_user_id
    )

    return {
        'phase': room.discussion_phase or ('ended' if (room.status or '') == 'ended' else 'waiting'),
        'turn_order': [participant.id for participant in participants],
        'current_turn_index': room.current_turn_index or 0,
        'turn_number': room.turn_number or 1,
        'active_speaker_participant_id': room.active_speaker_participant_id,
        'ai_turn_pending': bool(room.ai_turn_pending),
        'active_speaker': active_speaker,
        'seats': [participant.to_dict() for participant in participants],
        'messages': messages,
        'current_user': {
            'user_id': current_user_id,
            'is_owner': bool(current_user_id and room.owner_user_id == current_user_id),
            'can_speak_now': can_speak_now,
            'active_participant_id': active_participant.id if active_participant else None
        }
    }


def serialize_meeting_room_bundle(room, scenario=None, include_messages=False):
    scenario = scenario or db.session.get(ProductScenario, room.scenario_id)
    members = get_meeting_room_members(room)
    payload = {
        'room': room.to_dict(),
        'scenario': scenario.to_dict() if scenario else None,
        'members': [member.to_dict() for member in members],
        'discussion': build_meeting_discussion_payload(room, scenario)
    }
    if include_messages:
        messages = MeetingRoomMessage.query.filter_by(room_id=room.id).order_by(MeetingRoomMessage.created_at.asc()).all()
        payload['messages'] = [message.to_dict() for message in messages]
    return payload


def initialize_room_discussion(room, scenario, generated_personas=None):
    members = get_meeting_room_members(room)
    ordered_member_ids = normalize_meeting_turn_order(room, members)
    member_map = {member.id: member for member in members}
    room.turn_order = json_dumps(ordered_member_ids, [])

    ConversationRecord.query.filter_by(scenario_id=scenario.id).delete(synchronize_session=False)
    VirtualParticipant.query.filter_by(scenario_id=scenario.id, room_managed=True).delete(synchronize_session=False)
    db.session.flush()

    created_participants = []
    for seat_order, member_id in enumerate(ordered_member_ids):
        member = member_map.get(member_id)
        if not member:
            continue
        participant = build_meeting_member_participant(member, scenario, seat_order)
        db.session.add(participant)
        created_participants.append(participant)

    db.session.flush()

    ai_participants = materialize_room_ai_participants(
        scenario,
        generated_personas,
        len(created_participants)
    )
    created_participants.extend(ai_participants)

    room.status = 'active'
    room.discussion_phase = 'live'
    room.current_turn_index = 0
    room.turn_number = 1
    room.active_speaker_participant_id = created_participants[0].id if created_participants else None
    room.ai_turn_pending = False
    return created_participants


def build_room_discussion_history_text(records, participant_map):
    history_lines = []
    for record in records[-8:]:
        serialized = serialize_discussion_record(record, participant_map)
        history_lines.append(f"{serialized['participant_name']}：{serialized['content']}")
    return "\n".join(history_lines)


def generate_meeting_room_ai_response(room, scenario, participant, participant_map):
    records = ConversationRecord.query.filter_by(scenario_id=scenario.id).order_by(ConversationRecord.timestamp.asc(), ConversationRecord.id.asc()).all()
    history_text = build_room_discussion_history_text(records, participant_map)
    topics = safe_json_loads(scenario.discussion_topics, [])
    seat_count = max(len(participant_map), 1)
    round_num = max((room.turn_number or 1) - 1, 0) // seat_count
    current_topic = topics[min(round_num, len(topics) - 1)] if topics else (scenario.decision_problem or scenario.product_name)

    evidence_items = ExternalEvidence.query.filter_by(scenario_id=scenario.id).order_by(ExternalEvidence.created_at.desc()).all()
    evidence_text = "\n".join([
        f"- {(item.source_title or item.source_label or item.evidence_type or '外部资料')}：{item.content}"
        for item in evidence_items[:5]
    ])

    prompt, system_prompt, _ = build_quality_participant_prompt(
        participant,
        scenario,
        current_topic,
        history_text,
        evidence_text,
        opening=(room.turn_number or 1) <= seat_count
    )
    response = call_deepseek(prompt, system_prompt)
    if is_ai_error_text(response):
        response = build_quality_fallback_response(
            participant,
            scenario,
            current_topic,
            opening=(room.turn_number or 1) <= seat_count
        )
    return response


_room_ai_turn_lock = threading.Lock()
_room_ai_turn_jobs = set()


def get_active_room_participant(room, participants=None):
    participants = participants or get_room_managed_participants(room.scenario_id)
    if not room.active_speaker_participant_id:
        return None
    return next((participant for participant in participants if participant.id == room.active_speaker_participant_id), None)


def activate_next_room_speaker(room, participants=None):
    participants = participants or get_room_managed_participants(room.scenario_id)
    if not participants:
        room.active_speaker_participant_id = None
        room.ai_turn_pending = False
        return None

    next_index = ((room.current_turn_index or 0) + 1) % len(participants)
    room.current_turn_index = next_index
    room.turn_number = max(room.turn_number or 1, 1) + 1
    next_participant = participants[next_index]
    room.active_speaker_participant_id = next_participant.id
    return next_participant


def process_current_room_ai_turn(room, scenario, participants=None):
    participants = participants or get_room_managed_participants(room.scenario_id)
    if not participants:
        room.active_speaker_participant_id = None
        room.ai_turn_pending = False
        return False

    active_participant = get_active_room_participant(room, participants)
    if not active_participant or (active_participant.speaker_origin or 'ai') != 'ai':
        room.ai_turn_pending = False
        return False

    participant_map = {participant.id: participant for participant in participants}
    ai_response = generate_meeting_room_ai_response(room, scenario, active_participant, participant_map)
    ai_record = ConversationRecord(
        scenario_id=scenario.id,
        participant_id=active_participant.id,
        content=ai_response,
        is_host=False,
        message_type='turn'
    )
    db.session.add(ai_record)
    update_participant_state_v2(
        active_participant,
        ai_response,
        max((room.turn_number or 1) - 1, 0) // max(len(participants), 1)
    )
    return True


def resolve_room_ai_turns_inline(room, scenario):
    participants = get_room_managed_participants(room.scenario_id)
    if not participants:
        room.active_speaker_participant_id = None
        room.ai_turn_pending = False
        return 'empty'

    safety_guard = 0
    while safety_guard < len(participants):
        active_participant = get_active_room_participant(room, participants)
        if not active_participant or (active_participant.speaker_origin or 'ai') != 'ai':
            room.ai_turn_pending = False
            return 'human'

        room.ai_turn_pending = True
        if not process_current_room_ai_turn(room, scenario, participants):
            room.ai_turn_pending = False
            return 'human'

        safety_guard += 1
        next_participant = activate_next_room_speaker(room, participants)
        if not next_participant:
            room.ai_turn_pending = False
            return 'empty'
        if (next_participant.speaker_origin or 'ai') != 'ai':
            room.ai_turn_pending = False
            return 'human'

    room.ai_turn_pending = False
    return 'human'


def continue_room_discussion_after_turn(room, scenario, participants=None, async_ai=False):
    if (room.discussion_phase or 'waiting') != 'live':
        room.ai_turn_pending = False
        return 'inactive'

    participants = participants or get_room_managed_participants(room.scenario_id)
    if not participants:
        room.active_speaker_participant_id = None
        room.ai_turn_pending = False
        return 'empty'

    next_participant = activate_next_room_speaker(room, participants)
    if not next_participant:
        room.ai_turn_pending = False
        return 'empty'

    if (next_participant.speaker_origin or 'ai') != 'ai':
        room.ai_turn_pending = False
        return 'human'

    room.ai_turn_pending = True
    if async_ai:
        return 'ai_pending'
    return resolve_room_ai_turns_inline(room, scenario)


def run_room_ai_turn_worker(room_id):
    try:
        safety_guard = 0
        more_ai_turns = False
        while safety_guard < 24:
            with app.app_context():
                room = db.session.get(MeetingRoom, room_id)
                if not room or (room.discussion_phase or 'waiting') != 'live' or (room.status or 'waiting') == 'ended':
                    db.session.remove()
                    break

                scenario = db.session.get(ProductScenario, room.scenario_id)
                if not scenario:
                    room.ai_turn_pending = False
                    db.session.commit()
                    db.session.remove()
                    break

                participants = get_room_managed_participants(room.scenario_id)
                active_participant = get_active_room_participant(room, participants)
                if not active_participant or (active_participant.speaker_origin or 'ai') != 'ai':
                    room.ai_turn_pending = False
                    room.updated_at = datetime.now()
                    db.session.commit()
                    db.session.remove()
                    break

                room.ai_turn_pending = True
                process_current_room_ai_turn(room, scenario, participants)
                next_participant = activate_next_room_speaker(room, participants)
                more_ai_turns = bool(next_participant and (next_participant.speaker_origin or 'ai') == 'ai')
                room.ai_turn_pending = more_ai_turns
                room.updated_at = datetime.now()
                db.session.commit()
                db.session.remove()
                safety_guard += 1

            if not more_ai_turns:
                break
    finally:
        with _room_ai_turn_lock:
            _room_ai_turn_jobs.discard(room_id)


def enqueue_room_ai_turn_worker(room_id):
    with _room_ai_turn_lock:
        if room_id in _room_ai_turn_jobs:
            return False
        _room_ai_turn_jobs.add(room_id)

    worker = threading.Thread(
        target=run_room_ai_turn_worker,
        args=(room_id,),
        daemon=True
    )
    worker.start()
    return True


def advance_room_discussion(room, scenario):
    if (room.discussion_phase or 'waiting') != 'live':
        return

    continue_room_discussion_after_turn(room, scenario, async_ai=False)


def ensure_current_human_turn(room, scenario, async_ai=False):
    if (room.discussion_phase or 'waiting') != 'live':
        room.ai_turn_pending = False
        return 'inactive'
    participants = get_room_managed_participants(room.scenario_id)
    if not participants:
        room.active_speaker_participant_id = None
        room.ai_turn_pending = False
        return 'empty'
    if not room.active_speaker_participant_id:
        room.current_turn_index = 0
        room.turn_number = max(room.turn_number or 1, 1)
        room.active_speaker_participant_id = participants[0].id
    active_participant = get_active_room_participant(room, participants)
    if active_participant and (active_participant.speaker_origin or 'ai') == 'ai':
        room.ai_turn_pending = True
        if async_ai:
            return 'ai_pending'
        return resolve_room_ai_turns_inline(room, scenario)
    room.ai_turn_pending = False
    return 'human'


def validate_verification_request(target, purpose):
    normalized_target = normalize_auth_target(target)
    normalized_purpose = (purpose or 'login').strip()

    if normalized_purpose not in ('login', 'register', 'reset_password', 'change_password'):
        return None, None, "不支持的验证码用途", 400
    if not normalized_target:
        return None, None, "接收地址不能为空", 400
    if not validate_email(normalized_target):
        return None, None, "请输入有效的邮箱地址", 400

    existing_user = User.query.filter_by(email=normalized_target).first()
    if normalized_purpose == 'register' and existing_user:
        return None, None, "邮箱已被注册", 400
    if normalized_purpose in ('reset_password', 'change_password') and not existing_user:
        return None, None, "该邮箱尚未注册", 404

    return normalized_target, normalized_purpose, None, None

# 登录速率限制存储（内存 dict，单进程有效）
_login_attempts = {}
_login_lock = threading.Lock()

def check_login_rate_limit(key: str) -> bool:
    """检查登录速率限制，返回 True 表示可以继续"""
    now = time.time()
    with _login_lock:
        attempts = _login_attempts.get(key, [])
        attempts = [t for t in attempts if now - t < LOGIN_RATE_WINDOW]
        if len(attempts) >= LOGIN_RATE_LIMIT:
            _login_attempts[key] = attempts
            return False
        attempts.append(now)
        _login_attempts[key] = attempts
    return True

# ======================== 验证码工具函数 ========================

_verification_lock = threading.Lock()

def _gen_numeric_code(length=6):
    """生成指定长度的数字验证码"""
    return ''.join(secrets.choice(string.digits) for _ in range(length))

def _get_client_ip():
    """获取客户端 IP"""
    return (request.headers.get('X-Forwarded-For') or request.remote_addr or '').split(',')[0].strip()

def _parse_target_to_channel(target):
    """根据输入判断是邮箱还是手机号"""
    if validate_email(target):
        return 'email'
    digits = re.sub(r'\D', '', target)
    if len(digits) == 11 and digits.startswith('1'):
        return 'sms'
    return None

def _pick_channel(target):
    """根据配置和目标类型选择实际发送渠道"""
    auto = VERIFICATION_CHANNEL.lower()
    detected = _parse_target_to_channel(target)
    if auto == 'test':
        return 'test'
    if detected == 'sms':
        return 'sms'
    if detected == 'email' and auto in ('email', 'auto') and EMAIL_SMTP_HOST and EMAIL_SMTP_USER and EMAIL_SMTP_PASSWORD:
        return 'email'
    return detected

def _rate_limit_check(target, ip, purpose, cooldown):
    """检查频率限制"""
    threshold = datetime.now() - timedelta(seconds=cooldown)
    latest = VerificationCode.query.filter(
        VerificationCode.purpose == purpose,
        (VerificationCode.target == target) | (VerificationCode.ip == ip)
    ).order_by(VerificationCode.created_at.desc()).first()
    return latest is None or latest.created_at <= threshold

def _send_email(to_addr, code, purpose):
    """发送邮件验证码"""
    subject_map = {
        "login": "【焦点小组系统】登录验证码",
        "register": "【焦点小组系统】注册验证码",
        "reset_password": "【焦点小组系统】重置密码验证码",
        "change_password": "【焦点小组系统】修改密码验证码"
    }
    hint_map = {
        "login": "您正在进行登录操作",
        "register": "您正在进行账号注册",
        "reset_password": "您正在进行密码重置",
        "change_password": "您正在进行密码修改"
    }
    subject = subject_map.get(purpose, "【焦点小组系统】验证码")
    hint = hint_map.get(purpose, "您正在进行身份验证")

    html_body = f"""
    <div style="font-family: -apple-system, Arial, sans-serif; max-width: 560px; margin: 0 auto;">
        <div style="padding: 24px; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: #fff; border-radius: 8px 8px 0 0;">
            <h2 style="margin:0;">{EMAIL_FROM_NAME}</h2>
            <p style="margin:8px 0 0; opacity:.9;">{hint}</p>
        </div>
        <div style="padding: 32px; background: #fff; border: 1px solid #e5e7eb; border-top:none; border-radius: 0 0 8px 8px;">
            <p style="margin:0 0 16px; color:#374151;">您好，您的验证码为：</p>
            <div style="font-size:36px; font-weight:700; letter-spacing:12px; color:#667eea; text-align:center; padding:20px 0; background:#f9fafb; border-radius:8px; margin-bottom:16px;">
                {code}
            </div>
            <p style="margin:0; color:#6b7280; font-size:14px;">该验证码 {VERIFICATION_TTL//60} 分钟内有效，请勿转发给他人。</p>
            <p style="margin:8px 0 0; color:#6b7280; font-size:14px;">如非本人操作，请忽略此邮件。</p>
        </div>
        <p style="text-align:center; color:#9ca3af; font-size:12px; margin-top:16px;">此邮件由系统自动发送，请勿回复</p>
    </div>
    """
    try:
        msg = MIMEText(html_body, 'html', 'utf-8')
        msg['From'] = formataddr((str(Header(EMAIL_FROM_NAME, 'utf-8')), EMAIL_SMTP_USER))
        msg['To'] = to_addr
        msg['Subject'] = Header(subject, 'utf-8')

        if EMAIL_SMTP_USE_SSL:
            server = smtplib.SMTP_SSL(EMAIL_SMTP_HOST, int(EMAIL_SMTP_PORT), timeout=15)
        else:
            server = smtplib.SMTP(EMAIL_SMTP_HOST, int(EMAIL_SMTP_PORT), timeout=15)
            server.starttls()
        server.login(EMAIL_SMTP_USER, EMAIL_SMTP_PASSWORD)
        server.sendmail(EMAIL_SMTP_USER, [to_addr], msg.as_string())
        server.quit()
        return True, None
    except Exception as e:
        return False, str(e)

def _send_sms(phone, code, purpose):
    """发送短信验证码（预留接口）"""
    if not (SMS_ACCESS_KEY_ID and SMS_ACCESS_KEY_SECRET and SMS_SIGN_NAME and SMS_TEMPLATE_CODE):
        return False, "短信服务未配置"
    return False, "短信服务已配置但未接入实现"


def _smtp_config_ready():
    return bool(EMAIL_SMTP_HOST and EMAIL_SMTP_USER and EMAIL_SMTP_PASSWORD)

def send_verification_code(target, purpose='login'):
    """发送验证码统一入口"""
    if not target or not target.strip():
        return False, "接收地址不能为空", None, None

    target = normalize_auth_target(target)
    detected = _parse_target_to_channel(target)
    if detected == 'sms':
        return False, "当前版本暂不支持短信验证码，请使用邮箱", 'sms', None
    if detected != 'email':
        return False, "请输入有效的邮箱地址", None, None

    channel = _pick_channel(target)
    code = _gen_numeric_code(6)
    ip = _get_client_ip()

    with _verification_lock:
        if not _rate_limit_check(target, ip, purpose, VERIFICATION_COOLDOWN):
            return False, f"请求过于频繁，请 {VERIFICATION_COOLDOWN} 秒后再试", None, None

    sent_ok, err = False, None
    if channel == 'email':
        if not _smtp_config_ready():
            return False, "邮箱验证码服务未配置", channel, None
        sent_ok, err = _send_email(target, code, purpose)
    elif channel == 'test':
        sent_ok, err = True, None
    else:
        sent_ok, err = False, "邮箱验证码服务未配置"

    if not sent_ok:
        print(f"[VERIFICATION] 发送失败（target={target}, channel={channel}）：{err}")
        return False, err or "验证码发送失败", channel, None

    with _verification_lock:
        record = VerificationCode(
            channel=channel,
            target=target,
            code=code,
            purpose=purpose,
            expires_at=datetime.now() + timedelta(seconds=VERIFICATION_TTL),
            ip=ip
        )
        db.session.add(record)
        db.session.commit()

    if channel == 'test':
        print(f"[VERIFICATION-TEST] target={target}, purpose={purpose}, code={code}, expires={VERIFICATION_TTL}s")

    return True, "验证码已发送", channel, (code if channel == 'test' else None)

def verify_verification_code(target, code, purpose='login'):
    """验证验证码"""
    target = (target or '').strip().lower()
    code = (code or '').strip()
    if not target or not code:
        return False, "目标和验证码不能为空"

    now = datetime.now()
    records = VerificationCode.query.filter_by(
        target=target,
        purpose=purpose,
        used=False
    ).order_by(VerificationCode.created_at.desc()).all()

    if not records:
        return False, "未找到验证码，请先发送"

    matched = None
    for rec in records:
        if rec.expires_at < now:
            continue
        if rec.attempts >= VERIFICATION_MAX_ATTEMPTS:
            continue
        if rec.code == code:
            matched = rec
            break
        else:
            rec.attempts += 1
            db.session.commit()
            return False, f"验证码错误（还剩 {VERIFICATION_MAX_ATTEMPTS - rec.attempts} 次机会）"

    if not matched:
        return False, "验证码已过期或无效，请重新发送"

    matched.used = True
    db.session.commit()
    return True, "验证成功"

def build_ai_evidence_prompt(scenario):
    selling_points = safe_json_loads(scenario.core_selling_points, [])
    topics = safe_json_loads(scenario.discussion_topics, [])
    assumptions = normalize_text_list(safe_json_loads(scenario.validation_assumptions, scenario.validation_assumptions or []))
    return f"""请你作为企业研究助理，围绕下面这个产品场景，生成 3 条“可供内部决策讨论参考”的外部证据草案。

产品名称：{scenario.product_name}
产品概念：{scenario.product_concept}
核心卖点：{'、'.join(selling_points) if selling_points else '暂无'}
研究目标：{scenario.research_goal or '判断该概念是否值得继续推进'}
核心决策问题：{scenario.decision_problem or '当前阶段是否值得继续验证'}
目标用户：{scenario.target_user_profile or '暂无'}
竞品/替代方案：{scenario.competitor_context or '暂无'}
讨论主题：{'、'.join(topics) if topics else '暂无'}
待验证假设：{'、'.join(assumptions) if assumptions else '暂无'}

要求：
1. 输出 JSON 数组，且只能输出 JSON，不要输出解释文字。
2. 数组长度固定为 3。
3. 每个对象必须包含：title、evidence_type、source_label、strength_level、content。
4. evidence_type 只能是：market_data、user_signal、competitor_signal。
5. strength_level 只能是：低、中。
6. content 需要同时说明：这条线索说明什么、为什么会影响企业决策、为什么当前仍只是参考。
7. 不要伪造真实网址、论文编号、权威机构统计值；来源表述请保持为行业公开讨论、同类产品常见反馈、企业内部预判等参考级描述。
8. 语气专业，服务于企业判断，不要写成营销文案。"""

def ensure_extended_schema():
    tables = {
        'product_scenario': {
            'research_goal': "TEXT",
            'decision_problem': "TEXT",
            'target_user_profile': "TEXT",
            'competitor_context': "TEXT",
            'validation_assumptions': "TEXT",
            'research_plan': "TEXT",
            'research_plan_status': "VARCHAR(20) DEFAULT 'pending'",
            'meeting_status': "VARCHAR(20) DEFAULT 'standalone'"
        },
        'virtual_participant': {
            'usage_goal': "TEXT",
            'budget_sensitivity': "TEXT",
            'brand_preference': "TEXT",
            'risk_aversion': "TEXT",
            'decision_style': "TEXT",
            'deal_breakers': "TEXT",
            'stance_summary': "TEXT",
            'stance_state': "TEXT",
            'custom_params': "TEXT",
            'linked_user_id': "INTEGER",
            'meeting_member_id': "INTEGER",
            'speaker_origin': "VARCHAR(20) DEFAULT 'ai'",
            'seat_order': "INTEGER",
            'room_managed': "BOOLEAN DEFAULT 0"
        },
        'conversation_record': {
            'interaction_intent': "VARCHAR(50)"
        },
        'analysis_report': {
            'executive_summary': "TEXT",
            'key_assumptions': "TEXT",
            'evidence_items': "TEXT",
            'decision_risks': "TEXT",
            'recommended_actions': "TEXT",
            'confidence_level': "VARCHAR(20)",
            'source_breakdown': "TEXT",
            'discussion_summary': "TEXT",
            'presentation_payload': "TEXT",
            'presentation_generated_at': "DATETIME"
        },
        'external_evidence': {
            'source_url': "TEXT",
            'source_title': "VARCHAR(300)",
            'source_domain': "VARCHAR(200)",
            'generated_by': "VARCHAR(30) DEFAULT 'user_manual'"
        },
        'user': {
            'nickname': "VARCHAR(80)",
            'avatar': "VARCHAR(200)",
            'company': "VARCHAR(120)",
            'role': "VARCHAR(20) DEFAULT 'user'",
            'bio': "TEXT",
            'locale': "VARCHAR(10) DEFAULT 'zh-CN'",
            'last_login': "DATETIME",
            'login_count': "INTEGER DEFAULT 0",
            'ai_endpoint_url': "TEXT",
            'ai_model_name': "VARCHAR(120)",
            'ai_api_keys_encrypted': "TEXT",
            'ai_config_enabled': "BOOLEAN DEFAULT 0",
            'ai_last_test_status': "VARCHAR(20)",
            'ai_last_test_message': "TEXT",
            'ai_last_tested_at': "DATETIME"
        },
        'meeting_room': {
            'scenario_id': "INTEGER",
            'owner_user_id': "INTEGER NOT NULL",
            'room_name': "VARCHAR(160) NOT NULL",
            'room_code': "VARCHAR(20) NOT NULL",
            'invite_token': "VARCHAR(64) NOT NULL",
            'status': "VARCHAR(20) DEFAULT 'waiting'",
            'topic_title': "TEXT NOT NULL",
            'topic_notes': "TEXT",
            'target_count': "INTEGER DEFAULT 4",
            'created_at': "DATETIME",
            'started_at': "DATETIME",
            'updated_at': "DATETIME",
            'discussion_phase': "VARCHAR(20) DEFAULT 'waiting'",
            'turn_order': "TEXT",
            'current_turn_index': "INTEGER DEFAULT 0",
            'turn_number': "INTEGER DEFAULT 1",
            'active_speaker_participant_id': "INTEGER",
            'ai_turn_pending': "BOOLEAN DEFAULT 0"
        },
        'meeting_room_member': {
            'room_id': "INTEGER NOT NULL",
            'user_id': "INTEGER NOT NULL",
            'member_role': "VARCHAR(20) DEFAULT 'member'",
            'display_name': "VARCHAR(120)",
            'status': "VARCHAR(20) DEFAULT 'arrived'",
            'joined_at': "DATETIME",
            'last_seen_at': "DATETIME"
        },
        'meeting_room_message': {
            'room_id': "INTEGER NOT NULL",
            'user_id': "INTEGER",
            'sender_name': "VARCHAR(120) NOT NULL",
            'content': "TEXT NOT NULL",
            'message_type': "VARCHAR(20) DEFAULT 'member'",
            'created_at': "DATETIME"
        }
    }

    with db.engine.begin() as connection:
        for table_name, columns in tables.items():
            try:
                existing_columns = connection.exec_driver_sql(f"PRAGMA table_info({table_name})").fetchall()
                existing_names = {column[1] for column in existing_columns}
                for column_name, ddl in columns.items():
                    if column_name not in existing_names:
                        connection.exec_driver_sql(f"ALTER TABLE {table_name} ADD COLUMN {column_name} {ddl}")
            except Exception as e:
                print(f"[Schema] Failed to alter {table_name}: {e}")

    ExternalEvidence.__table__.create(bind=db.engine, checkfirst=True)
    ResearchRound.__table__.create(bind=db.engine, checkfirst=True)
    VerificationCode.__table__.create(bind=db.engine, checkfirst=True)
    MeetingRoom.__table__.create(bind=db.engine, checkfirst=True)
    MeetingRoomMember.__table__.create(bind=db.engine, checkfirst=True)
    MeetingRoomMessage.__table__.create(bind=db.engine, checkfirst=True)

    scenarios = ProductScenario.query.all()
    for scenario in scenarios:
        room = MeetingRoom.query.filter_by(scenario_id=scenario.id).first()
        sync_scenario_meeting_status(scenario, room)
    db.session.commit()

def extract_json_block(text):
    cleaned = (text or '').strip()
    cleaned = re.sub(r'^```json\s*', '', cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r'^```\s*', '', cleaned)
    cleaned = re.sub(r'\s*```$', '', cleaned)
    match = re.search(r'(\[[\s\S]*\])', cleaned)
    return match.group(1) if match else cleaned


def extract_json_object_block(text):
    cleaned = (text or '').strip()
    cleaned = re.sub(r'^```json\s*', '', cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r'^```\s*', '', cleaned)
    cleaned = re.sub(r'\s*```$', '', cleaned)
    match = re.search(r'(\{[\s\S]*\})', cleaned)
    return match.group(1) if match else cleaned

def build_search_queries(scenario):
    queries = []
    base_terms = [scenario.product_name, scenario.product_concept, scenario.target_user_profile or '', scenario.competitor_context or '']
    joined = ' '.join(term for term in base_terms if term).strip()
    if joined:
        queries.append(f"{joined} 市场分析")
        queries.append(f"{joined} 用户需求")
    selling_points = safe_json_loads(scenario.core_selling_points, [])
    if selling_points:
        queries.append(f"{scenario.product_name} {' '.join(selling_points[:2])} 竞品")
    topics = safe_json_loads(scenario.discussion_topics, [])
    for topic in topics[:2]:
        queries.append(f"{scenario.product_name} {topic}")
    seen = set()
    result = []
    for query in queries:
        normalized = query.strip()
        if normalized and normalized not in seen:
            seen.add(normalized)
            result.append(normalized)
    return result[:5]

def search_duckduckgo(query, limit=5):
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8"
    }
    response = requests.get(
        "https://html.duckduckgo.com/html/",
        params={"q": query},
        headers=headers,
        timeout=20
    )
    response.raise_for_status()
    html = response.text
    pattern = re.compile(r'<a[^>]*class="result__a"[^>]*href="(?P<href>[^"]+)"[^>]*>(?P<title>.*?)</a>', re.I)
    results = []
    for match in pattern.finditer(html):
        href = re.sub(r'^//', 'https://', match.group('href'))
        href = re.sub(r'&amp;', '&', href)
        title = re.sub(r'<.*?>', '', match.group('title'))
        domain_match = re.search(r'https?://([^/]+)', href)
        results.append({
            "title": title.strip(),
            "url": href.strip(),
            "domain": domain_match.group(1).strip() if domain_match else ''
        })
        if len(results) >= limit:
            break
    return results

def fetch_page_excerpt(url):
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8"
    }
    response = requests.get(url, headers=headers, timeout=20)
    response.raise_for_status()
    text = re.sub(r'<script[\s\S]*?</script>', ' ', response.text, flags=re.I)
    text = re.sub(r'<style[\s\S]*?</style>', ' ', text, flags=re.I)
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text[:1200]

def build_evidence_synthesis_prompt(scenario, sources):
    source_lines = []
    for idx, source in enumerate(sources, start=1):
        source_lines.append(
            f"[{idx}] 标题：{source['title']}\n来源：{source['domain']}\n链接：{source['url']}\n摘要：{source['excerpt']}"
        )
    return f"""你是企业研究助理。请根据下面的真实检索结果，整理 3 条适合企业内部决策讨论的外部证据。

产品名称：{scenario.product_name}
产品概念：{scenario.product_concept}
研究目标：{scenario.research_goal or '判断是否值得继续推进'}
核心决策问题：{scenario.decision_problem or '当前阶段是否值得继续验证'}

检索结果：
{chr(10).join(source_lines)}

要求：
1. 只输出 JSON 数组，不要解释文字。
2. 数组长度固定为 3。
3. 每项必须包含：title、evidence_type、source_label、source_url、source_title、source_domain、strength_level、content。
4. strength_level 只能是“低”或“中”。
5. content 必须明确写出：观察到的外部信号、它为什么影响企业决策、当前仍需验证的地方。
6. 只能引用上面的检索结果，不要杜撰新来源。"""

def parse_report_sections(content):
    sections = []
    current_title = None
    current_lines = []
    for raw_line in (content or '').splitlines():
        line = raw_line.strip()
        heading = re.match(r'^##\s+(.+?)\s*$', line)
        if heading:
            if current_title:
                sections.append((current_title, '\n'.join(current_lines).strip()))
            current_title = heading.group(1).strip()
            current_lines = []
            continue
        current_lines.append(raw_line.rstrip())
    if current_title:
        sections.append((current_title, '\n'.join(current_lines).strip()))
    return sections

def first_matching_lines(body, pattern, limit=3):
    items = []
    for line in (body or '').splitlines():
        text = line.strip()
        if not text:
            continue
        if re.search(pattern, text):
            items.append(re.sub(r'^[-*\d\.\s、]+', '', text).strip())
        if len(items) >= limit:
            break
    return items

def extract_report_metadata(report_content, external_evidence=None):
    sections = parse_report_sections(report_content)
    evidence_items = []
    recommended_actions = []
    decision_risks = []
    key_assumptions = []
    source_breakdown = []

    for title, body in sections:
        short_lines = [line.strip() for line in body.splitlines() if line.strip()]
        source_breakdown.append({
            'section': title,
            'summary': short_lines[0][:180] if short_lines else '',
            'source_type': 'discussion'
        })

        if '总结' in title or '观点' in title:
            key_assumptions.extend(first_matching_lines(body, r'假设|前提|如果|需要验证', 4))
        if '分析' in title:
            evidence_items.extend(first_matching_lines(body, r'讨论证据|参与者|王磊|李思琪|张晨|陈晓峰|数据|案例', 6))
        if '分歧' in title:
            decision_risks.extend(first_matching_lines(body, r'风险|顾虑|挑战|停止条件', 5))
        if '建议' in title or '结论' in title:
            recommended_actions.extend(first_matching_lines(body, r'负责人|建议|动作|优先级|下一步', 8))

    external_evidence = external_evidence or []
    for item in external_evidence:
        evidence_items.append(f"{item.title}｜{item.source_label or '外部来源'}｜{item.content[:120]}")
        source_breakdown.append({
            'section': item.title,
            'summary': item.content[:180],
            'source_type': 'external',
            'strength_level': item.strength_level
        })

    evidence_strength = '低'
    if external_evidence:
        evidence_strength = '高'
    elif len(evidence_items) >= 3:
        evidence_strength = '中'

    return {
        'key_assumptions': list(dict.fromkeys(key_assumptions))[:6],
        'evidence_items': list(dict.fromkeys(evidence_items))[:8],
        'decision_risks': list(dict.fromkeys(decision_risks))[:6],
        'recommended_actions': list(dict.fromkeys(recommended_actions))[:8],
        'confidence_level': evidence_strength,
        'source_breakdown': source_breakdown[:12]
    }

def build_research_plan(scenario):
    assumptions = normalize_text_list(safe_json_loads(scenario.validation_assumptions, scenario.validation_assumptions or []))
    topics = safe_json_loads(scenario.discussion_topics, [])
    lines = [
        "## 一、研究目标",
        scenario.research_goal or "明确该产品概念是否值得继续推进，并识别最需要验证的风险与机会。",
        "",
        "## 二、核心决策问题",
        scenario.decision_problem or "当前阶段是否应推进下一轮验证或试点。",
        "",
        "## 三、优先讨论主题"
    ]
    lines.extend([f"- {topic}" for topic in topics[:5]] or ["- 当前未配置讨论主题"])
    lines.extend(["", "## 四、需要重点验证的假设"])
    lines.extend([f"- {item}" for item in assumptions[:6]] or ["- 用户是否真正认可核心卖点", "- 是否存在可接受的价格区间"])
    lines.extend(["", "## 五、建议研究动作", "- 先生成稳定画像的虚拟参与者", "- 完成一轮讨论后输出管理层报告", "- 必要时补充外部证据以提升结论可信度"])
    return '\n'.join(lines)

def participant_state_prompt(participant):
    return f"""姓名：{participant.persona_name}
背景：{participant.background}
性格：{participant.personality}
标签：{'/'.join(safe_json_loads(participant.persona_tags, []))}
使用目标：{participant.usage_goal or '待观察'}
预算敏感度：{participant.budget_sensitivity or '中'}
品牌偏好：{participant.brand_preference or '暂无明显偏好'}
风险厌恶：{participant.risk_aversion or '中'}
决策风格：{participant.decision_style or '谨慎讨论后决定'}
反感点：{'、'.join(safe_json_loads(participant.deal_breakers, [])) or '暂无'}
稳定立场：{participant.stance_summary or '尚未形成稳定立场'}"""


AI_ERROR_PREFIXES = ("API Error", "Request Error", "请求失败")


def is_ai_error_text(text):
    return isinstance(text, str) and text.startswith(AI_ERROR_PREFIXES)


def normalize_participant_profile(participant, scenario):
    tags = safe_json_loads(participant.persona_tags, [])
    tags_text = "、".join(tags) if tags else "普通用户"
    product_name = scenario.product_name or "该产品"

    if not participant.usage_goal:
        participant.usage_goal = f"判断 {product_name} 是否真的能解决自己在 {tags_text} 场景下的实际问题"
    if not participant.budget_sensitivity:
        participant.budget_sensitivity = "中等，只有在价值足够明确时才愿意付费"
    if not participant.brand_preference:
        participant.brand_preference = "更看重实际体验和可信度，不会只因为新奇概念就买单"
    if not participant.risk_aversion:
        participant.risk_aversion = "中等偏高，担心功能不成熟或体验不稳定"
    if not participant.decision_style:
        participant.decision_style = "会先比较收益和风险，再决定是否尝试"

    deal_breakers = safe_json_loads(participant.deal_breakers, [])
    if not deal_breakers:
        deal_breakers = [
            "价格明显高于替代方案",
            "使用场景不清晰",
            "宣传概念大于真实价值"
        ]
        participant.deal_breakers = json_dumps(deal_breakers, [])

    if not participant.stance_summary:
        participant.stance_summary = (
            f"对 {product_name} 保持谨慎兴趣，会先关注真实使用场景、支付意愿和可替代方案，再决定是否支持推进。"
        )

    stance_state = safe_json_loads(participant.stance_state, {})
    if not isinstance(stance_state, dict):
        stance_state = {}
    if not stance_state.get("current_position"):
        stance_state["current_position"] = participant.stance_summary
    if not stance_state.get("confidence"):
        stance_state["confidence"] = "中"
    if "last_updated_round" not in stance_state:
        stance_state["last_updated_round"] = 0
    participant.stance_state = json_dumps(stance_state, {})


def build_quality_host_intro(scenario, selling_points, first_topic):
    selling_text = "、".join(selling_points[:4]) if selling_points else "核心卖点待进一步明确"
    return (
        f"各位好，今天我们围绕“{scenario.product_name}”做一场面向企业决策的虚拟焦点访谈。"
        f"这个概念的核心是：{scenario.product_concept}。"
        f"目前我们重点关注的卖点有：{selling_text}。"
        f"这轮讨论不是为了凑热闹，而是为了判断这个概念值不值得继续推进，先从“{first_topic}”开始。"
    )


def build_quality_participant_prompt(participant, scenario, topic, history_text, evidence_text, opening=False):
    interaction_types = ["回应", "补充", "追问", "质疑", "举例", "保留意见", "要求举证", "改变观点"]
    interaction_type = random.choice(interaction_types)
    research_goal = scenario.research_goal or "判断该产品概念是否值得继续推进"
    decision_problem = scenario.decision_problem or "当前阶段是否应该继续投入资源验证"

    base_context = f"""你正在参加一场企业研究型焦点访谈，请稳定扮演这个参与者，不要脱离人设。

研究目标：{research_goal}
核心决策问题：{decision_problem}
当前话题：{topic}
参与者画像：
{participant_state_prompt(participant)}

外部证据：
{evidence_text or '暂无外部证据，仅可基于现有讨论谨慎判断。'}"""

    if opening:
        prompt = f"""{base_context}

请你先做第一轮发言，像真实参会者一样表达，不要写成报告。
要求：
1. 明确说出你最在意的一个价值点或一个担忧点。
2. 最好结合自己的使用场景、预算、风险顾虑或过往经验。
3. 不要泛泛而谈，不要只说“有潜力”“挺不错”。
4. 控制在 70-140 字。
只输出发言内容。"""
    else:
        prompt = f"""{base_context}

最近讨论记录：
{history_text}

这次你的互动意图是：{interaction_type}
请继续参与讨论，并满足下面要求：
1. 尽量回应某个具体观点，而不是重新说一遍自己的总看法。
2. 可以质疑、补充例子、要求举证，或者说明自己为什么被说服/没有被说服。
3. 如果提到支持推进或反对推进，必须带上原因。
4. 控制在 80-170 字。
只输出发言内容。"""

    system_prompt = (
        "你不是客服，也不是泛泛而谈的助手。你是企业研究讨论里的真实参与者，"
        "要有稳定立场、真实顾虑和具体场景，宁可表达有限但明确的判断，也不要说空话套话。"
    )
    return prompt, system_prompt, interaction_type


def build_quality_fallback_response(participant, scenario, topic, opening=False):
    stance = participant.stance_summary or f"对 {scenario.product_name} 保持谨慎态度"
    usage_goal = participant.usage_goal or f"确认 {scenario.product_name} 是否值得尝试"
    risk = participant.risk_aversion or "中等"
    budget = participant.budget_sensitivity or "中等"

    if opening:
        return (
            f"如果只看 {topic}，我现在不会直接下结论。对我来说，先要看 {scenario.product_name} 能不能真正支撑“{usage_goal}”，"
            f"其次才是新鲜感。我的风险顾虑是{risk}，预算敏感度也是{budget}，所以如果只是概念好听、落地场景却不清楚，我不会支持马上推进。"
        )

    return (
        f"我想把话题拉回到 {topic}。结合刚才大家的说法，我现在最缺的还是一个能支持决策的依据："
        f"{stance}。如果下一步不能证明它在真实场景里比现有替代方案更有效，我倾向于先小范围验证，而不是直接扩大投入。"
    )


def build_quality_host_prompt(scenario, topic, conversation_history):
    recent_lines = "\n".join(
        [f"- {item['participant']}：{item['content']}" for item in conversation_history[-5:]]
    )
    research_goal = scenario.research_goal or "判断该产品概念是否值得继续推进"
    decision_problem = scenario.decision_problem or "当前阶段是否继续投入验证"
    prompt = f"""你是企业研究型焦点访谈的主持人，请根据最近一轮讨论继续推进。

研究目标：{research_goal}
核心决策问题：{decision_problem}
当前话题：{topic}
最近发言：
{recent_lines}

请输出 1 段 50-90 字的主持人引导语，满足以下要求：
1. 不要空泛总结，要推动讨论更具体。
2. 优先追问模糊观点、要求举例，或者请尚未充分表达的人说明分歧。
3. 如果大家过早达成一致，要提醒他们补充反例和风险。
只输出主持人发言内容。"""
    system_prompt = "你是一位研究型主持人，目标是帮团队得到更能支撑决策的讨论证据，而不是制造表面热闹。"
    return prompt, system_prompt


def build_quality_host_fallback(scenario, topic):
    decision_problem = scenario.decision_problem or "当前阶段是否继续推进"
    return (
        f"我先追问一句，大家刚才的观点里哪些是个人感觉，哪些是能支持“{decision_problem}”的依据？"
        f"围绕“{topic}”再各补一个具体场景或风险点。"
    )


def update_participant_state_v2(participant, response_text, round_num):
    stance_state = safe_json_loads(participant.stance_state, {})
    if not isinstance(stance_state, dict):
        stance_state = {}
    stance_state["current_position"] = (response_text or "")[:200]
    stance_state["last_updated_round"] = round_num + 1
    if re.search(r"改观|被说服|现在更倾向|重新考虑|我收回", response_text or ""):
        stance_state["confidence"] = "变化中"
    elif not stance_state.get("confidence"):
        stance_state["confidence"] = "中"
    participant.stance_state = json_dumps(stance_state, {})
    if response_text and len(response_text.strip()) > 12:
        participant.stance_summary = response_text[:160]
    db.session.add(participant)


def normalize_report_confidence_level(value, has_external_evidence=False, evidence_count=0):
    text = (value or '').strip()
    if text in {'高', '中', '低'}:
        return text
    if has_external_evidence:
        return '高'
    if evidence_count >= 3:
        return '中'
    return '低'


def sanitize_report_text(text):
    cleaned = (text or '').strip()
    if not cleaned:
        return ''
    start_marker = "=====REPORT_START====="
    end_marker = "=====REPORT_END====="
    start_idx = cleaned.find(start_marker)
    end_idx = cleaned.find(end_marker)
    if start_idx == -1 or end_idx == -1 or end_idx <= start_idx:
        return ''
    cleaned = cleaned[start_idx + len(start_marker):end_idx].strip()
    return cleaned.strip()


REPORT_SECTION_SPECS = [
    {
        'code': 'section_summary',
        'title': '一、核心观点总结',
        'heading': '## 一、核心观点总结',
        'phase': '生成核心观点总结',
        'pct': 24,
    },
    {
        'code': 'section_analysis',
        'title': '二、关键讨论点分析',
        'heading': '## 二、关键讨论点分析',
        'phase': '生成关键讨论点分析',
        'pct': 42,
    },
    {
        'code': 'section_divergence',
        'title': '三、主要分歧点归纳',
        'heading': '## 三、主要分歧点归纳',
        'phase': '生成主要分歧点归纳',
        'pct': 60,
    },
    {
        'code': 'section_recommendation',
        'title': '四、可行性结论与建议',
        'heading': '## 四、可行性结论与建议',
        'phase': '生成可行性结论与建议',
        'pct': 78,
    },
]


def strip_report_markers(text):
    cleaned = (text or '').strip()
    if not cleaned:
        return ''
    cleaned = cleaned.replace("=====REPORT_START=====", '').replace("=====REPORT_END=====", '').strip()
    return cleaned


def build_report_context_text(report_inputs):
    scenario = report_inputs['scenario']
    selling_points = safe_json_loads(scenario.core_selling_points, [])
    assumptions = safe_json_loads(scenario.validation_assumptions, [])
    summary_text = report_inputs.get('discussion_summary_text', '')
    summary_block = ''
    if summary_text:
        summary_block = (
            "结构化讨论摘要（优先基于此生成正式报告）：\n"
            f"{summary_text}\n\n"
        )

    return (
        f"请根据下面的虚拟焦点小组讨论材料，产出面向企业决策的中文分析报告。\n\n"
        f"产品名称：{scenario.product_name}\n"
        f"产品概念：{scenario.product_concept}\n"
        f"核心卖点：{'、'.join(selling_points) or '待补充'}\n"
        f"研究目标：{scenario.research_goal or '判断该概念是否值得继续推进'}\n"
        f"核心决策问题：{scenario.decision_problem or '当前阶段是否建议继续投入验证'}\n"
        f"目标用户范围：{scenario.target_user_profile or '待补充'}\n"
        f"竞品/替代方案：{scenario.competitor_context or '待补充'}\n"
        f"待验证假设：{'、'.join(assumptions) or '待补充'}\n\n"
        f"参与者画像：\n{chr(10).join(report_inputs['participant_profiles']) or '暂无参与者画像'}\n\n"
        f"外部证据：\n{report_inputs['external_evidence_text'] or '暂无外部证据输入。'}\n\n"
        f"{summary_block}"
        f"原始讨论内容（仅作兜底参考）：\n{report_inputs['conversation_text']}"
    )


def build_report_prompt(scenario, participant_profiles, conversation_text, external_evidence_text):
    selling_points = safe_json_loads(scenario.core_selling_points, [])
    assumptions = safe_json_loads(scenario.validation_assumptions, [])
    return f"""请根据以下虚拟焦点访谈材料，生成一份可直接支持企业决策的结构化分析报告。

产品信息：
产品名称：{scenario.product_name}
产品概念：{scenario.product_concept}
核心卖点：{'、'.join(selling_points) or '待补充'}
研究目标：{scenario.research_goal or '判断该概念是否值得继续推进'}
核心决策问题：{scenario.decision_problem or '当前阶段是否应继续投入验证'}
目标用户范围：{scenario.target_user_profile or '待补充'}
竞品/替代方案：{scenario.competitor_context or '待补充'}
待验证假设：{'、'.join(assumptions) or '待补充'}

参与者画像：
{chr(10).join(participant_profiles)}

外部证据：
{external_evidence_text}

讨论内容：
{conversation_text}

只允许使用以下 4 个大标题，并且严格按顺序输出：
## 一、核心观点总结
## 二、关键讨论点分析
## 三、主要分歧点归纳
## 四、可行性结论与建议

要求：
1. 第一部分要像管理层摘要，明确给出：核心结论、目标人群、推进判断、关键风险、建议下一步。
2. 第二部分至少覆盖 4 个关键议题，每个议题都说明：讨论中出现的观点、系统推导出的判断、仍需外部验证的假设。
3. 第三部分先给出分歧表格，再用文字解释，并明确当前更适合优先采取哪一侧方案。
4. 第四部分必须写成可执行建议，每条建议尽量包含：依据、风险、下一步验证动作、停止条件。
5. 要尽量引用参与者姓名、具体观点、讨论片段或外部证据，避免空泛总结。
6. 不要新增其他二级标题，可以使用项目符号、编号、小标题加粗和表格。
7. 在报告正文开始处输出 =====REPORT_START===== ，结束处输出 =====REPORT_END===== 。
8. 只输出报告正文，不要输出说明文字。
请使用中文。"""


def build_discussion_summary_prompt(report_inputs):
    scenario = report_inputs['scenario']
    selling_points = safe_json_loads(scenario.core_selling_points, [])
    assumptions = safe_json_loads(scenario.validation_assumptions, [])
    return f"""请先把下面的焦点小组讨论材料压缩成一份结构化讨论摘要，供后续正式报告生成使用。

产品名称：{scenario.product_name}
产品概念：{scenario.product_concept}
核心卖点：{'、'.join(selling_points) or '待补充'}
研究目标：{scenario.research_goal or '判断该概念是否值得继续推进'}
核心决策问题：{scenario.decision_problem or '当前阶段是否应继续投入验证'}
目标用户范围：{scenario.target_user_profile or '待补充'}
竞品/替代方案：{scenario.competitor_context or '待补充'}
待验证假设：{'、'.join(assumptions) or '待补充'}

参与者画像：
{chr(10).join(report_inputs['participant_profiles']) or '暂无参与者画像'}

外部证据：
{report_inputs['external_evidence_text'] or '暂无外部证据输入。'}

原始讨论内容：
{report_inputs['conversation_text']}

请严格输出合法 JSON 对象，不要输出解释文字，不要输出 Markdown。

字段结构固定如下：
{{
  "core_conclusion": "一句话核心结论",
  "continue_recommendation": "必须明确表达建议继续推进/暂停推进/不建议推进",
  "target_users": ["目标用户1", "目标用户2"],
  "key_topics": [
    {{
      "topic": "议题名称",
      "discussion_points": ["讨论中出现的关键观点"],
      "system_judgment": "系统判断",
      "open_questions": ["仍待验证的问题"]
    }}
  ],
  "divergences": [
    {{
      "dimension": "分歧维度",
      "side_a": "观点A",
      "side_b": "观点B",
      "current_bias": "当前更偏向哪一侧",
      "pending_validation": "仍待验证点"
    }}
  ],
  "key_risks": ["关键风险1", "关键风险2"],
  "next_steps": ["下一步动作1", "下一步动作2"]
}}

要求：
1. key_topics 至少 4 项。
2. continue_recommendation 必须明确表达继续/暂停/不建议推进。
3. 尽量引用真实讨论中的具体顾虑、场景、判断，不要空泛。
4. 不要编造不存在的人名、证据或场景。"""


def validate_discussion_summary(summary_obj):
    if not isinstance(summary_obj, dict):
        return False, "讨论摘要不是合法对象"

    continue_recommendation = str(summary_obj.get('continue_recommendation') or '').strip()
    if not continue_recommendation or not re.search(r'继续推进|继续验证|继续投入|暂停推进|不建议推进', continue_recommendation):
        return False, "讨论摘要必须明确表达是否继续推进"

    key_topics = summary_obj.get('key_topics')
    if not isinstance(key_topics, list) or len(key_topics) < 4:
        return False, "讨论摘要至少需要 4 个关键议题"

    for item in key_topics[:4]:
        if not isinstance(item, dict):
            return False, "关键议题格式无效"
        if not str(item.get('topic') or '').strip():
            return False, "关键议题缺少 topic"
        if not normalize_text_list(item.get('discussion_points')):
            return False, "关键议题缺少 discussion_points"
        if not str(item.get('system_judgment') or '').strip():
            return False, "关键议题缺少 system_judgment"
        if not normalize_text_list(item.get('open_questions')):
            return False, "关键议题缺少 open_questions"

    divergences = summary_obj.get('divergences')
    if not isinstance(divergences, list) or not divergences:
        return False, "讨论摘要至少需要 1 个分歧点"

    key_risks = normalize_text_list(summary_obj.get('key_risks'))
    next_steps = normalize_text_list(summary_obj.get('next_steps'))
    if not key_risks:
        return False, "讨论摘要缺少关键风险"
    if not next_steps:
        return False, "讨论摘要缺少下一步动作"
    return True, ""


def generate_discussion_summary(report_inputs):
    prompt = build_discussion_summary_prompt(report_inputs)
    try:
        result = call_deepseek(
            prompt,
            "你是一位企业研究总监，擅长先把冗长讨论压缩成可靠的结构化摘要，并且必须严格输出合法 JSON。"
        )
    except Exception as exc:
        return None, str(exc)
    if is_ai_error_text(result):
        return None, result

    parsed = safe_json_loads(extract_json_object_block(result), None)
    if not isinstance(parsed, dict):
        return None, "讨论摘要无法解析为合法 JSON"

    is_valid, error_message = validate_discussion_summary(parsed)
    if not is_valid:
        return None, error_message
    return parsed, None


def format_discussion_summary_for_report(summary_obj):
    if not isinstance(summary_obj, dict):
        return ''

    topic_lines = []
    for index, item in enumerate(summary_obj.get('key_topics', []), start=1):
        if not isinstance(item, dict):
            continue
        discussion_points = normalize_text_list(item.get('discussion_points'))
        open_questions = normalize_text_list(item.get('open_questions'))
        topic_lines.append(
            "\n".join(
                [
                    f"议题{index}：{str(item.get('topic') or '').strip()}",
                    f"- 讨论观点：{'；'.join(discussion_points) or '待补充'}",
                    f"- 系统判断：{str(item.get('system_judgment') or '').strip() or '待补充'}",
                    f"- 待验证问题：{'；'.join(open_questions) or '待补充'}",
                ]
            )
        )

    divergence_lines = []
    for item in summary_obj.get('divergences', []):
        if not isinstance(item, dict):
            continue
        divergence_lines.append(
            f"- {str(item.get('dimension') or '').strip()}：{str(item.get('side_a') or '').strip()} vs "
            f"{str(item.get('side_b') or '').strip()}；当前偏向={str(item.get('current_bias') or '').strip()}；"
            f"待验证={str(item.get('pending_validation') or '').strip()}"
        )

    return "\n".join(
        [
            f"核心结论：{str(summary_obj.get('core_conclusion') or '').strip() or '待补充'}",
            f"推进建议：{str(summary_obj.get('continue_recommendation') or '').strip() or '待补充'}",
            f"目标用户：{'、'.join(normalize_text_list(summary_obj.get('target_users'))) or '待补充'}",
            "关键议题：",
            "\n\n".join(topic_lines) or "待补充",
            "主要分歧：",
            "\n".join(divergence_lines) or "待补充",
            f"关键风险：{'；'.join(normalize_text_list(summary_obj.get('key_risks'))) or '待补充'}",
            f"下一步动作：{'；'.join(normalize_text_list(summary_obj.get('next_steps'))) or '待补充'}",
        ]
    )


def build_report_section_prompt(report_inputs, section_spec, completed_sections=None, retry_reason=''):
    completed_sections = completed_sections or []
    requirements_map = {
        'section_summary': [
            '必须明确给出核心结论、目标用户、推进判断、关键风险、下一步动作。',
            '请使用“**1）核心结论** / **2）目标用户** / **3）推进判断** / **4）关键风险** / **5）下一步动作**”的编号结构。',
            '每个编号标题必须单独成行，每个编号小节至少展开 2 句，不要只写一句判断。',
            '要像管理层摘要，判断要直接，不能模糊。',
        ],
        'section_analysis': [
            '至少覆盖 4 个关键议题。',
            '每个议题都写出讨论中的观点、系统判断、待验证假设。',
            '优先使用“**议题1：...** / **议题2：...** / **议题3：...** / **议题4：...**”的编号结构。',
            '每个议题下的三项内容都尽量写成 2 句左右，说明观点背后的业务含义或验证方向。',
        ],
        'section_divergence': [
            '必须写明分歧双方、当前偏向、仍待验证点。',
            '优先先给出分歧表格，再给文字解释。',
            '表格列名尽量简短，优先使用“维度 / 观点A / 观点B / 当前偏向 / 验证重点”，每个单元格控制为短句，避免一格塞入整段长文。',
            '表格下方请再使用“**1）分歧焦点** / **2）当前偏向** / **3）验证重点**”三段编号说明，每段至少 2 句。',
        ],
        'section_recommendation': [
            '必须明确回答是否建议继续推进。',
            '必须写出原因、下一步动作、关键风险、停止条件。',
            '请使用“**1）是否继续推进** / **2）核心原因** / **3）建议动作** / **4）关键风险** / **5）停止条件**”的编号结构。',
            '每个编号标题单独成行，每个编号小节至少展开 2 句，不要把多个编号内容挤在同一行。',
            '建议要可执行，最好带时间、范围或量化验证方式。',
        ],
    }
    previous_sections_text = "\n\n".join(
        [f"{item['heading']}\n{item['body']}" for item in completed_sections]
    ).strip()
    requirements_text = "\n".join(
        [f"{index}. {item}" for index, item in enumerate(requirements_map.get(section_spec['code'], []), start=1)]
    )
    retry_text = f"\n上一次输出未通过校验，问题是：{retry_reason}\n请只修正当前章节，不要重写其他章节。\n" if retry_reason else ""

    prompt_parts = [
        build_report_context_text(report_inputs),
        "",
        f"当前只生成这一节：{section_spec['heading']}",
        "输出要求：",
        requirements_text,
        "如果使用编号，请让每个编号标题独占一行，再在下一行开始写具体解释。",
        "不要输出其他一级标题。",
        "不要输出 =====REPORT_START===== 或 =====REPORT_END=====。",
        "只输出当前章节正文，不要写解释说明。",
    ]
    if previous_sections_text:
        prompt_parts.extend(["", "前面已经确认通过的章节：", previous_sections_text])
    if retry_text:
        prompt_parts.extend(["", retry_text.strip()])
    return "\n".join(prompt_parts)


def extract_report_section_body(text, section_spec):
    cleaned = strip_report_markers(text)
    if not cleaned:
        return ''

    cleaned = re.sub(rf'^\s*{re.escape(section_spec["heading"])}\s*', '', cleaned, count=1, flags=re.M).strip()
    parsed_sections = parse_report_sections(cleaned)
    if parsed_sections:
        for title, body in parsed_sections:
            if section_spec['title'] in title:
                return (body or '').strip()
        if len(parsed_sections) == 1:
            return (parsed_sections[0][1] or '').strip()
    return cleaned


def count_report_topics(text):
    matches = re.findall(r'议题\s*\d+|主题\s*\d+|\*\*议题\s*\d+|\*\*主题\s*\d+', text or '')
    normalized = {item.replace('*', '').replace(' ', '') for item in matches}
    return len(normalized)


def validate_report_section(section_spec, body):
    text = (body or '').strip()
    if not text:
        return False, "当前章节内容为空"

    if section_spec['code'] == 'section_summary':
        for pattern in [r'核心结论|结论', r'推进判断|建议继续|建议推进|是否建议', r'风险', r'下一步']:
            if not re.search(pattern, text):
                return False, "核心观点总结缺少关键结论、风险或下一步动作"
        return True, ""

    if section_spec['code'] == 'section_analysis':
        if count_report_topics(text) < 4:
            return False, "关键讨论点分析至少需要 4 个议题"
        for pattern, error_message in [
            (r'讨论中的观点|讨论观点|观点', "关键讨论点分析缺少讨论观点"),
            (r'系统判断|判断', "关键讨论点分析缺少系统判断"),
            (r'待验证假设|待验证|假设', "关键讨论点分析缺少待验证假设"),
        ]:
            if not re.search(pattern, text):
                return False, error_message
        return True, ""

    if section_spec['code'] == 'section_divergence':
        if not re.search(r'分歧', text):
            return False, "主要分歧点归纳缺少分歧描述"
        if '|' not in text:
            return False, "主要分歧点归纳缺少分歧表格"
        if not re.search(r'当前偏向|偏向', text):
            return False, "主要分歧点归纳缺少当前偏向"
        if not re.search(r'待验证点|待验证', text):
            return False, "主要分歧点归纳缺少待验证点"
        return True, ""

    if section_spec['code'] == 'section_recommendation':
        if not re.search(r'是否建议继续推进|建议继续推进|建议推进|继续进行|继续验证|继续投入|不建议推进|暂停推进', text):
            return False, "可行性结论与建议必须明确回答是否继续推进"
        for pattern, error_message in [
            (r'下一步动作|下一步|验证动作|立即', "可行性结论与建议缺少下一步动作"),
            (r'风险|关键风险', "可行性结论与建议缺少关键风险"),
            (r'停止条件|停止|终止条件', "可行性结论与建议缺少停止条件"),
        ]:
            if not re.search(pattern, text):
                return False, error_message
        return True, ""

    return True, ""


def assemble_report_content(section_bodies):
    chunks = ["=====REPORT_START====="]
    for section_spec in REPORT_SECTION_SPECS:
        chunks.append(section_spec['heading'])
        chunks.append((section_bodies.get(section_spec['code']) or '').strip())
    chunks.append("=====REPORT_END=====")
    return "\n\n".join(chunks)


def validate_report_structure(report_content):
    sections = parse_report_sections(report_content or '')
    if len(sections) != len(REPORT_SECTION_SPECS):
        return False, "报告缺少必要章节"

    for index, section_spec in enumerate(REPORT_SECTION_SPECS):
        title, body = sections[index]
        if section_spec['title'] not in title:
            return False, "报告章节顺序不正确"
        is_valid, error_message = validate_report_section(section_spec, body)
        if not is_valid:
            return False, error_message
    return True, ""


def generate_report_section(report_inputs, section_spec, completed_sections=None):
    completed_sections = completed_sections or []
    last_error = ''
    for _ in range(3):
        prompt = build_report_section_prompt(
            report_inputs,
            section_spec,
            completed_sections=completed_sections,
            retry_reason=last_error,
        )
        result = call_deepseek(
            prompt,
            "你是一位面向企业管理层的研究分析顾问，擅长把讨论材料整理成可执行的决策报告。"
        )
        if is_ai_error_text(result):
            return None, result

        section_body = extract_report_section_body(result, section_spec)
        is_valid, error_message = validate_report_section(section_spec, section_body)
        if is_valid:
            return section_body, None
        last_error = error_message

    return None, f"{section_spec['title']} 生成失败：{last_error or '内容校验未通过'}"


def build_report_generation_context(report_inputs):
    return {
        'scenario': report_inputs['scenario'],
        'participant_profiles': report_inputs['participant_profiles'],
        'conversation_text': report_inputs['conversation_text'],
        'external_evidence_text': report_inputs['external_evidence_text'],
        'external_evidence': report_inputs['external_evidence'],
        'discussion_summary': report_inputs.get('discussion_summary'),
        'discussion_summary_text': report_inputs.get('discussion_summary_text', ''),
        'current_date': report_inputs['current_date'],
        'report_title': f"{report_inputs['scenario'].product_name} - 产品分析报告 ({report_inputs['current_date']})",
        'section_bodies': {},
        'completed_sections': [],
    }


def build_report_progress_dispatcher(progress_callback=None):
    emitted_stage_codes = set()

    def emit(stage_code, phase, pct, allow_repeat=False):
        if not progress_callback:
            return
        if not allow_repeat and stage_code in emitted_stage_codes:
            return
        emitted_stage_codes.add(stage_code)
        progress_callback(stage_code, phase, pct)

    return emit


def build_local_fallback_report_bundle(report_inputs, progress_callback=None, reason='no_runtime_ai_config', upstream_error=None, emit_transition=False, emit_prepare_stage=True):
    scenario = report_inputs['scenario']
    scenario_snapshot = report_inputs.get('scenario_snapshot') or {}
    scenario_id = report_inputs.get('scenario_id') or report_inputs.get('scenario_db_id') or scenario_snapshot.get('id')
    log_report_event(
        'report_generation_fallback_started',
        scenario_id=scenario_id,
        mode='local_fallback',
        reason=reason,
        upstream_error=upstream_error
    )
    if emit_transition and progress_callback:
        progress_callback("fallback_started", "AI 结果未通过校验，切换为稳定报告生成", 82)
    topics = normalize_text_list(scenario_snapshot.get('discussion_topics'))
    assumptions = normalize_text_list(scenario_snapshot.get('validation_assumptions'))
    participants = report_inputs.get('participants') or []
    records = report_inputs.get('records') or []
    external_evidence = report_inputs.get('external_evidence') or []
    product_name = scenario_snapshot.get('product_name') or '当前产品'
    research_plan = scenario_snapshot.get('research_plan') or ''
    report_title = f"{product_name} - 产品分析报告 ({report_inputs['current_date']})"

    if progress_callback and emit_prepare_stage:
        progress_callback("prepare_inputs", "准备讨论材料", 8)

    topic_lines = []
    for index, topic in enumerate(topics[:4], start=1):
        matching_records = [record.content.strip() for record in records if topic and topic in (record.content or '')]
        sample_point = matching_records[0] if matching_records else (records[min(index - 1, len(records) - 1)].content.strip() if records else "当前讨论记录较少，建议补充更多一线反馈。")
        topic_lines.append(
            "\n".join(
                [
            f"**议题{index}：{topic}**",
                    f"- 讨论中的观点：{sample_point[:120]}",
                    f"- 系统判断：该议题已经被参与者明确提及，需要在后续验证中继续量化。",
                    f"- 待验证假设：围绕“{topic}”的判断仍需要更多真实场景数据支持。",
                ]
            )
        )
    while len(topic_lines) < 4:
        fallback_topic = assumptions[len(topic_lines)] if len(assumptions) > len(topic_lines) else f"关键验证点{len(topic_lines) + 1}"
        topic_lines.append(
            "\n".join(
                [
                    f"**议题{len(topic_lines) + 1}：{fallback_topic}**",
                    f"- 讨论中的观点：现有材料对“{fallback_topic}”涉及较少，说明这部分还需要进一步访谈。",
                    f"- 系统判断：该议题会直接影响继续推进的判断。",
                    f"- 待验证假设：若该议题无法被验证，当前结论的可信度会明显下降。",
                ]
            )
        )

    risk_seed = assumptions[:2] or topics[:2] or ["使用场景尚未验证", "投入产出比还不稳定"]
    action_seed = topics[:2] or assumptions[:2] or ["补充一线用户访谈", "安排小范围试点"]
    evidence_seed = external_evidence[0].content[:80] if external_evidence else (records[0].content[:80] if records else "现有证据主要来自有限讨论记录")
    participant_names = "、".join([item.persona_name for item in participants[:3]]) or "当前参与者"

    summary_section = "\n".join(
        [
            f"- 核心结论：围绕“{product_name}”的讨论显示该方案具备继续验证价值。",
            f"- 推进判断：建议继续推进，但以小范围验证而不是直接全面投入为主。",
            f"- 关键风险：{risk_seed[0]}；如果处理不当，会直接削弱继续推进的把握。",
            f"- 下一步：优先围绕“{action_seed[0]}”安排补充验证，并整理更稳定的证据链。",
        ]
    )

    divergence_table = "\n".join(
        [
            "| 分歧维度 | 观点A | 观点B | 当前偏向 | 待验证点 |",
            "| --- | --- | --- | --- | --- |",
            f"| 推进节奏 | 先快速落地 | 先补充验证 | 先补充验证 | {risk_seed[0]}是否能在试点中被压低 |",
            "",
            f"- 分歧双方：部分参与者倾向尽快推进，另一部分参与者更担心{risk_seed[-1]}。",
            "- 当前偏向：先在可控范围内验证，再决定是否扩大投入。",
            f"- 待验证点：{evidence_seed}",
        ]
    )

    recommendation_section = "\n".join(
        [
            "- 是否建议继续推进：建议继续推进，但仅限于下一轮小范围验证。",
            f"- 原因：现有讨论已经表明 {product_name} 具备明确价值点，但证据强度仍不足以支持直接全面扩张。",
            f"- 下一步动作：1) 围绕“{action_seed[0]}”安排补充访谈；2) 用试点观察{action_seed[-1]}的真实表现。",
            f"- 关键风险：{risk_seed[0]}；{risk_seed[-1]}。",
            "- 停止条件：如果下一轮验证仍无法证明用户接受度和使用价值，则暂停继续投入。",
        ]
    )

    section_map = {
        'section_summary': summary_section,
        'section_analysis': "\n\n".join(topic_lines[:4]),
        'section_divergence': divergence_table,
        'section_recommendation': recommendation_section,
    }

    for section_spec in REPORT_SECTION_SPECS:
        if progress_callback:
            progress_callback(section_spec['code'], section_spec['phase'], section_spec['pct'])

    report_content = sanitize_report_text(assemble_report_content(section_map))
    metadata = enrich_report_metadata(
        extract_report_metadata(report_content, external_evidence),
        external_evidence
    )
    executive_summary = (
        f"基于 {participant_names} 的讨论，系统建议继续推进 {product_name}，"
        f"但应先完成小范围验证，重点确认 {risk_seed[0]} 与 {action_seed[0]}。"
    )

    if progress_callback:
        progress_callback("finalize_report", "整理摘要并保存报告", 92)

    report = persist_analysis_report(
        scenario_id,
        report_title,
        report_content,
        executive_summary,
        metadata,
        research_plan=research_plan
    )
    log_report_event(
        'report_generation_fallback_completed',
        scenario_id=scenario_id,
        mode='local_fallback',
        report_id=getattr(report, 'id', None)
    )
    return report, None


def advance_report_generation(state, progress_callback=None):
    if progress_callback and not state.get('_prepare_emitted'):
        progress_callback("prepare_inputs", "准备讨论材料", 8)
        state['_prepare_emitted'] = True

    section_index = len(state['completed_sections'])
    if section_index >= len(REPORT_SECTION_SPECS):
        return None, None, None

    section_spec = REPORT_SECTION_SPECS[section_index]
    if progress_callback:
        progress_callback(section_spec['code'], section_spec['phase'], section_spec['pct'])

    section_body, section_error = generate_report_section(
        state,
        section_spec,
        completed_sections=state['completed_sections'],
    )
    if section_error:
        return None, section_error, None

    state['section_bodies'][section_spec['code']] = section_body
    state['completed_sections'].append({
        'heading': section_spec['heading'],
        'body': section_body,
    })
    return section_spec, None, section_body


def build_report_from_generation_state(report_inputs, state, progress_callback=None):
    report_content = sanitize_report_text(assemble_report_content(state['section_bodies']))
    if not report_content:
        return None, "报告正文结构不完整，请稍后重试"

    structure_ok, structure_error = validate_report_structure(report_content)
    if not structure_ok:
        return None, structure_error

    if progress_callback:
        progress_callback("finalize_report", "整理摘要并保存报告", 92)

    executive_summary = build_executive_summary_from_report(report_content)
    raw_metadata = extract_report_metadata(report_content, report_inputs['external_evidence'])
    metadata = enrich_report_metadata(raw_metadata, report_inputs['external_evidence'])
    report = persist_analysis_report(
        report_inputs['scenario'].id,
        state['report_title'],
        report_content,
        executive_summary,
        metadata,
        research_plan=report_inputs['scenario'].research_plan or '',
        discussion_summary=report_inputs.get('discussion_summary')
    )
    return report, None


def build_report_inputs(scenario_id):
    scenario, error_response = get_scenario_or_403(scenario_id, allow_meeting_member=True)
    if error_response:
        return None, error_response

    ensure_extended_schema()

    records = ConversationRecord.query.filter_by(scenario_id=scenario_id).order_by(ConversationRecord.timestamp).all()
    if not records:
        return None, app.response_class(
            response=json.dumps({"error": "请先进行对话模拟"}, ensure_ascii=False),
            status=400,
            mimetype='application/json; charset=utf-8'
        )

    participants = VirtualParticipant.query.filter_by(scenario_id=scenario_id).all()
    participant_profiles = []
    participant_map = {}
    for participant in participants:
        participant_map[participant.id] = participant
        normalize_participant_profile(participant, scenario)
        participant_profiles.append(
            f"{participant.persona_name}：标签={'/'.join(safe_json_loads(participant.persona_tags, [])) or '待补充'}；"
            f"背景={participant.background}；使用目标={participant.usage_goal or '待观察'}；"
            f"预算敏感度={participant.budget_sensitivity or '中'}；品牌偏好={participant.brand_preference or '暂无明显偏好'}；"
            f"风险厌恶={participant.risk_aversion or '中'}；决策风格={participant.decision_style or '谨慎比较'}；"
            f"稳定立场={participant.stance_summary or '待形成'}"
        )
    db.session.commit()

    conversation_lines = []
    for record in records:
        if record.is_host:
            participant_name = "主持人"
        elif record.participant_id == -1:
            participant_name = "用户"
        else:
            participant = participant_map.get(record.participant_id)
            participant_name = participant.persona_name if participant else "未知参与者"
        conversation_lines.append(f"{participant_name}：{record.content}")

    external_evidence = ExternalEvidence.query.filter_by(scenario_id=scenario_id).order_by(ExternalEvidence.created_at.desc()).all()
    external_evidence_lines = [
        f"- {item.title}｜来源：{item.source_label or item.source_title or item.evidence_type}｜强度：{normalize_strength_level(item.strength_level)}｜内容：{item.content}"
        for item in external_evidence[:8]
    ]

    return {
        "scenario": scenario,
        "scenario_id": scenario_id,
        "scenario_db_id": scenario.id,
        "scenario_snapshot": {
            "id": scenario.id,
            "product_name": scenario.product_name,
            "discussion_topics": safe_json_loads(scenario.discussion_topics, []),
            "validation_assumptions": safe_json_loads(scenario.validation_assumptions, []),
            "research_plan": scenario.research_plan or '',
        },
        "records": records,
        "participants": participants,
        "participant_profiles": participant_profiles,
        "conversation_text": "\n".join(conversation_lines),
        "external_evidence": external_evidence,
        "external_evidence_text": "\n".join(external_evidence_lines) or "暂无外部证据输入。",
        "current_date": datetime.now().strftime("%Y年%m月%d日")
    }, None


def build_executive_summary_from_report(report_content):
    if not report_content:
        return ''
    prompt = f"""请根据以下报告正文，写一段 90-150 字的管理层执行摘要。
要求：
1. 明确是否建议继续推进。
2. 点出最大机会或最大风险。
3. 点出最该先做的动作。
4. 只输出摘要正文。

报告内容：
{report_content}"""
    summary = call_deepseek(prompt, "你是一位企业研究总监，擅长把长报告压缩成管理层摘要。")
    if is_ai_error_text(summary):
        return ''
    return (summary or '').strip()


def enrich_report_metadata(metadata, external_evidence):
    metadata = dict(metadata or {})
    evidence_items = metadata.get('evidence_items') or []
    decision_risks = metadata.get('decision_risks') or []
    recommended_actions = metadata.get('recommended_actions') or []
    key_assumptions = metadata.get('key_assumptions') or []
    source_breakdown = metadata.get('source_breakdown') or []

    normalized_source_breakdown = []
    for item in source_breakdown:
        if isinstance(item, dict):
            normalized_source_breakdown.append({
                'section': item.get('section', ''),
                'summary': item.get('summary', ''),
                'source_type': item.get('source_type', 'discussion'),
                'strength_level': normalize_strength_level(item.get('strength_level') or '')
            })

    for item in external_evidence or []:
        if getattr(item, 'content', None):
            evidence_items.append(item.content[:180])
        normalized_source_breakdown.append({
            'section': item.title,
            'summary': item.content[:180],
            'source_type': 'external',
            'strength_level': normalize_strength_level(item.strength_level)
        })

    dedup = lambda seq: [x for i, x in enumerate(seq) if x and x not in seq[:i]]

    normalized_evidence = dedup([str(item).strip() for item in evidence_items if str(item).strip()])[:8]
    normalized_risks = dedup([str(item).strip() for item in decision_risks if str(item).strip()])[:6]
    normalized_actions = dedup([str(item).strip() for item in recommended_actions if str(item).strip()])[:8]
    normalized_assumptions = dedup([str(item).strip() for item in key_assumptions if str(item).strip()])[:6]

    return {
        'key_assumptions': normalized_assumptions,
        'evidence_items': normalized_evidence,
        'decision_risks': normalized_risks,
        'recommended_actions': normalized_actions,
        'confidence_level': normalize_report_confidence_level(
            metadata.get('confidence_level'),
            has_external_evidence=bool(external_evidence),
            evidence_count=len(normalized_evidence)
        ),
        'source_breakdown': normalized_source_breakdown[:12]
    }


def persist_analysis_report(scenario_id, report_title, report_content, executive_summary, metadata, research_plan='', discussion_summary=None):
    existing_report = AnalysisReport.query.filter_by(scenario_id=scenario_id).first()
    if existing_report:
        existing_report.report_title = report_title
        existing_report.content = report_content
        existing_report.executive_summary = executive_summary
        existing_report.key_assumptions = json_dumps(metadata['key_assumptions'], [])
        existing_report.evidence_items = json_dumps(metadata['evidence_items'], [])
        existing_report.decision_risks = json_dumps(metadata['decision_risks'], [])
        existing_report.recommended_actions = json_dumps(metadata['recommended_actions'], [])
        existing_report.confidence_level = metadata['confidence_level']
        existing_report.source_breakdown = json_dumps(metadata['source_breakdown'], [])
        existing_report.discussion_summary = json_dumps(discussion_summary, {}) if discussion_summary else None
        existing_report.generated_at = datetime.now()
        report = existing_report
    else:
        report = AnalysisReport(
            scenario_id=scenario_id,
            report_title=report_title,
            content=report_content,
            executive_summary=executive_summary,
            key_assumptions=json_dumps(metadata['key_assumptions'], []),
            evidence_items=json_dumps(metadata['evidence_items'], []),
            decision_risks=json_dumps(metadata['decision_risks'], []),
            recommended_actions=json_dumps(metadata['recommended_actions'], []),
            confidence_level=metadata['confidence_level'],
            source_breakdown=json_dumps(metadata['source_breakdown'], []),
            discussion_summary=json_dumps(discussion_summary, {}) if discussion_summary else None
        )
        db.session.add(report)
    db.session.commit()

    next_round_index = ResearchRound.query.filter_by(scenario_id=scenario_id).count() + 1
    round_item = ResearchRound(
        scenario_id=scenario_id,
        round_index=next_round_index,
        plan_snapshot=research_plan or '',
        report_snapshot=report_content,
        summary=executive_summary or report_title
    )
    db.session.add(round_item)
    db.session.commit()
    log_report_event(
        'report_persisted',
        scenario_id=scenario_id,
        report_id=getattr(report, 'id', None),
        round_index=next_round_index,
        confidence_level=metadata.get('confidence_level')
    )
    return report


# ---------------------------
# 数据库模型定义
# ---------------------------

class ProductScenario(db.Model):
    """产品场景模型 - 存储焦点小组讨论的产品信息和讨论配置"""
    __tablename__ = 'product_scenario'
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)  # 关联用户ID
    product_name = db.Column(db.String(100), nullable=False)  # 产品名称
    product_concept = db.Column(db.Text, nullable=False)  # 产品概念描述
    core_selling_points = db.Column(db.Text, nullable=False)  # 核心卖点（JSON数组）
    discussion_topics = db.Column(db.Text, nullable=False)  # 讨论话题（JSON数组）
    occasion_type = db.Column(db.String(50), nullable=False, default='focus_group')  # 场景类型
    occasion_description = db.Column(db.Text, nullable=False, default='标准焦点小组讨论')  # 场景描述
    research_goal = db.Column(db.Text, nullable=True)  # 研究目标
    decision_problem = db.Column(db.Text, nullable=True)  # 决策问题
    target_user_profile = db.Column(db.Text, nullable=True)  # 目标用户范围
    competitor_context = db.Column(db.Text, nullable=True)  # 竞品/替代方案
    validation_assumptions = db.Column(db.Text, nullable=True)  # 需要验证的假设（JSON数组）
    research_plan = db.Column(db.Text, nullable=True)  # 研究计划（Markdown）
    research_plan_status = db.Column(db.String(20), nullable=False, default='pending')  # 研究计划状态
    created_at = db.Column(db.DateTime, default=datetime.now)  # 创建时间

    meeting_status = db.Column(db.String(20), nullable=False, default='standalone')

    def to_dict(self):
        """将模型转换为字典格式，便于JSON序列化"""
        return {
            'id': self.id,
            'user_id': self.user_id,
            'product_name': self.product_name,
            'product_concept': self.product_concept,
            'core_selling_points': safe_json_loads(self.core_selling_points, []),
            'discussion_topics': safe_json_loads(self.discussion_topics, []),
            'occasion_type': self.occasion_type,
            'occasion_description': self.occasion_description,
            'research_goal': self.research_goal or '',
            'decision_problem': self.decision_problem or '',
            'target_user_profile': self.target_user_profile or '',
            'competitor_context': self.competitor_context or '',
            'validation_assumptions': safe_json_loads(self.validation_assumptions, []),
            'research_plan': self.research_plan or '',
            'research_plan_status': self.research_plan_status or 'pending',
            'meeting_status': normalize_meeting_status_value(self.meeting_status, 'standalone'),
            'created_at': self.created_at.isoformat()
        }

class MeetingRoom(db.Model):
    __tablename__ = 'meeting_room'
    id = db.Column(db.Integer, primary_key=True)
    scenario_id = db.Column(db.Integer, db.ForeignKey('product_scenario.id'), nullable=False)
    owner_user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    room_name = db.Column(db.String(160), nullable=False)
    room_code = db.Column(db.String(20), unique=True, nullable=False, index=True)
    invite_token = db.Column(db.String(64), unique=True, nullable=False, index=True)
    status = db.Column(db.String(20), nullable=False, default='waiting')
    topic_title = db.Column(db.Text, nullable=False)
    topic_notes = db.Column(db.Text, nullable=True)
    target_count = db.Column(db.Integer, nullable=False, default=4)
    created_at = db.Column(db.DateTime, default=datetime.now)
    started_at = db.Column(db.DateTime, nullable=True)
    updated_at = db.Column(db.DateTime, default=datetime.now, onupdate=datetime.now)
    discussion_phase = db.Column(db.String(20), nullable=False, default='waiting')
    turn_order = db.Column(db.Text, nullable=True)
    current_turn_index = db.Column(db.Integer, nullable=False, default=0)
    turn_number = db.Column(db.Integer, nullable=False, default=1)
    active_speaker_participant_id = db.Column(db.Integer, nullable=True)
    ai_turn_pending = db.Column(db.Boolean, nullable=False, default=False)

    def to_dict(self):
        members = MeetingRoomMember.query.filter_by(room_id=self.id).order_by(MeetingRoomMember.joined_at.asc()).all()
        current_count = len(members)
        arrived_count = sum(1 for member in members if (member.status or 'arrived') == 'arrived')
        return {
            'id': self.id,
            'scenario_id': self.scenario_id,
            'owner_user_id': self.owner_user_id,
            'room_name': self.room_name,
            'room_code': self.room_code,
            'invite_token': self.invite_token,
            'invite_link': build_room_invite_link(self.invite_token),
            'status': self.status or 'waiting',
            'discussion_phase': self.discussion_phase or ('ended' if (self.status or '') == 'ended' else 'waiting'),
            'topic_title': self.topic_title,
            'topic_notes': self.topic_notes or '',
            'target_count': self.target_count or 0,
            'current_count': current_count,
            'arrived_count': arrived_count,
            'can_start': (self.status or 'waiting') == 'waiting' and current_count > 0,
            'turn_order': safe_json_loads(self.turn_order, []),
            'current_turn_index': self.current_turn_index or 0,
            'turn_number': self.turn_number or 1,
            'active_speaker_participant_id': self.active_speaker_participant_id,
            'ai_turn_pending': bool(self.ai_turn_pending),
            'created_at': self.created_at.isoformat() if self.created_at else '',
            'started_at': self.started_at.isoformat() if self.started_at else '',
            'updated_at': self.updated_at.isoformat() if self.updated_at else ''
        }

class MeetingRoomMember(db.Model):
    __tablename__ = 'meeting_room_member'
    id = db.Column(db.Integer, primary_key=True)
    room_id = db.Column(db.Integer, db.ForeignKey('meeting_room.id'), nullable=False, index=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False, index=True)
    member_role = db.Column(db.String(20), nullable=False, default='member')
    display_name = db.Column(db.String(120), nullable=True)
    status = db.Column(db.String(20), nullable=False, default='arrived')
    joined_at = db.Column(db.DateTime, default=datetime.now)
    last_seen_at = db.Column(db.DateTime, default=datetime.now)

    __table_args__ = (
        db.UniqueConstraint('room_id', 'user_id', name='uq_meeting_room_member_room_user'),
    )

    def to_dict(self):
        return {
            'id': self.id,
            'room_id': self.room_id,
            'user_id': self.user_id,
            'member_role': self.member_role or 'member',
            'display_name': self.display_name or '',
            'status': self.status or 'arrived',
            'joined_at': self.joined_at.isoformat() if self.joined_at else '',
            'last_seen_at': self.last_seen_at.isoformat() if self.last_seen_at else ''
        }

class MeetingRoomMessage(db.Model):
    __tablename__ = 'meeting_room_message'
    id = db.Column(db.Integer, primary_key=True)
    room_id = db.Column(db.Integer, db.ForeignKey('meeting_room.id'), nullable=False, index=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)
    sender_name = db.Column(db.String(120), nullable=False)
    content = db.Column(db.Text, nullable=False)
    message_type = db.Column(db.String(20), nullable=False, default='member')
    created_at = db.Column(db.DateTime, default=datetime.now)

    def to_dict(self):
        return {
            'id': self.id,
            'room_id': self.room_id,
            'user_id': self.user_id,
            'sender_name': self.sender_name,
            'content': self.content,
            'message_type': self.message_type or 'member',
            'created_at': self.created_at.isoformat() if self.created_at else ''
        }

class VirtualParticipant(db.Model):
    """虚拟参与者模型 - 存储AI生成的参与者角色画像"""
    __tablename__ = 'virtual_participant'
    id = db.Column(db.Integer, primary_key=True)
    scenario_id = db.Column(db.Integer, db.ForeignKey('product_scenario.id'), nullable=False)  # 关联场景ID
    persona_name = db.Column(db.String(50), nullable=False)  # 角色名称
    persona_tags = db.Column(db.Text, nullable=False)  # 角色标签（JSON数组）
    personality = db.Column(db.Text, nullable=False)  # 性格特点描述
    background = db.Column(db.Text, nullable=False)  # 背景描述
    usage_goal = db.Column(db.Text, nullable=True)  # 使用目标
    budget_sensitivity = db.Column(db.Text, nullable=True)  # 预算敏感度
    brand_preference = db.Column(db.Text, nullable=True)  # 品牌偏好
    risk_aversion = db.Column(db.Text, nullable=True)  # 风险厌恶程度
    decision_style = db.Column(db.Text, nullable=True)  # 决策风格
    deal_breakers = db.Column(db.Text, nullable=True)  # 反感点（JSON数组）
    stance_summary = db.Column(db.Text, nullable=True)  # 立场摘要
    stance_state = db.Column(db.Text, nullable=True)  # 动态状态（JSON）
    is_custom = db.Column(db.Boolean, default=False)  # 是否为自定义角色
    is_ai_generated = db.Column(db.Boolean, default=False)  # 是否由AI生成
    custom_params = db.Column(db.Text, nullable=True)  # 自定义参数（JSON）
    linked_user_id = db.Column(db.Integer, nullable=True)
    meeting_member_id = db.Column(db.Integer, nullable=True)
    speaker_origin = db.Column(db.String(20), nullable=False, default='ai')
    seat_order = db.Column(db.Integer, nullable=True)
    room_managed = db.Column(db.Boolean, default=False)

    def to_dict(self):
        """将模型转换为字典格式"""
        return {
            'id': self.id,
            'scenario_id': self.scenario_id,
            'persona_name': self.persona_name,
            'persona_tags': safe_json_loads(self.persona_tags, []),
            'personality': self.personality,
            'background': self.background,
            'usage_goal': self.usage_goal or '',
            'budget_sensitivity': self.budget_sensitivity or '',
            'brand_preference': self.brand_preference or '',
            'risk_aversion': self.risk_aversion or '',
            'decision_style': self.decision_style or '',
            'deal_breakers': safe_json_loads(self.deal_breakers, []),
            'stance_summary': self.stance_summary or '',
            'stance_state': safe_json_loads(self.stance_state, {}),
            'is_custom': self.is_custom,
            'is_ai_generated': self.is_ai_generated,
            'custom_params': safe_json_loads(self.custom_params, {}),
            'linked_user_id': self.linked_user_id,
            'meeting_member_id': self.meeting_member_id,
            'speaker_origin': self.speaker_origin or ('human' if self.linked_user_id else 'ai'),
            'seat_order': self.seat_order,
            'room_managed': bool(self.room_managed)
        }

class ConversationRecord(db.Model):
    """对话记录模型 - 存储焦点小组讨论的所有对话内容"""
    __tablename__ = 'conversation_record'
    id = db.Column(db.Integer, primary_key=True)
    scenario_id = db.Column(db.Integer, db.ForeignKey('product_scenario.id'), nullable=False)  # 关联场景ID
    participant_id = db.Column(db.Integer, nullable=False)  # 参与者ID（0为主持人，-1为用户）
    content = db.Column(db.Text, nullable=False)  # 对话内容
    is_host = db.Column(db.Boolean, default=False)  # 是否为主持人发言
    message_type = db.Column(db.String(20), default='normal')  # 消息类型：intro/topic/guide/summary/conclusion/normal/user
    interaction_intent = db.Column(db.String(50), nullable=True)  # 发言意图
    timestamp = db.Column(db.DateTime, default=datetime.now)  # 时间戳

class AnalysisReport(db.Model):
    """分析报告模型 - 存储AI生成的决策分析报告"""
    __tablename__ = 'analysis_report'
    id = db.Column(db.Integer, primary_key=True)
    scenario_id = db.Column(db.Integer, db.ForeignKey('product_scenario.id'), nullable=False)  # 关联场景ID
    report_title = db.Column(db.String(200), nullable=False)  # 报告标题
    content = db.Column(db.Text, nullable=False)  # 报告内容（Markdown格式）
    executive_summary = db.Column(db.Text, nullable=True)  # 执行摘要
    key_assumptions = db.Column(db.Text, nullable=True)  # 关键假设（JSON数组）
    evidence_items = db.Column(db.Text, nullable=True)  # 证据条目（JSON数组）
    decision_risks = db.Column(db.Text, nullable=True)  # 决策风险（JSON数组）
    recommended_actions = db.Column(db.Text, nullable=True)  # 建议动作（JSON数组）
    confidence_level = db.Column(db.String(20), nullable=True)  # 置信度
    source_breakdown = db.Column(db.Text, nullable=True)  # 信息来源拆解（JSON数组）
    discussion_summary = db.Column(db.Text, nullable=True)  # 结构化讨论摘要（JSON对象）
    generated_at = db.Column(db.DateTime, default=datetime.now)  # 生成时间

    def to_dict(self):
        """将模型转换为字典格式"""
        return {
            'id': self.id,
            'scenario_id': self.scenario_id,
            'report_title': self.report_title,
            'content': self.content,
            'executive_summary': self.executive_summary or '',
            'key_assumptions': safe_json_loads(self.key_assumptions, []),
            'evidence_items': safe_json_loads(self.evidence_items, []),
            'decision_risks': safe_json_loads(self.decision_risks, []),
            'recommended_actions': safe_json_loads(self.recommended_actions, []),
            'confidence_level': self.confidence_level or '中',
            'source_breakdown': safe_json_loads(self.source_breakdown, []),
            'generated_at': self.generated_at.isoformat()
        }

class ExternalEvidence(db.Model):
    """外部证据模型 - 存储用户补充的外部数据或调研材料"""
    __tablename__ = 'external_evidence'
    id = db.Column(db.Integer, primary_key=True)
    scenario_id = db.Column(db.Integer, db.ForeignKey('product_scenario.id'), nullable=False)
    title = db.Column(db.String(200), nullable=False)
    evidence_type = db.Column(db.String(50), nullable=False, default='market_data')
    source_label = db.Column(db.String(200), nullable=True)
    source_url = db.Column(db.Text, nullable=True)
    source_title = db.Column(db.String(300), nullable=True)
    source_domain = db.Column(db.String(200), nullable=True)
    generated_by = db.Column(db.String(30), nullable=False, default='user_manual')
    content = db.Column(db.Text, nullable=False)
    strength_level = db.Column(db.String(20), nullable=False, default='高')
    created_at = db.Column(db.DateTime, default=datetime.now)

    def to_dict(self):
        return {
            'id': self.id,
            'scenario_id': self.scenario_id,
            'title': self.title,
            'evidence_type': self.evidence_type,
            'source_label': self.source_label or '',
            'source_url': self.source_url or '',
            'source_title': self.source_title or '',
            'source_domain': self.source_domain or '',
            'generated_by': self.generated_by or 'user_manual',
            'content': self.content,
            'strength_level': self.strength_level,
            'created_at': self.created_at.isoformat()
        }

class ResearchRound(db.Model):
    """研究轮次模型 - 存储每轮研究计划和报告快照"""
    __tablename__ = 'research_round'
    id = db.Column(db.Integer, primary_key=True)
    scenario_id = db.Column(db.Integer, db.ForeignKey('product_scenario.id'), nullable=False)
    round_index = db.Column(db.Integer, nullable=False, default=1)
    plan_snapshot = db.Column(db.Text, nullable=True)
    report_snapshot = db.Column(db.Text, nullable=True)
    summary = db.Column(db.Text, nullable=True)
    created_at = db.Column(db.DateTime, default=datetime.now)

    def to_dict(self):
        return {
            'id': self.id,
            'scenario_id': self.scenario_id,
            'round_index': self.round_index,
            'plan_snapshot': self.plan_snapshot or '',
            'report_snapshot': self.report_snapshot or '',
            'summary': self.summary or '',
            'created_at': self.created_at.isoformat()
        }

class User(db.Model):
    """用户模型 - 存储系统用户信息"""
    __tablename__ = 'user'
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)  # 用户名（唯一）
    email = db.Column(db.String(120), unique=True, nullable=False)  # 邮箱（唯一）
    password_hash = db.Column(db.String(200), nullable=False)  # 密码哈希值（扩展长度存pbkdf2）
    # 新增字段
    nickname = db.Column(db.String(80), nullable=True)        # 显示名
    avatar = db.Column(db.String(200), nullable=True)         # 头像URL
    company = db.Column(db.String(120), nullable=True)        # 公司
    role = db.Column(db.String(20), nullable=False, default='user')  # 用户角色
    bio = db.Column(db.Text, nullable=True)                   # 个人简介
    locale = db.Column(db.String(10), nullable=False, default='zh-CN')
    last_login = db.Column(db.DateTime, nullable=True)        # 最近登录时间
    login_count = db.Column(db.Integer, nullable=False, default=0)
    ai_endpoint_url = db.Column(db.Text, nullable=True)
    ai_model_name = db.Column(db.String(120), nullable=True)
    ai_api_keys_encrypted = db.Column(db.Text, nullable=True)
    ai_config_enabled = db.Column(db.Boolean, nullable=False, default=False)
    ai_last_test_status = db.Column(db.String(20), nullable=True)
    ai_last_test_message = db.Column(db.Text, nullable=True)
    ai_last_tested_at = db.Column(db.DateTime, nullable=True)
    created_at = db.Column(db.DateTime, default=datetime.now)  # 创建时间

    def to_dict(self):
        """将模型转换为字典格式（不包含密码）"""
        return {
            'id': self.id,
            'username': self.username,
            'email': self.email,
            'nickname': self.nickname or '',
            'avatar': self.avatar or '',
            'company': self.company or '',
            'role': self.role,
            'bio': self.bio or '',
            'locale': self.locale,
            'last_login': self.last_login.isoformat() if self.last_login else '',
            'login_count': self.login_count,
            'created_at': self.created_at.isoformat()
        }

class VerificationCode(db.Model):
    """验证码模型 - 存储短信/邮箱验证码"""
    __tablename__ = 'verification_code'
    id = db.Column(db.Integer, primary_key=True)
    channel = db.Column(db.String(10), nullable=False)       # email / sms / test
    target = db.Column(db.String(120), nullable=False)       # 邮箱地址或手机号
    code = db.Column(db.String(12), nullable=False)          # 验证码
    purpose = db.Column(db.String(30), nullable=False)       # login / register / reset_password
    expires_at = db.Column(db.DateTime, nullable=False)      # 过期时间
    attempts = db.Column(db.Integer, nullable=False, default=0)  # 尝试次数
    used = db.Column(db.Boolean, nullable=False, default=False)
    ip = db.Column(db.String(60), nullable=True)             # 请求IP
    created_at = db.Column(db.DateTime, default=datetime.now)

    def to_dict(self):
        return {
            'id': self.id,
            'channel': self.channel,
            'target': self.target,
            'purpose': self.purpose,
            'created_at': self.created_at.isoformat(),
            'expires_at': self.expires_at.isoformat()
        }

# ---------------------------
# 核心函数：DeepSeek API调用
# ---------------------------

def call_deepseek(prompt, system_prompt="????????AI??", max_retries=3):
    configs = resolve_runtime_ai_configs()
    if not configs:
        return "API Error: no AI runtime config available"

    last_error = None
    for config in configs:
        result, error = _call_ai_with_config(
            config,
            prompt,
            system_prompt,
            max_retries=max_retries,
            max_tokens=3500
        )
        if result is not None:
            return result
        last_error = error

    return last_error or "??????????"


# ---------------------------
# 前端页面路由
# ---------------------------

@app.route('/')
def landing_page():
    return app.send_static_file('index.html')

@app.route('/index.html')
@app.route('/product-studio')
@app.route('/product-studio.html')
def product_studio_page():
    """首页路由 - 返回前端主页面"""
    return app.send_static_file('index.html')

@app.route('/meetings')
@app.route('/meetings.html')
def meetings_page():
    return redirect('/meeting-room')

@app.route('/report')
@app.route('/report.html')
def report():
    """报告页面路由 - 返回报告查看页面"""
    return app.send_static_file('report.html')

@app.route('/meeting-room')
@app.route('/meeting-room.html')
def meeting_room_page():
    return app.send_static_file('meeting-room.html')

@app.route('/presentation')
@app.route('/presentation.html')
def presentation():
    return app.send_static_file('presentation.html')

@app.route('/register')
@app.route('/register.html')
def register_page():
    """注册页面路由 - 返回注册页面"""
    return app.send_static_file('register.html')

# ---------------------------
# 场景管理API
# ---------------------------

@app.route('/api/scenarios', methods=['POST'])
def create_scenario():
    """
    创建产品讨论场景
    
    POST /api/scenarios
    请求体：{ product_name, product_concept, core_selling_points[], discussion_topics[], occasion_type, occasion_description }
    返回：创建的场景对象
    
    功能：创建一个新的焦点小组讨论场景，包含产品信息和讨论配置
    """
    try:
        user_id = session.get('user_id')
        if not user_id:
            return app.response_class(
                response=json.dumps({"error": "请先登录"}, ensure_ascii=False),
                status=401,
                mimetype='application/json; charset=utf-8'
            )
        
        data = request.json
        validation_assumptions = normalize_text_list(data.get('validation_assumptions', []))
        scenario = ProductScenario(
            user_id=user_id,
            product_name=data['product_name'],
            product_concept=data['product_concept'],
            core_selling_points=json.dumps(data['core_selling_points'], ensure_ascii=False),
            discussion_topics=json.dumps(data['discussion_topics'], ensure_ascii=False),
            occasion_type=data.get('occasion_type', 'focus_group'),
            occasion_description=data.get('occasion_description', '标准焦点小组讨论'),
            research_goal=(data.get('research_goal') or '').strip(),
            decision_problem=(data.get('decision_problem') or '').strip(),
            target_user_profile=(data.get('target_user_profile') or '').strip(),
            competitor_context=(data.get('competitor_context') or '').strip(),
            validation_assumptions=json_dumps(validation_assumptions, []),
            research_plan_status='pending',
            meeting_status='standalone'
        )
        scenario.research_plan = build_research_plan(scenario)
        db.session.add(scenario)
        db.session.commit()
        result = scenario.to_dict()
        return app.response_class(
            response=json.dumps(result, ensure_ascii=False),
            status=201,
            mimetype='application/json; charset=utf-8'
        )
    except Exception as e:
        return app.response_class(
            response=json.dumps({'error': str(e)}, ensure_ascii=False),
            status=500,
            mimetype='application/json; charset=utf-8'
        )

@app.route('/api/scenarios/<int:id>/research-plan', methods=['GET'])
def get_research_plan(id):
    scenario, error_response = get_scenario_or_403(id)
    if error_response:
        return error_response

    if not scenario.research_plan:
        scenario.research_plan = build_research_plan(scenario)
        db.session.commit()

    return app.response_class(
        response=json.dumps({
            "scenario_id": scenario.id,
            "research_plan": scenario.research_plan,
            "research_plan_status": scenario.research_plan_status or 'pending'
        }, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/research-plan', methods=['POST'])
def refresh_research_plan(id):
    scenario, error_response = get_scenario_or_403(id)
    if error_response:
        return error_response

    scenario.research_plan = build_research_plan(scenario)
    scenario.research_plan_status = 'confirmed' if request.json and request.json.get('confirm') else 'pending'
    db.session.commit()

    return app.response_class(
        response=json.dumps({
            "scenario_id": scenario.id,
            "research_plan": scenario.research_plan,
            "research_plan_status": scenario.research_plan_status
        }, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/evidence', methods=['GET'])
def get_external_evidence(id):
    _, error_response = get_scenario_or_403(id)
    if error_response:
        return error_response

    items = ExternalEvidence.query.filter_by(scenario_id=id).order_by(ExternalEvidence.created_at.desc()).all()
    return app.response_class(
        response=json.dumps([item.to_dict() for item in items], ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/evidence', methods=['POST'])
def create_external_evidence(id):
    _, error_response = get_scenario_or_403(id)
    if error_response:
        return error_response

    data = request.json or {}
    title = (data.get('title') or '').strip()
    content = (data.get('content') or '').strip()
    if not title or not content:
        return app.response_class(
            response=json.dumps({"error": "证据标题和内容不能为空"}, ensure_ascii=False),
            status=400,
            mimetype='application/json; charset=utf-8'
        )

    item = ExternalEvidence(
        scenario_id=id,
        title=title,
        evidence_type=(data.get('evidence_type') or 'market_data').strip(),
        source_label=(data.get('source_label') or '').strip(),
        content=content,
        strength_level=(data.get('strength_level') or '高').strip() or '高'
    )
    item.strength_level = normalize_strength_level(item.strength_level)
    db.session.add(item)
    db.session.commit()

    return app.response_class(
        response=json.dumps(item.to_dict(), ensure_ascii=False),
        status=201,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/evidence/auto-generate', methods=['POST'])
def auto_generate_external_evidence(id):
    scenario, error_response = get_scenario_or_403(id)
    if error_response:
        return error_response

    ensure_extended_schema()

    existing_items = ExternalEvidence.query.filter_by(scenario_id=id).order_by(ExternalEvidence.created_at.desc()).all()
    if existing_items:
        return app.response_class(
            response=json.dumps({
                "items": [item.to_dict() for item in existing_items],
                "message": "当前场景已有外部证据，可继续手动补充。"
            }, ensure_ascii=False),
            status=200,
            mimetype='application/json; charset=utf-8'
        )

    result = call_deepseek(
        build_ai_evidence_prompt(scenario),
        "你是企业研究助理，负责生成可供内部讨论参考的研究线索。输出必须是合法 JSON。"
    )
    if result.startswith(("API Error", "Request Error", "请求失败")):
        return app.response_class(
            response=json.dumps({"error": "AI 参考证据生成失败，请稍后重试"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    cleaned = result.strip()
    cleaned = re.sub(r'^```json\s*', '', cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r'^```\s*', '', cleaned)
    cleaned = re.sub(r'\s*```$', '', cleaned)

    try:
        items = json.loads(cleaned)
    except Exception:
        return app.response_class(
            response=json.dumps({"error": "AI 返回的证据格式无法解析，请重试"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    if not isinstance(items, list) or not items:
        return app.response_class(
            response=json.dumps({"error": "AI 未生成可用证据，请重试"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    created_items = []
    for raw in items[:3]:
        if not isinstance(raw, dict):
            continue

        title = str(raw.get('title') or '').strip()
        content = str(raw.get('content') or '').strip()
        if not title or not content:
            continue

        item = ExternalEvidence(
            scenario_id=id,
            title=title[:200],
            evidence_type=(str(raw.get('evidence_type') or 'market_data').strip() or 'market_data')[:50],
            source_label=(str(raw.get('source_label') or 'AI参考生成').strip() or 'AI参考生成')[:200],
            content=content,
            strength_level=normalize_strength_level(raw.get('strength_level') or '中')
        )
        db.session.add(item)
        created_items.append(item)

    if not created_items:
        return app.response_class(
            response=json.dumps({"error": "AI 未生成有效证据内容，请重试"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    db.session.commit()

    return app.response_class(
        response=json.dumps({
            "items": [item.to_dict() for item in created_items],
            "message": "AI 已生成参考证据，用户补充证据可选。"
        }, ensure_ascii=False),
        status=201,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/evidence/auto-generate-v2', methods=['POST'])
def auto_generate_external_evidence_v2(id):
    scenario, error_response = get_scenario_or_403(id)
    if error_response:
        return error_response

    ensure_extended_schema()

    queries = build_search_queries(scenario)
    print(f"[evidence-v2] scenario={id} queries={queries}")
    if not queries:
        return app.response_class(
            response=json.dumps({"error": "当前场景信息不足，无法生成检索关键词"}, ensure_ascii=False),
            status=400,
            mimetype='application/json; charset=utf-8'
        )

    try:
        search_hits = []
        seen_urls = set()
        for query in queries:
            for item in search_duckduckgo(query, limit=4):
                if item["url"] in seen_urls:
                    continue
                seen_urls.add(item["url"])
                search_hits.append(item)
                if len(search_hits) >= 6:
                    break
            if len(search_hits) >= 6:
                break
    except Exception as exc:
        print(f"[evidence-v2] scenario={id} stage=search_failed error={exc}")
        return app.response_class(
            response=json.dumps({"error": "外部检索失败，请稍后重试"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    if not search_hits:
        return app.response_class(
            response=json.dumps({"error": "未检索到足够外部资料，请调整场景信息后重试"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    source_materials = []
    for hit in search_hits[:5]:
        try:
            excerpt = fetch_page_excerpt(hit["url"])
        except Exception as exc:
            print(f"[evidence-v2] scenario={id} stage=fetch_failed url={hit['url']} error={exc}")
            continue
        if excerpt:
            source_materials.append({
                "title": hit["title"],
                "url": hit["url"],
                "domain": hit["domain"],
                "excerpt": excerpt
            })

    print(f"[evidence-v2] scenario={id} fetched_sources={len(source_materials)}")
    if len(source_materials) < 2:
        return app.response_class(
            response=json.dumps({"error": "检索结果不足，暂时无法整理出可靠参考证据"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    synthesis = call_deepseek(
        build_evidence_synthesis_prompt(scenario, source_materials),
        "你是企业研究助理，负责基于外部资料整理可追溯证据。输出必须是合法 JSON。"
    )
    if synthesis.startswith(("API Error", "Request Error", "请求失败")):
        return app.response_class(
            response=json.dumps({"error": "AI 整理检索结果失败，请稍后重试"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    try:
        items = json.loads(extract_json_block(synthesis))
    except Exception as exc:
        print(f"[evidence-v2] scenario={id} stage=parse_failed error={exc}")
        return app.response_class(
            response=json.dumps({"error": "AI 返回的证据格式无法解析，请重试"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    if not isinstance(items, list) or not items:
        return app.response_class(
            response=json.dumps({"error": "AI 未整理出可用证据，请重试"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    try:
        ExternalEvidence.query.filter_by(scenario_id=id, generated_by='ai_search').delete()
    except Exception:
        pass

    created_items = []
    for raw in items[:3]:
        if not isinstance(raw, dict):
            continue

        title = str(raw.get('title') or '').strip()
        content = str(raw.get('content') or '').strip()
        source_url = str(raw.get('source_url') or '').strip()
        if not title or not content or not source_url:
            continue

        item = ExternalEvidence(
            scenario_id=id,
            title=title[:200],
            evidence_type=(str(raw.get('evidence_type') or 'market_data').strip() or 'market_data')[:50],
            source_label=(str(raw.get('source_label') or 'AI联网检索').strip() or 'AI联网检索')[:200],
            source_url=source_url,
            source_title=(str(raw.get('source_title') or title).strip() or title)[:300],
            source_domain=(str(raw.get('source_domain') or '').strip())[:200],
            generated_by='ai_search',
            content=content,
            strength_level='中'
        )
        item.strength_level = normalize_strength_level(item.strength_level)
        db.session.add(item)
        created_items.append(item)

    if not created_items:
        db.session.rollback()
        return app.response_class(
            response=json.dumps({"error": "AI 未整理出带来源链接的有效证据"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    try:
        db.session.commit()
    except Exception as exc:
        db.session.rollback()
        print(f"[evidence-v2] scenario={id} stage=save_failed error={exc}")
        return app.response_class(
            response=json.dumps({"error": "证据保存失败，请稍后重试"}, ensure_ascii=False),
            status=500,
            mimetype='application/json; charset=utf-8'
        )

    return app.response_class(
        response=json.dumps({
            "items": [item.to_dict() for item in created_items],
            "message": "AI 已联网获取参考证据，旧的 AI 证据已覆盖，手动证据已保留。"
        }, ensure_ascii=False),
        status=201,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios', methods=['GET'])
def get_scenarios():
    """
    获取当前用户的所有场景列表
    
    GET /api/scenarios
    返回：场景列表数组
    
    功能：查询当前登录用户创建的所有产品讨论场景
    """
    user_id = session.get('user_id')
    if not user_id:
        return app.response_class(
            response=json.dumps({"error": "请先登录"}, ensure_ascii=False),
            status=401,
            mimetype='application/json; charset=utf-8'
        )
    
    scenarios = (
        ProductScenario.query
        .filter_by(user_id=user_id)
        .filter_by(meeting_status='standalone')
        .order_by(ProductScenario.created_at.desc(), ProductScenario.id.desc())
        .all()
    )
    result = [s.to_dict() for s in scenarios]
    return app.response_class(
        response=json.dumps(result, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

def get_scenario_or_403(id, allow_meeting_member=False):
    """
    场景权限检查辅助函数
    
    Args:
        id: 场景ID
    
    Returns:
        (scenario对象, None) - 成功
        (None, error_response) - 失败（未登录/场景不存在/无权访问）
    
    功能：验证场景存在性和用户访问权限，用于所有场景相关API的前置校验
    """
    user_id = session.get('user_id')
    if not user_id:
        return None, app.response_class(
            response=json.dumps({"error": "请先登录"}, ensure_ascii=False),
            status=401,
            mimetype='application/json; charset=utf-8'
        )
    
    scenario = db.session.get(ProductScenario, id)
    if not scenario:
        return None, app.response_class(
            response=json.dumps({"error": "场景不存在"}, ensure_ascii=False),
            status=404,
            mimetype='application/json; charset=utf-8'
        )
    
    sync_scenario_meeting_status(scenario)

    if scenario.user_id != user_id:
        if allow_meeting_member:
            room = MeetingRoom.query.filter_by(scenario_id=scenario.id).first()
            if room and room.owner_user_id == user_id:
                return scenario, None
            if room:
                membership = MeetingRoomMember.query.filter_by(room_id=room.id, user_id=user_id).first()
                if membership:
                    return scenario, None
        return None, app.response_class(
            response=json.dumps({"error": "无权访问该场景"}, ensure_ascii=False),
            status=403,
            mimetype='application/json; charset=utf-8'
        )
    
    return scenario, None

@app.route('/api/scenarios/<int:id>', methods=['GET'])
def get_scenario(id):
    """
    获取单个场景详情
    
    GET /api/scenarios/{id}
    返回：场景对象详情
    
    功能：查询指定场景的详细信息
    """
    scenario, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response
    
    result = scenario.to_dict()
    return app.response_class(
        response=json.dumps(result, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

# ---------------------------
# AI参与者生成API
# ---------------------------

@app.route('/api/scenarios/<int:id>/generate-participants', methods=['POST'])
def generate_participants_with_ai(id):
    """
    使用AI生成虚拟参与者角色
    
    POST /api/scenarios/{id}/generate-participants
    请求体：{ count, target_audience, regenerate, profile }
    返回：参与者列表
    
    功能说明：
        - 根据场景类型（焦点小组/产品团队/销售对话/深度访谈/头脑风暴）生成相应角色
        - 每个角色包含：姓名、标签、性格、背景
        - 默认生成4个参与者，支持自定义数量
        - 如果已有参与者且未设置regenerate，则直接返回现有参与者
        - 支持通过profile参数自定义参与者画像特征
    """
    scenario, error_response = get_scenario_or_403(id)
    if error_response:
        return error_response
    standalone_error = reject_if_not_standalone(scenario, "会议室内部场景不能使用普通人设生成，请前往会议室系统继续。")
    if standalone_error:
        return standalone_error

    data = request.get_json() or {}
    try:
        count = max(1, int(data.get('count', 4)))
    except (TypeError, ValueError):
        count = 4

    target_audience = str(data.get('target_audience') or '').strip()
    regenerate = bool(data.get('regenerate', False))
    preserve_custom = bool(data.get('preserve_custom', True))
    persona_instructions = str(data.get('persona_instructions') or '').strip()
    legacy_profile = data.get('profile', {}) if isinstance(data.get('profile', {}), dict) else {}
    if target_audience and not persona_instructions:
        persona_instructions = f"参与者应优先贴近这类目标人群：{target_audience}"

    raw_seed_participants = data.get('seed_participants', [])
    seed_participants = []
    if isinstance(raw_seed_participants, list):
        for index, seed in enumerate(raw_seed_participants[:count], start=1):
            seed_participants.append(sanitize_participant_seed(seed, default_name=f"种子参与者{index}"))

    existing_participants = get_standalone_participants(id)
    if existing_participants and not regenerate:
        return app.response_class(
            response=json.dumps([p.to_dict() for p in existing_participants], ensure_ascii=False),
            status=200,
            mimetype='application/json; charset=utf-8'
        )

    existing_custom_participants = []
    if preserve_custom:
        existing_custom_participants = [
            participant for participant in get_standalone_participants(id)
            if participant.is_custom
        ]

    if regenerate:
        query = VirtualParticipant.query.filter(VirtualParticipant.scenario_id == id)
        query = query.filter(or_(VirtualParticipant.room_managed.is_(False), VirtualParticipant.room_managed.is_(None)))
        if preserve_custom:
            query = query.filter_by(is_ai_generated=True)
        query.delete()
        db.session.commit()

    generation_count = count
    prompt_seed_participants = seed_participants
    merge_seed_participants = seed_participants
    if preserve_custom and existing_custom_participants:
        generation_count = max(count - len(existing_custom_participants), 0)
        merge_seed_participants = []

    if generation_count == 0:
        current_participants = get_standalone_participants(id)
        return app.response_class(
            response=json.dumps({
                "status": "success",
                "ai_response": "",
                "participants": [p.to_dict() for p in current_participants]
            }, ensure_ascii=False),
            status=200,
            mimetype='application/json; charset=utf-8'
        )

    role_mapping = {
        "product_team": "企业产品团队成员",
        "sales_conversation": "销售相关角色",
        "focus_group": "目标用户",
        "user_interview": "深度访谈用户",
        "brainstorming": "创意团队成员"
    }
    role_type = role_mapping.get(scenario.occasion_type, "目标用户")
    prompt, system_prompt = build_generated_participant_prompt(
        scenario,
        role_type,
        generation_count,
        persona_instructions,
        prompt_seed_participants,
        legacy_profile
    )

    result = call_deepseek(prompt, system_prompt)
    try:
        personas = json.loads(result)
        if not isinstance(personas, list):
            personas = [personas]
    except Exception:
        personas = build_fallback_personas(merge_seed_participants, generation_count)

    merged_personas = merge_generated_personas(personas, merge_seed_participants, generation_count)
    generation_metadata = {
        "persona_instructions": persona_instructions,
        "seed_participants": seed_participants,
        "preserve_custom": preserve_custom,
        "target_audience": target_audience,
        "legacy_profile": legacy_profile
    }

    created_participants = []
    for person in merged_personas:
        participant = VirtualParticipant(
            scenario_id=id,
            persona_name=person.get('persona_name') or '未知参与者',
            persona_tags=json.dumps(person.get('persona_tags', []) or ['普通用户'], ensure_ascii=False),
            personality=person.get('personality') or '性格待补充',
            background=person.get('background') or '背景待补充',
            usage_goal=person.get('usage_goal', ''),
            budget_sensitivity=person.get('budget_sensitivity', ''),
            brand_preference=person.get('brand_preference', ''),
            risk_aversion=person.get('risk_aversion', ''),
            decision_style=person.get('decision_style', ''),
            deal_breakers=json_dumps(person.get('deal_breakers', []), []),
            stance_summary=person.get('stance_summary', ''),
            stance_state=build_participant_stance_state(person.get('stance_summary', '')),
            is_custom=False,
            is_ai_generated=True,
            custom_params=json.dumps(generation_metadata, ensure_ascii=False)
        )
        db.session.add(participant)
        created_participants.append(participant)

    db.session.commit()

    return app.response_class(
        response=json.dumps({
            "status": "success",
            "ai_response": result[:500] if isinstance(result, str) else '',
            "participants": [p.to_dict() for p in created_participants]
        }, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )
    
    data = request.json or {}
    count = data.get('count', 4)  # 默认生成4个参与者
    target_audience = data.get('target_audience', '普通消费者')
    regenerate = data.get('regenerate', False)  # 是否重新生成
    profile = data.get('profile', {})  # 自定义画像参数

    # 如果已有参与者且不需要重新生成，直接返回
    existing_participants = VirtualParticipant.query.filter_by(scenario_id=id).all()
    if existing_participants and not regenerate:
        result = [p.to_dict() for p in existing_participants]
        return app.response_class(
            response=json.dumps(result, ensure_ascii=False),
            status=200,
            mimetype='application/json; charset=utf-8'
        )

    # 删除现有参与者
    VirtualParticipant.query.filter_by(scenario_id=id).delete()
    db.session.commit()

    # 根据场景类型设置角色类型
    occasion_type = scenario.occasion_type
    occasion_desc = scenario.occasion_description or ""

    role_mapping = {
        "product_team": ("企业产品团队成员", "产品经理、资深研发工程师、用户体验设计师、市场分析师、运营负责人", "你需要扮演企业内部产品团队成员，从专业角度分析产品"),
        "sales_conversation": ("销售人员", "销售主管、大客户经理、渠道专员、区域销售代表", "你需要扮演销售人员，从销售角度探测客户需求"),
        "focus_group": ("目标用户", "硬核极客、价格敏感型宝妈、颜控女大学生、职场白领、运动爱好者", "你需要扮演目标用户，表达真实使用体验和需求"),
        "user_interview": ("深度访谈用户", "资深用户、潜在客户、竞品用户、行业专家", "你需要扮演被访谈用户，深入表达观点"),
        "brainstorming": ("创意团队成员", "产品经理、创意总监、市场营销专家、技术骨干、用户研究员", "你需要扮演创意团队成员，激发创新想法")
    }

    role_type, role_examples, role_desc = role_mapping.get(occasion_type, ("普通消费者", "普通用户、潜在客户", "你需要扮演普通消费者"))

    # 构建自定义画像约束
    profile_constraints = []
    if profile.get('professional_level'):
        profile_constraints.append(f"• 专业级别：{profile['professional_level']}")
    if profile.get('income_level'):
        profile_constraints.append(f"• 收入水平：{profile['income_level']}")
    if profile.get('industry'):
        profile_constraints.append(f"• 所在行业：{profile['industry']}")
    if profile.get('age_range'):
        profile_constraints.append(f"• 年龄段：{profile['age_range']}")
    if profile.get('education'):
        profile_constraints.append(f"• 教育背景：{profile['education']}")
    if profile.get('city_tier'):
        profile_constraints.append(f"• 所在城市：{profile['city_tier']}")
    if profile.get('needs') and isinstance(profile['needs'], list):
        profile_constraints.append(f"• 核心需求：{', '.join(profile['needs'])}")
    if profile.get('budget_sensitivity'):
        profile_constraints.append(f"• 预算敏感度：{profile['budget_sensitivity']}")
    if profile.get('risk_aversion'):
        profile_constraints.append(f"• 风险偏好：{profile['risk_aversion']}")
    if profile.get('decision_style'):
        profile_constraints.append(f"• 决策风格：{profile['decision_style']}")
    if profile.get('brand_preference'):
        profile_constraints.append(f"• 品牌偏好：{profile['brand_preference']}")
    if profile.get('deal_breakers') and isinstance(profile['deal_breakers'], list):
        profile_constraints.append(f"• 反感点：{', '.join(profile['deal_breakers'])}")

    profile_text = "\n".join(profile_constraints) if profile_constraints else "无特殊约束"
    diversity_mix = profile.get('diversity_mix', False)

    # 构建AI提示词
    diversity_note = "且角色之间要体现多样性（年龄、背景、观点差异）" if diversity_mix else ""
    prompt = f"""请为【{occasion_desc}】场景生成{count}个专业的{role_type}角色画像{diversity_note}，用于产品分析讨论。

用户画像约束：
{profile_text}

要求：
1. 每个角色要有独特的名字、背景、性格特点
2. 背景要具体，包括职业经历、专业领域
3. 性格特点要与背景相符
4. 角色之间要有差异化，适合进行产品讨论
5. 必须符合上述用户画像约束
6. 还要补充以下字段：usage_goal, budget_sensitivity, brand_preference, risk_aversion, decision_style, deal_breakers, stance_summary
7. deal_breakers输出数组，其他字段输出字符串

请直接输出JSON数组格式，例如：
[{{"name":"张三","tags":["技术","创新"],"personality":"理性严谨","background":"背景描述","usage_goal":"解决什么问题","budget_sensitivity":"高","brand_preference":"偏好成熟品牌","risk_aversion":"中高","decision_style":"先比较再决定","deal_breakers":["价格过高","难以上手"],"stance_summary":"重视稳定性和性价比"}}]

直接输出JSON数组，不要加任何其他内容。"""

    system_prompt = f"你是一个专业的企业管理顾问，擅长创建真实、专业的职场角色画像。现在需要为【{occasion_desc}】场景生成{count}个{role_type}角色{diversity_note}，每个角色必须有明确的职位和专业背景，且符合给定的用户画像约束。"

    # 调用AI生成角色
    result = call_deepseek(prompt, system_prompt)

    # 解析AI响应
    try:
        personas = json.loads(result)
        if not isinstance(personas, list):
            personas = [personas]
    except:
        # 解析失败时使用默认角色
        personas = [
            {"name": "随机用户A", "tags": ["普通用户"], "personality": "普通用户性格", "background": "普通背景"}
        ]

    # 保存参与者到数据库
    participants = []
    profile_json = json.dumps(profile, ensure_ascii=False) if profile else None
    for p in personas[:count]:
        participant = VirtualParticipant(
            scenario_id=id,
            persona_name=p.get('name', '未知用户'),
            persona_tags=json.dumps(p.get('tags', ['普通用户']), ensure_ascii=False),
            personality=p.get('personality', '性格未知'),
            background=p.get('background', '背景未知'),
            usage_goal=p.get('usage_goal', ''),
            budget_sensitivity=p.get('budget_sensitivity', ''),
            brand_preference=p.get('brand_preference', ''),
            risk_aversion=p.get('risk_aversion', ''),
            decision_style=p.get('decision_style', ''),
            deal_breakers=json_dumps(p.get('deal_breakers', []), []),
            stance_summary=p.get('stance_summary', ''),
            stance_state=json_dumps({
                "current_position": p.get('stance_summary', ''),
                "confidence": "中",
                "last_updated_round": 0
            }, {}),
            is_custom=False,
            is_ai_generated=True,
            custom_params=profile_json
        )
        db.session.add(participant)
        participants.append(participant)

    db.session.commit()

    result_data = {
        "status": "success",
        "ai_response": result[:500],
        "participants": [p.to_dict() for p in participants]
    }
    return app.response_class(
        response=json.dumps(result_data, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

# ---------------------------
# 对话模拟API
# ---------------------------

def build_local_fallback_report_bundle(report_inputs, progress_callback=None, reason='no_runtime_ai_config', upstream_error=None, emit_transition=False, emit_prepare_stage=True):
    scenario = report_inputs['scenario']
    scenario_snapshot = report_inputs.get('scenario_snapshot') or {}
    scenario_id = report_inputs.get('scenario_id') or report_inputs.get('scenario_db_id') or scenario_snapshot.get('id')
    log_report_event(
        'report_generation_fallback_started',
        scenario_id=scenario_id,
        mode='local_fallback',
        reason=reason,
        upstream_error=upstream_error
    )
    if emit_transition and progress_callback:
        progress_callback("fallback_started", "AI 结果未通过校验，切换为稳定报告生成", 82)

    topics = normalize_text_list(scenario_snapshot.get('discussion_topics'))
    assumptions = normalize_text_list(scenario_snapshot.get('validation_assumptions'))
    participants = report_inputs.get('participants') or []
    records = report_inputs.get('records') or []
    external_evidence = report_inputs.get('external_evidence') or []
    non_host_records = [
        record for record in records
        if not getattr(record, 'is_host', False) and str(getattr(record, 'content', '') or '').strip()
    ]
    product_name = scenario_snapshot.get('product_name') or '当前产品'
    research_plan = scenario_snapshot.get('research_plan') or ''
    report_title = f"{product_name} - 产品分析报告 ({report_inputs['current_date']})"

    if progress_callback and emit_prepare_stage:
        progress_callback("prepare_inputs", "准备讨论材料", 8)

    def unique_texts(items):
        result = []
        seen = set()
        for item in items:
            text = str(item or '').strip()
            if not text or text in seen:
                continue
            seen.add(text)
            result.append(text)
        return result

    joined_records = "\n".join([record.content for record in non_host_records])
    topic_candidates = unique_texts(topics + assumptions)
    keyword_topic_map = [
        ("噪音", "噪音与训练体验"),
        ("维护", "维护复杂度与稳定性"),
        ("场景", "场景适配与稳定性"),
        ("室内", "场地环境适配"),
        ("室外", "场地环境适配"),
        ("价格", "价格接受度与投入回报"),
        ("付费", "价格接受度与投入回报"),
        ("效率", "训练效率提升价值"),
    ]
    for keyword, label in keyword_topic_map:
        if keyword in joined_records:
            topic_candidates.append(label)
    topic_candidates.extend([
        "使用体验与维护顾虑",
        "场景适配与稳定性",
        "价格接受度与投入回报",
        "后续验证重点",
    ])
    topic_candidates = unique_texts(topic_candidates)[:4]

    risk_candidates = list(assumptions)
    if "噪音" in joined_records:
        risk_candidates.append("噪音影响训练体验")
    if "维护" in joined_records:
        risk_candidates.append("维护复杂度和持续运维成本")
    if "价格" in joined_records or "付费" in joined_records:
        risk_candidates.append("价格接受度与价值感知不匹配")
    if "场景" in joined_records or "室外" in joined_records or "室内" in joined_records:
        risk_candidates.append("不同场地环境下的稳定性仍需验证")
    risk_candidates.extend(["维护复杂度和持续运维成本", "不同场地环境下的稳定性仍需验证"])
    risk_candidates = unique_texts(risk_candidates)[:2]

    action_candidates = unique_texts(
        topics + assumptions + [
            "安排小范围试点",
            "补充一线用户访谈",
            "验证噪音与维护表现",
            "确认目标用户的价格接受区间",
        ]
    )[:2]

    primary_risk = risk_candidates[0] if risk_candidates else "维护复杂度和持续运维成本"
    secondary_risk = risk_candidates[1] if len(risk_candidates) > 1 else "不同场地环境下的稳定性仍需验证"
    primary_action = action_candidates[0] if action_candidates else "安排小范围试点"
    secondary_action = action_candidates[1] if len(action_candidates) > 1 else "补充一线用户访谈"
    evidence_seed = external_evidence[0].content[:80] if external_evidence else (non_host_records[0].content[:80] if non_host_records else "现有证据主要来自有限讨论记录")
    participant_names = "、".join([item.persona_name for item in participants[:3]]) or "当前参与者"

    topic_lines = []
    for index, topic in enumerate(topic_candidates, start=1):
        matching_records = [record.content.strip() for record in non_host_records if topic and topic in (record.content or '')]
        fallback_index = min(index - 1, max(len(non_host_records) - 1, 0))
        sample_point = matching_records[0] if matching_records else (non_host_records[fallback_index].content.strip() if non_host_records else "当前讨论记录较少，建议补充更多一线反馈。")
        topic_lines.append(
            "\n".join(
                [
                    f"**议题{index}：{topic}**",
                    f"- 讨论中的观点：{sample_point[:120]}",
                    f"- 系统判断：该议题已经在现有讨论中出现，后续需要继续量化验证。",
                    f"- 待验证假设：围绕“{topic}”的判断仍需要更多真实场景数据支持。",
                ]
            )
        )

    summary_section = "\n".join(
        [
            "**1）核心结论**",
            f"围绕“{product_name}”的现有讨论显示，这个方案已经具备继续验证的业务价值，不再停留在抽象概念层面。",
            "参与者的反馈并非一边倒乐观，但已经足以说明它解决的问题是真实存在的，下一步重点应该从“是否有价值”转向“价值能否稳定兑现”。",
            "",
            "**2）目标用户**",
            f"当前更适合优先聚焦的不是所有潜在用户，而是对效率提升最敏感、愿意容忍初期打磨成本的一线使用者，例如 {participant_names} 所代表的高频场景用户。",
            "这些用户既能快速反馈真实问题，也最容易帮助团队判断产品价值是否足以支撑后续投入。",
            "",
            "**3）推进判断**",
            "建议继续推进，但推进方式应当是小范围、可控、可复盘的验证，而不是直接扩大投入或过早承诺规模化。",
            "换句话说，当前阶段更像是进入“验证设计质量”的窗口，而不是进入“全面铺开”的窗口。",
            "",
            "**4）关键风险**",
            f"当前最需要防范的核心风险仍然是“{primary_risk}”，因为它会直接影响用户对产品稳定性和实际价值的判断。",
            f"如果这个风险在试点里反复出现，即便用户认可概念本身，也会明显削弱继续推进 {product_name} 的把握。",
            "",
            "**5）下一步动作**",
            f"下一步应优先围绕“{primary_action}”安排补充验证，并尽量把验证对象、使用场景和反馈口径先固定下来。",
            "同时要同步整理更稳定的证据链，让后续结论不只是基于单点印象，而是能够支撑管理层继续决策的连续观察结果。",
        ]
    )
    divergence_table = "\n".join(
        [
            "| 维度 | 观点A | 观点B | 当前偏向 | 验证重点 |",
            "| --- | --- | --- | --- | --- |",
            f"| 推进节奏 | 先快速落地 | 先补充验证 | 先补充验证 | {primary_risk}能否被有效控制 |",
            "",
            "**1）分歧焦点**",
            f"当前分歧并不在于产品有没有潜力，而在于是否应该尽快推进，还是先把 {secondary_risk} 这类关键不确定性压低后再进入更大范围测试。",
            "一部分参与者更关注抢先验证和动作速度，另一部分参与者更关注早期体验失真会不会让后续判断被噪音带偏。",
            "",
            "**2）当前偏向**",
            "从现有证据看，系统更偏向先补充验证再扩大投入，而不是直接把试点当成规模化前奏。",
            "这种偏向的原因不是保守，而是因为当前证据仍不足以证明产品在复杂场景下可以稳定交付同样的价值体验。",
            "",
            "**3）验证重点**",
            f"下一轮最值得盯住的验证重点仍然是：{evidence_seed}。",
            f"如果与 {primary_risk} 相关的问题能在试点中被清晰量化并且可控，当前这组分歧会明显收敛。",
        ]
    )
    recommendation_section = "\n".join(
        [
            "**1）是否继续推进**",
            "建议继续推进，但仅限于下一轮小范围验证，不建议直接进入更大规模的资源投入。",
            f"当前更合理的动作是把 {product_name} 放进可控试点中继续观察，而不是基于一轮讨论就提前下结论。",
            "",
            "**2）核心原因**",
            f"现有讨论已经表明 {product_name} 具备明确价值点，尤其在效率提升和连续使用体验上已经出现了真实需求信号。",
            "但这些信号仍然更像“值得继续验证”的证据，而不是“已经足够支持全面扩张”的证据，因此结论需要保持积极但克制。",
            "",
            "**3）建议动作**",
            f"第一，围绕“{primary_action}”安排补充访谈或试点观察，优先选择最接近真实付费和高频使用的用户群体。",
            f"第二，用下一轮试点继续验证“{secondary_action}”的真实表现，并同步记录失败场景、复购意愿和替代方案比较结果。",
            "",
            "**4）关键风险**",
            f"当前最值得持续追踪的风险包括：{primary_risk}，以及 {secondary_risk}。",
            "如果这两类问题在真实场景中频繁出现，那么用户即使认可概念，也很可能不会形成持续使用或采购决策。",
            "",
            "**5）停止条件**",
            "如果下一轮验证仍然无法证明用户接受度、使用价值和持续意愿，团队就不应该继续追加更大投入。",
            "更具体地说，一旦试点中的负面反馈主要集中在稳定性、体验中断或价值感不足，并且短期内看不到改善路径，就应当暂停推进。",
        ]
    )

    section_map = {
        'section_summary': summary_section,
        'section_analysis': "\n\n".join(topic_lines[:4]),
        'section_divergence': divergence_table,
        'section_recommendation': recommendation_section,
    }

    for section_spec in REPORT_SECTION_SPECS:
        if progress_callback:
            progress_callback(section_spec['code'], section_spec['phase'], section_spec['pct'])

    report_content = sanitize_report_text(assemble_report_content(section_map))
    metadata = enrich_report_metadata(
        extract_report_metadata(report_content, external_evidence),
        external_evidence
    )
    executive_summary = (
        f"基于 {participant_names} 的讨论，系统建议继续推进 {product_name}，"
        f"但应先完成小范围验证，重点确认{primary_risk} 与 {primary_action}。"
    )

    if progress_callback:
        progress_callback("finalize_report", "整理摘要并保存报告", 92)

    report = persist_analysis_report(
        scenario_id,
        report_title,
        report_content,
        executive_summary,
        metadata,
        research_plan=research_plan,
        discussion_summary=report_inputs.get('discussion_summary')
    )
    log_report_event(
        'report_generation_fallback_completed',
        scenario_id=scenario_id,
        mode='local_fallback',
        report_id=getattr(report, 'id', None)
    )
    return report, None


def generate_direct_report_bundle(report_inputs, progress_callback=None, fallback_reason='ai_generation_failed', emit_transition_on_fallback=True, emit_prepare_stage=True):
    log_report_event(
        'report_generation_ai_started',
        scenario_id=report_inputs.get('scenario_id') or report_inputs.get('scenario_db_id'),
        mode='ai'
    )
    state = build_report_generation_context(report_inputs)
    if not emit_prepare_stage:
        state['_prepare_emitted'] = True
    while len(state['completed_sections']) < len(REPORT_SECTION_SPECS):
        current_index = len(state['completed_sections'])
        current_section = REPORT_SECTION_SPECS[current_index] if current_index < len(REPORT_SECTION_SPECS) else None
        _, section_error, _ = advance_report_generation(state, progress_callback=progress_callback)
        if section_error:
            log_report_event(
                'report_generation_section_failed',
                scenario_id=report_inputs.get('scenario_id') or report_inputs.get('scenario_db_id'),
                mode='ai',
                section_code=current_section.get('code') if current_section else None,
                error=section_error,
                level='ERROR'
            )
            return build_local_fallback_report_bundle(
                report_inputs,
                progress_callback=progress_callback,
                reason=fallback_reason,
                upstream_error=section_error,
                emit_transition=emit_transition_on_fallback,
                emit_prepare_stage=emit_prepare_stage
            )

    report, report_error = build_report_from_generation_state(
        report_inputs,
        state,
        progress_callback=progress_callback
    )
    if report_error:
        log_report_event(
            'report_generation_finalize_failed',
            scenario_id=report_inputs.get('scenario_id') or report_inputs.get('scenario_db_id'),
            mode='ai',
            error=report_error,
            level='ERROR'
        )
        return build_local_fallback_report_bundle(
            report_inputs,
            progress_callback=progress_callback,
            reason='report_finalize_failed',
            upstream_error=report_error,
            emit_transition=emit_transition_on_fallback,
            emit_prepare_stage=emit_prepare_stage
        )
    return report, None


def generate_analysis_report_bundle(report_inputs, progress_callback=None):
    emit_progress = build_report_progress_dispatcher(progress_callback)
    if not resolve_runtime_ai_configs():
        return build_local_fallback_report_bundle(
            report_inputs,
            progress_callback=emit_progress,
            reason='no_runtime_ai_config'
        )

    scenario_id = report_inputs.get('scenario_id') or report_inputs.get('scenario_db_id')
    emit_progress("discussion_summary", "生成讨论摘要", 14)
    log_report_event('discussion_summary_started', scenario_id=scenario_id, mode='ai')
    discussion_summary, summary_error = generate_discussion_summary(report_inputs)
    if summary_error:
        log_report_event(
            'discussion_summary_failed',
            scenario_id=scenario_id,
            mode='ai',
            error=summary_error,
            level='WARNING'
        )
        emit_progress("fallback_to_direct_report", "摘要生成失败，切换为原始对话报告生成", 18)
        log_report_event('fallback_to_direct_report', scenario_id=scenario_id, mode='ai', reason='discussion_summary_failed')
        return generate_direct_report_bundle(
            report_inputs,
            progress_callback=emit_progress,
            fallback_reason='ai_generation_failed',
            emit_transition_on_fallback=True,
            emit_prepare_stage=False
        )

    report_inputs = dict(report_inputs)
    report_inputs['discussion_summary'] = discussion_summary
    report_inputs['discussion_summary_text'] = format_discussion_summary_for_report(discussion_summary)

    report, report_error = generate_direct_report_bundle(
        report_inputs,
        progress_callback=emit_progress,
        fallback_reason='summary_guided_report_failed',
        emit_transition_on_fallback=True,
        emit_prepare_stage=False
    )
    if report_error:
        return report, report_error
    return report, None

@app.route('/api/scenarios/<int:id>/simulate', methods=['POST'])
def simulate_conversation_with_ai(id):
    """
    模拟焦点小组对话（同步模式）
    
    POST /api/scenarios/{id}/simulate
    请求体：{ rounds, message_count }
    返回：模拟结果统计
    
    功能说明：
        - 根据场景配置和参与者角色，模拟多轮焦点小组讨论
        - 支持两种模式：rounds模式（按轮数）和message_count模式（按消息数量）
        - AI主持人会引导讨论、总结话题、推进对话
        - 每轮讨论覆盖不同的讨论话题
        - 支持AI调用失败时的备用回复机制
    """
    scenario, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response
    
    data = request.json or {}
    rounds = data.get('rounds', 3)  # 默认3轮讨论
    message_count = data.get('message_count', 0)  # 按消息数量模式（0表示使用rounds模式）

    participants = get_standalone_participants(id)

    if not participants:
        return app.response_class(
            response=json.dumps({"error": "请先生成参与者"}, ensure_ascii=False),
            status=400,
            mimetype='application/json; charset=utf-8'
        )

    selling_points = json.loads(scenario.core_selling_points)
    topics = json.loads(scenario.discussion_topics)

    # 清空之前的对话记录
    ConversationRecord.query.filter_by(scenario_id=id).delete()

    # 添加主持人开场介绍
    host_intro = ConversationRecord(
        scenario_id=id,
        participant_id=0,
        content=f"好啦，各位朋友大家好！欢迎来参加今天的焦点小组讨论。\n\n今天我们主要想聊聊「{scenario.product_name}」这个产品。\n\n先简单说一下这个产品：{scenario.product_concept}\n\n它主打的卖点是：{','.join(selling_points)}\n\n我们就从第一个话题开始吧：{topics[0]}\n\n大家有什么想法就说什么，不用客气啊！",
        is_host=True,
        message_type='intro'
    )
    db.session.add(host_intro)
    db.session.commit()

    topic_msg = ConversationRecord(
        scenario_id=id,
        participant_id=0,
        content=f"话题 1：{topics[0]}",
        is_host=True,
        message_type='topic'
    )
    db.session.add(topic_msg)
    db.session.commit()

    conversation_history = []
    research_goal = scenario.research_goal or "识别该产品是否值得推进，以及最需要验证的风险和机会。"
    decision_problem = scenario.decision_problem or "当前阶段是否应继续推进该概念。"
    evidence_items = ExternalEvidence.query.filter_by(scenario_id=id).order_by(ExternalEvidence.created_at.desc()).all()
    evidence_text = "\n".join([f"- {item.title}（{item.source_label or item.evidence_type}）：{item.content}" for item in evidence_items[:5]])

    def build_participant_prompt(participant, topic, history_text, opening=False):
        interaction_types = ["回应", "赞同", "补充", "提问", "建议", "质疑", "要求举证", "保留意见", "改变观点"]
        interaction_type = random.choice(interaction_types)
        base_intro = f"""你是虚拟焦点访谈中的真实参与者，请始终保持角色稳定，不要脱离人设。

当前研究目标：{research_goal}
核心决策问题：{decision_problem}
当前讨论话题：{topic}
角色状态：
{participant_state_prompt(participant)}
外部证据（如果有）：
{evidence_text or '暂无外部证据'}
"""
        if opening:
            return (
                f"""{base_intro}
请根据你的长期立场，先做第一轮自然发言。你可以表达期待、顾虑、使用场景或反对点。
要求：
- 50到120字
- 口语化，但保留职业和消费习惯
- 可以有犹豫和保留，不要一味正面
- 不要总结成报告语气，只输出发言内容""",
                "你扮演的是一个真实用户，会受预算、风险偏好、品牌偏好和个人经历影响，说法要稳定。"
            )

        return (
            f"""{base_intro}
最近讨论内容：
{history_text}

这一次你的互动意图是：{interaction_type}
请基于你的稳定立场和刚才大家的发言继续讨论。你可以同意、反驳、追问、要求举证、举例、改变观点或保留意见。
要求：
- 60到150字
- 尽量点到某个人的观点或某个证据
- 如果改变观点，要说清为什么
- 不要像报告，保持真实对话感
- 只输出你的发言内容""",
            "你不是客服，也不是分析师。你是在开会，会有情绪、偏见、犹豫和立场变化，但不能前后自相矛盾。"
        )

    def update_participant_state(participant, response_text, round_num):
        stance_state = safe_json_loads(participant.stance_state, {})
        stance_state["current_position"] = response_text[:180]
        stance_state["last_updated_round"] = round_num + 1
        if re.search(r'我改|我现在更倾向|被说服|我接受', response_text):
            stance_state["confidence"] = "变化中"
        participant.stance_state = json_dumps(stance_state, {})
        db.session.add(participant)

    # message_count模式：生成指定数量的消息
    if message_count > 0:
        target_count = message_count
        current_count = 0
        participant_index = 0
        
        while current_count < target_count:
            time.sleep(1)  # 控制对话节奏
            
            participant = participants[participant_index % len(participants)]
            round_num = current_count // len(participants)
            
            # 第一轮发言使用初始prompt，后续轮次使用带历史记录的prompt
            if current_count == 0:
                prompt, system_prompt = build_participant_prompt(participant, topics[round_num % len(topics)], '', True)
            else:
                history_text = "\n".join([f"{h['participant']}：{h['content']}" for h in conversation_history[-8:]])
                prompt, system_prompt = build_participant_prompt(participant, topics[round_num % len(topics)], history_text, False)

            response = call_deepseek(prompt, system_prompt)

            # AI调用失败时使用备用回复
            if response.startswith(("API Error", "Request Error", "请求失败")):
                error_responses = [
                    f"这个问题我有点想法...（停顿）嗯，我觉得可以从另一个角度来看待。",
                    f"我这边网络不太稳定，不过关于这个话题，我的看法是...",
                    f"抱歉，刚才信号不太好。我想补充一点，{participant.persona_name}认为这个产品还是很有潜力的。",
                    f"让我整理一下思路...我觉得大家刚才的讨论很有价值。",
                    f"（清了清嗓子）关于这个问题，我想分享一下我的看法。"
                ]
                response = random.choice(error_responses)
                print(f"AI请求失败，使用备用回复: {response}")

            # 保存对话记录
            record = ConversationRecord(
                scenario_id=id,
                participant_id=participant.id,
                content=response,
                is_host=False,
                message_type='normal',
                interaction_intent='discussion'
            )
            db.session.add(record)
            update_participant_state(participant, response, round_num)
            db.session.commit()

            conversation_history.append({"participant": participant.persona_name, "content": response})
            current_count += 1
            participant_index += 1

            # 每轮结束后主持人引导
            if current_count < target_count and current_count % len(participants) == 0 and len(conversation_history) >= 2:
                time.sleep(1)
                recent_talkers = [h['participant'] for h in conversation_history[-4:]]
                host_prompt = f"""你是这场企业研究型焦点小组讨论的主持人。刚才发言的朋友们聊得很热烈：

最近发言的朋友：{', '.join(recent_talkers)}
研究目标：{research_goal}
核心决策问题：{decision_problem}

你需要像研究主持人一样推进讨论：
- 可以邀请还没怎么发言的朋友说说看法
- 可以追问某个模糊观点背后的场景、原因或证据
- 如果大家过早达成一致，要提醒他们说出顾虑和反例
- 如果有人提到外部事实，可以要求更具体一点

请用口语化、亲切但专业的方式说几句，30-80字就够了。不要像在做报告。"""

                host_response = call_deepseek(host_prompt, "你是一位亲切自然的主持人，说话不像是在念稿，更像是和一群朋友聊天。语气轻松但专业，能够让每个人都感到舒适和被尊重。")

                if host_response.startswith(("API Error", "Request Error", "请求失败")):
                    host_response = random.choice([
                        "好的，那我们继续下一个话题吧。",
                        "大家说得都很有道理！那接下来..."
                    ])
                    print(f"主持人请求失败，使用备用回复: {host_response}")

                guide_msg = ConversationRecord(
                    scenario_id=id,
                    participant_id=0,
                    content=host_response,
                    is_host=True,
                    message_type='guide'
                )
                db.session.add(guide_msg)
                db.session.commit()
                conversation_history.append({"participant": "主持人", "content": host_response})
    else:
        # rounds模式：按轮数进行讨论
        for round_num in range(rounds):
            time.sleep(1)

            if round_num == 0:
                # 第一轮：每位参与者初次发言
                for participant in participants:
                    time.sleep(1.5)

                    prompt = f"""想象你就是一个真实的消费者"{participant.persona_name}"，正在参加一个产品讨论会。

你的背景：{participant.persona_name}，{participant.background}
你的性格特点：{participant.personality}
别人对你的印象标签：{','.join(json.loads(participant.persona_tags))}

我们今天在讨论的产品是「{scenario.product_name}」，它的概念是：{scenario.product_concept}
这个产品的核心卖点包括：{','.join(selling_points)}

现在大家正在聊的话题是：{topics[round_num % len(topics)]}

请像平时和朋友聊天一样，说说你对这个产品的看法。可以是你的期待、疑惑、或者你觉得需要改进的地方。用你自己最舒服的方式表达就好，不用刻意组织语言，真实自然一点。
（请控制在50-120字，直接输出你说的话即可）"""

                    system_prompt = "你扮演的是一个真实的、有血有肉的消费者，不是AI。在讨论中会自然地流露情感，有时候会有语气词，有时候会说一些看似跑题但实际相关的话。"

                    response = call_deepseek(prompt, system_prompt)

                    if response.startswith(("API Error", "Request Error", "请求失败")):
                        response = random.choice([
                            f"我觉得这个产品挺有意思的，{participant.persona_name}表示很期待。",
                            f"关于这个话题，我有一些想法想分享...",
                            f"嗯，让我想想，我觉得这个产品的定位很清晰。"
                        ])
                        print(f"AI请求失败，使用备用回复: {response}")

                    record = ConversationRecord(
                        scenario_id=id,
                        participant_id=participant.id,
                        content=response,
                        is_host=False,
                        message_type='normal',
                        interaction_intent='discussion'
                    )
                    db.session.add(record)
                    update_participant_state(participant, response, round_num)
                    db.session.commit()

                    conversation_history.append({"participant": participant.persona_name, "content": response})
            else:
                # 后续轮次：基于历史记录进行互动
                for participant in participants:
                    time.sleep(1.5)

                    history_text = "\n".join([f"{h['participant']}：{h['content']}" for h in conversation_history[-8:]])

                    interaction_types = ["追问", "反驳", "补充", "同意", "质疑", "举例"]
                    interaction_type = random.choice(interaction_types)

                    prompt = f"""你"{participant.persona_name}"正在和一个产品团队开会讨论产品。

你的身份背景：{participant.background}
你的性格：{participant.personality}
你的标签：{','.join(json.loads(participant.persona_tags))}

今天讨论的产品是「{scenario.product_name}」：
- 产品概念：{scenario.product_concept}
- 核心卖点：{','.join(selling_points)}

之前的讨论内容：
{history_text}

现在大家正在讨论的话题是：{topics[round_num % len(topics)]}

作为"{participant.persona_name}"，你听了大家的发言后，也想说几句。用你自己的话说，不用拘泥于格式，可以：
- 同意某人的观点并补充
- 提出不同看法
- 分享一个相关的亲身经历
- 问个问题

请自然地参与讨论，60-150字，用口语化的方式表达。"""

                    system_prompt = "你扮演的是一个真实的消费者，有自己的观点和情绪。会因为同意别人的话而点头，也会因为不认同的观点而皱眉。说话时会不自觉地带出自己的职业习惯和个人经历。"

                    response = call_deepseek(prompt, system_prompt)

                    if response.startswith(("API Error", "Request Error", "请求失败")):
                        response = random.choice([
                            f"我同意刚才那位朋友的看法，{participant.persona_name}也有类似的感受。",
                            f"这个点我想补充一下...",
                            f"（思考了一下）我觉得可以从另一个角度来看这个问题。"
                        ])
                        print(f"AI请求失败，使用备用回复: {response}")

                    record = ConversationRecord(
                        scenario_id=id,
                        participant_id=participant.id,
                        content=response,
                        is_host=False,
                        message_type='normal'
                    )
                    db.session.add(record)
                    db.session.commit()

                    conversation_history.append({"participant": participant.persona_name, "content": response})

            # 主持人引导（非最后一轮）
            if round_num < rounds - 1 and len(conversation_history) >= 2:
                time.sleep(1)
                recent_talkers = [h['participant'] for h in conversation_history[-4:]]
                host_prompt = f"""你是这场焦点小组讨论的主持人。刚才发言的朋友们聊得很热烈：

最近发言的朋友：{', '.join(recent_talkers)}

你作为一个有经验的主持人，需要把讨论继续推进下去。不要说太多，自然地引导一下：
- 可以邀请还没怎么发言的朋友说说看法
- 或者针对刚才聊到的某个点，追问一下
- 也可以简单总结一下大家说的方向

请用口语化、亲切的方式说几句，30-60字就够了。就像朋友间聊天那样自然，不要像在做报告。"""

                host_response = call_deepseek(host_prompt, "你是一位亲切自然的主持人，说话不像是在念稿，更像是和一群朋友聊天。语气轻松但专业，能够让每个人都感到舒适和被尊重。")

                if host_response.startswith(("API Error", "Request Error", "请求失败")):
                    host_response = random.choice([
                        "好的，那我们继续下一个话题吧。",
                        "大家说得都很有道理！那接下来..."
                    ])
                    print(f"主持人请求失败，使用备用回复: {host_response}")

                guide_msg = ConversationRecord(
                    scenario_id=id,
                    participant_id=0,
                    content=host_response,
                    is_host=True,
                    message_type='guide'
                )
                db.session.add(guide_msg)
                db.session.commit()
                conversation_history.append({"participant": "主持人", "content": host_response})

            # 每轮结束小结（非最后一轮）
            if round_num < rounds - 1:
                time.sleep(0.5)
                summary_msg = ConversationRecord(
                    scenario_id=id,
                    participant_id=0,
                    content="好，刚才大家聊得挺好的，我先简单小结一下...（停顿）嗯，总的来说大家对这个产品的功能还是比较认可的，尤其是某某方面...那我们接下来聊聊下一个话题吧。",
                    is_host=True,
                    message_type='summary'
                )
                db.session.add(summary_msg)
                db.session.commit()

    # 讨论结束，主持人总结
    conclusion_msg = ConversationRecord(
        scenario_id=id,
        participant_id=0,
        content="好啦，今天的讨论就到这里吧！非常感谢各位的积极参与，大家聊得都很真实、很深入，给我们提供了很多有价值的意见。等我整理一下今天的讨论，很快就会给大家一份完整的分析报告。辛苦各位了！",
        is_host=True,
        message_type='conclusion'
    )
    db.session.add(conclusion_msg)
    db.session.commit()

    result_data = {
        "status": "success",
        "message": f"对话模拟完成，共{rounds}轮讨论，{len(participants)}位参与者",
        "rounds": rounds,
        "participants_count": len(participants),
        "total_messages": ConversationRecord.query.filter_by(scenario_id=id).count()
    }
    return app.response_class(
        response=json.dumps(result_data, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/simulate/stream', methods=['GET'])
def simulate_conversation_stream(id):
    scenario, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response

    rounds = int(request.args.get('rounds', 2))
    message_count = int(request.args.get('message_count', 0))

    participants = VirtualParticipant.query.filter_by(scenario_id=id).all()
    if not participants:
        return app.response_class(
            response=json.dumps({"error": "请先生成参与者"}, ensure_ascii=False),
            status=400,
            mimetype='application/json; charset=utf-8'
        )

    def event_stream():
        selling_points = json.loads(scenario.core_selling_points)
        topics = json.loads(scenario.discussion_topics)

        ConversationRecord.query.filter_by(scenario_id=id).delete()

        host_intro = ConversationRecord(
            scenario_id=id, participant_id=0,
            content=f"好啦，各位朋友大家好！欢迎来参加今天的焦点小组讨论。\n\n今天我们主要想聊聊「{scenario.product_name}」这个产品。\n\n先简单说一下这个产品：{scenario.product_concept}\n\n它主打的卖点是：{','.join(selling_points)}\n\n我们就从第一个话题开始吧：{topics[0]}\n\n大家有什么想法就说什么，不用客气啊！",
            is_host=True, message_type='intro'
        )
        db.session.add(host_intro)
        db.session.commit()
        yield f"data: {json.dumps({'type': 'message', 'participant_name': 'AI主持人', 'content': host_intro.content, 'is_host': True, 'message_type': 'intro', 'participant_id': 0, 'timestamp': host_intro.timestamp.isoformat()}, ensure_ascii=False)}\n\n"

        topic_msg = ConversationRecord(
            scenario_id=id, participant_id=0,
            content=f"话题 1：{topics[0]}",
            is_host=True, message_type='topic'
        )
        db.session.add(topic_msg)
        db.session.commit()

        conversation_history = []
        total_participants = len(participants)

        if message_count > 0:
            target_count = message_count
            current_count = 0
            participant_index = 0
            total_steps = target_count + (target_count // total_participants)

            while current_count < target_count:
                yield f"data: {json.dumps({'type': 'progress', 'current': current_count + 1, 'total': target_count, 'phase': '模拟讨论中', 'pct': int((current_count + 1) / target_count * 100)}, ensure_ascii=False)}\n\n"
                time.sleep(0.5)

                participant = participants[participant_index % total_participants]
                round_num = current_count // total_participants

                if current_count == 0:
                    prompt = f"""想象你就是一个真实的消费者"{participant.persona_name}"，正在参加一个产品讨论会。\n\n你的背景：{participant.persona_name}，{participant.background}\n你的性格特点：{participant.personality}\n别人对你的印象标签：{','.join(json.loads(participant.persona_tags))}\n\n我们今天在讨论的产品是「{scenario.product_name}」，它的概念是：{scenario.product_concept}\n这个产品的核心卖点包括：{','.join(selling_points)}\n\n现在大家正在聊的话题是：{topics[round_num % len(topics)]}\n\n请像平时和朋友聊天一样，说说你对这个产品的看法。可以是你的期待、疑惑、或者你觉得需要改进的地方。用你自己最舒服的方式表达就好，不用刻意组织语言，真实自然一点。\n（请控制在50-120字，直接输出你说的话即可）"""
                    system_prompt = "你扮演的是一个真实的、有血有肉的消费者，不是AI。在讨论中会自然地流露情感，有时候会有语气词，有时候会说一些看似跑题但实际相关的话。"
                else:
                    history_text = "\n".join([f"{h['participant']}：{h['content']}" for h in conversation_history[-8:]])
                    prompt = f"""你"{participant.persona_name}"正在和一个产品团队开会讨论产品。\n\n你的身份背景：{participant.background}\n你的性格：{participant.personality}\n你的标签：{','.join(json.loads(participant.persona_tags))}\n\n今天讨论的产品是「{scenario.product_name}」：\n- 产品概念：{scenario.product_concept}\n- 核心卖点：{','.join(selling_points)}\n\n之前的讨论内容：\n{history_text}\n\n现在大家正在讨论的话题是：{topics[round_num % len(topics)]}\n\n作为"{participant.persona_name}"，你听了大家的发言后，也想说几句。用你自己的话说，不用拘泥于格式，可以：\n- 同意某人的观点并补充\n- 提出不同看法\n- 分享一个相关的亲身经历\n- 问个问题\n\n请自然地参与讨论，60-150字，用口语化的方式表达。"""
                    system_prompt = "你扮演的是一个真实的消费者，有自己的观点和情绪。会因为同意别人的话而点头，也会因为不认同的观点而皱眉。说话时会不自觉地带出自己的职业习惯和个人经历。"

                response = call_deepseek(prompt, system_prompt)
                if response.startswith(("API Error", "Request Error", "请求失败")):
                    response = random.choice([
                        f"这个问题我有点想法...（停顿）嗯，我觉得可以从另一个角度来看待。",
                        f"我这边网络不太稳定，不过关于这个话题，我的看法是...",
                        f"{participant.persona_name}认为这个产品还是很有潜力的。",
                        f"让我整理一下思路...我觉得大家刚才的讨论很有价值。",
                        f"（清了清嗓子）关于这个问题，我想分享一下我的看法。"
                    ])

                record = ConversationRecord(
                    scenario_id=id, participant_id=participant.id,
                    content=response, is_host=False, message_type='normal'
                )
                db.session.add(record)
                db.session.commit()
                conversation_history.append({"participant": participant.persona_name, "content": response})

                yield f"data: {json.dumps({'type': 'message', 'participant_name': participant.persona_name, 'content': response, 'is_host': False, 'message_type': 'normal', 'participant_id': participant.id, 'timestamp': record.timestamp.isoformat()}, ensure_ascii=False)}\n\n"

                current_count += 1
                participant_index += 1

                if current_count < target_count and current_count % total_participants == 0 and len(conversation_history) >= 2:
                    recent_talkers = [h['participant'] for h in conversation_history[-4:]]
                    host_prompt = f"""你是这场焦点小组讨论的主持人。刚才发言的朋友们聊得很热烈：\n\n最近发言的朋友：{', '.join(recent_talkers)}\n\n你作为一个有经验的主持人，需要把讨论继续推进下去。不要说太多，自然地引导一下：\n- 可以邀请还没怎么发言的朋友说说看法\n- 或者针对刚才聊到的某个点，追问一下\n- 也可以简单总结一下大家说的方向\n\n请用口语化、亲切的方式说几句，30-60字就够了。"""
                    host_response = call_deepseek(host_prompt, "你是一位亲切自然的主持人，说话不像是在念稿，更像是和一群朋友聊天。")
                    if host_response.startswith(("API Error", "Request Error", "请求失败")):
                        host_response = random.choice(["好的，那我们继续下一个话题吧。", "大家说得都很有道理！那接下来..."])

                    guide_msg = ConversationRecord(
                        scenario_id=id, participant_id=0,
                        content=host_response, is_host=True, message_type='guide'
                    )
                    db.session.add(guide_msg)
                    db.session.commit()
                    conversation_history.append({"participant": "主持人", "content": host_response})
                    yield f"data: {json.dumps({'type': 'message', 'participant_name': 'AI主持人', 'content': host_response, 'is_host': True, 'message_type': 'guide', 'participant_id': 0, 'timestamp': guide_msg.timestamp.isoformat()}, ensure_ascii=False)}\n\n"
        else:
            total_steps = rounds * total_participants + rounds
            step = 0
            for round_num in range(rounds):
                if round_num == 0:
                    for participant in participants:
                        step += 1
                        yield f"data: {json.dumps({'type': 'progress', 'current': step, 'total': total_steps, 'phase': f'第 {round_num + 1} 轮讨论', 'pct': int(step / total_steps * 100)}, ensure_ascii=False)}\n\n"
                        time.sleep(1)

                        prompt = f"""想象你就是一个真实的消费者"{participant.persona_name}"，正在参加一个产品讨论会。\n\n你的背景：{participant.persona_name}，{participant.background}\n你的性格特点：{participant.personality}\n别人对你的印象标签：{','.join(json.loads(participant.persona_tags))}\n\n我们今天在讨论的产品是「{scenario.product_name}」，它的概念是：{scenario.product_concept}\n这个产品的核心卖点包括：{','.join(selling_points)}\n\n现在大家正在聊的话题是：{topics[round_num % len(topics)]}\n\n请像平时和朋友聊天一样，说说你对这个产品的看法。可以是你的期待、疑惑、或者你觉得需要改进的地方。用你自己最舒服的方式表达就好，不用刻意组织语言，真实自然一点。\n（请控制在50-120字，直接输出你说的话即可）"""
                        system_prompt = "你扮演的是一个真实的、有血有肉的消费者，不是AI。在讨论中会自然地流露情感，有时候会有语气词，有时候会说一些看似跑题但实际相关的话。"
                        response = call_deepseek(prompt, system_prompt)
                        if response.startswith(("API Error", "Request Error", "请求失败")):
                            response = random.choice([f"我觉得这个产品挺有意思的，{participant.persona_name}表示很期待。", f"关于这个话题，我有一些想法想分享...", f"嗯，让我想想，我觉得这个产品的定位很清晰。"])

                        record = ConversationRecord(
                            scenario_id=id, participant_id=participant.id,
                            content=response, is_host=False, message_type='normal'
                        )
                        db.session.add(record)
                        db.session.commit()
                        conversation_history.append({"participant": participant.persona_name, "content": response})

                        yield f"data: {json.dumps({'type': 'message', 'participant_name': participant.persona_name, 'content': response, 'is_host': False, 'message_type': 'normal', 'participant_id': participant.id, 'timestamp': record.timestamp.isoformat()}, ensure_ascii=False)}\n\n"
                else:
                    for participant in participants:
                        step += 1
                        yield f"data: {json.dumps({'type': 'progress', 'current': step, 'total': total_steps, 'phase': f'第 {round_num + 1} 轮讨论', 'pct': int(step / total_steps * 100)}, ensure_ascii=False)}\n\n"
                        time.sleep(1)

                        history_text = "\n".join([f"{h['participant']}：{h['content']}" for h in conversation_history[-8:]])
                        prompt = f"""你"{participant.persona_name}"正在和一个产品团队开会讨论产品。\n\n你的身份背景：{participant.background}\n你的性格：{participant.personality}\n你的标签：{','.join(json.loads(participant.persona_tags))}\n\n今天讨论的产品是「{scenario.product_name}」：\n- 产品概念：{scenario.product_concept}\n- 核心卖点：{','.join(selling_points)}\n\n之前的讨论内容：\n{history_text}\n\n现在大家正在讨论的话题是：{topics[round_num % len(topics)]}\n\n作为"{participant.persona_name}"，你听了大家的发言后，也想说几句。用你自己的话说，不用拘泥于格式，可以：\n- 同意某人的观点并补充\n- 提出不同看法\n- 分享一个相关的亲身经历\n- 问个问题\n\n请自然地参与讨论，60-150字，用口语化的方式表达。"""
                        system_prompt = "你扮演的是一个真实的消费者，有自己的观点和情绪。会因为同意别人的话而点头，也会因为不认同的观点而皱眉。说话时会不自觉地带出自己的职业习惯和个人经历。"
                        response = call_deepseek(prompt, system_prompt)
                        if response.startswith(("API Error", "Request Error", "请求失败")):
                            response = random.choice([f"我同意刚才那位朋友的看法，{participant.persona_name}也有类似的感受。", f"这个点我想补充一下...", f"（思考了一下）我觉得可以从另一个角度来看这个问题。"])

                        record = ConversationRecord(
                            scenario_id=id, participant_id=participant.id,
                            content=response, is_host=False, message_type='normal'
                        )
                        db.session.add(record)
                        db.session.commit()
                        conversation_history.append({"participant": participant.persona_name, "content": response})

                        yield f"data: {json.dumps({'type': 'message', 'participant_name': participant.persona_name, 'content': response, 'is_host': False, 'message_type': 'normal', 'participant_id': participant.id, 'timestamp': record.timestamp.isoformat()}, ensure_ascii=False)}\n\n"

                if round_num < rounds - 1 and len(conversation_history) >= 2:
                    step += 1
                    yield f"data: {json.dumps({'type': 'progress', 'current': step, 'total': total_steps, 'phase': '主持人引导', 'pct': int(step / total_steps * 100)}, ensure_ascii=False)}\n\n"
                    recent_talkers = [h['participant'] for h in conversation_history[-4:]]
                    host_prompt = f"""你是这场焦点小组讨论的主持人。刚才发言的朋友们聊得很热烈：\n\n最近发言的朋友：{', '.join(recent_talkers)}\n\n你作为一个有经验的主持人，需要把讨论继续推进下去。不要说太多，自然地引导一下：\n- 可以邀请还没怎么发言的朋友说说看法\n- 或者针对刚才聊到的某个点，追问一下\n- 也可以简单总结一下大家说的方向\n\n请用口语化、亲切的方式说几句，30-60字就够了。"""
                    host_response = call_deepseek(host_prompt, "你是一位亲切自然的主持人，说话不像是在念稿，更像是和一群朋友聊天。")
                    if host_response.startswith(("API Error", "Request Error", "请求失败")):
                        host_response = random.choice(["好的，那我们继续下一个话题吧。", "大家说得都很有道理！那接下来..."])
                    guide_msg = ConversationRecord(
                        scenario_id=id, participant_id=0,
                        content=host_response, is_host=True, message_type='guide'
                    )
                    db.session.add(guide_msg)
                    db.session.commit()
                    conversation_history.append({"participant": "主持人", "content": host_response})
                    yield f"data: {json.dumps({'type': 'message', 'participant_name': 'AI主持人', 'content': host_response, 'is_host': True, 'message_type': 'guide', 'participant_id': 0, 'timestamp': guide_msg.timestamp.isoformat()}, ensure_ascii=False)}\n\n"

        conclusion_msg = ConversationRecord(
            scenario_id=id, participant_id=0,
            content="好啦，今天的讨论就到这里吧！非常感谢各位的积极参与，大家聊得都很真实、很深入，给我们提供了很多有价值的意见。等我整理一下今天的讨论，很快就会给大家一份完整的分析报告。辛苦各位了！",
            is_host=True, message_type='conclusion'
        )
        db.session.add(conclusion_msg)
        db.session.commit()
        yield f"data: {json.dumps({'type': 'message', 'participant_name': 'AI主持人', 'content': conclusion_msg.content, 'is_host': True, 'message_type': 'conclusion', 'participant_id': 0, 'timestamp': conclusion_msg.timestamp.isoformat()}, ensure_ascii=False)}\n\n"

        total_msgs = ConversationRecord.query.filter_by(scenario_id=id).count()
        yield f"data: {json.dumps({'type': 'done', 'total_messages': total_msgs}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(event_stream()),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive'
        }
    )


@app.route('/api/scenarios/<int:id>/simulate/stream-v2', methods=['GET'])
def simulate_conversation_stream_v2(id):
    scenario, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response
    standalone_error = reject_if_not_standalone(scenario, "会议室内部场景不能使用普通自动模拟，请前往会议室页面继续。")
    if standalone_error:
        return standalone_error

    rounds = max(1, int(request.args.get('rounds', 2)))
    message_count = max(0, int(request.args.get('message_count', 0)))

    participants = get_standalone_participants(id)
    if not participants:
        return app.response_class(
            response=json.dumps({"error": "请先生成参与者"}, ensure_ascii=False),
            status=400,
            mimetype='application/json; charset=utf-8'
        )

    def event_stream():
        selling_points = safe_json_loads(scenario.core_selling_points, [])
        topics = safe_json_loads(scenario.discussion_topics, [])
        if not topics:
            topics = ["用户真实需求", "使用门槛与风险", "是否值得继续推进"]

        evidence_items = ExternalEvidence.query.filter_by(scenario_id=id).order_by(ExternalEvidence.created_at.desc()).all()
        evidence_lines = []
        for item in evidence_items[:5]:
            source_name = item.source_title or item.source_label or item.evidence_type or "外部资料"
            evidence_lines.append(f"- {source_name}：{item.content}")
        evidence_text = "\n".join(evidence_lines)

        ConversationRecord.query.filter_by(scenario_id=id).delete()
        db.session.commit()

        for participant in participants:
            normalize_participant_profile(participant, scenario)
        db.session.commit()

        intro_text = build_quality_host_intro(scenario, selling_points, topics[0])
        host_intro = ConversationRecord(
            scenario_id=id,
            participant_id=0,
            content=intro_text,
            is_host=True,
            message_type='intro',
            interaction_intent='host_intro'
        )
        db.session.add(host_intro)
        db.session.commit()
        yield f"data: {json.dumps({'type': 'message', 'participant_name': 'AI主持人', 'content': intro_text, 'is_host': True, 'message_type': 'intro', 'participant_id': 0, 'timestamp': host_intro.timestamp.isoformat()}, ensure_ascii=False)}\n\n"

        topic_msg = ConversationRecord(
            scenario_id=id,
            participant_id=0,
            content=f"当前讨论主题：{topics[0]}",
            is_host=True,
            message_type='topic',
            interaction_intent='host_topic'
        )
        db.session.add(topic_msg)
        db.session.commit()
        yield f"data: {json.dumps({'type': 'message', 'participant_name': 'AI主持人', 'content': topic_msg.content, 'is_host': True, 'message_type': 'topic', 'participant_id': 0, 'timestamp': topic_msg.timestamp.isoformat()}, ensure_ascii=False)}\n\n"

        conversation_history = []
        total_participants = len(participants)

        if message_count > 0:
            target_count = message_count
        else:
            target_count = rounds * total_participants

        host_insert_points = set()
        if total_participants > 0 and target_count > total_participants:
            for idx in range(total_participants, target_count, total_participants):
                host_insert_points.add(idx)

        current_count = 0
        participant_index = 0
        total_steps = target_count + len(host_insert_points) + 1
        step = 0

        while current_count < target_count:
            step += 1
            current_topic = topics[min(current_count // total_participants, len(topics) - 1)]
            yield f"data: {json.dumps({'type': 'progress', 'current': step, 'total': total_steps, 'phase': '正在生成高质量讨论', 'pct': int(step / total_steps * 100)}, ensure_ascii=False)}\n\n"
            time.sleep(0.3)

            participant = participants[participant_index % total_participants]
            round_num = current_count // total_participants
            history_text = "\n".join([f"{h['participant']}：{h['content']}" for h in conversation_history[-8:]])
            prompt, system_prompt, interaction_type = build_quality_participant_prompt(
                participant,
                scenario,
                current_topic,
                history_text,
                evidence_text,
                opening=(current_count < total_participants)
            )
            response = call_deepseek(prompt, system_prompt)
            if is_ai_error_text(response):
                response = build_quality_fallback_response(
                    participant,
                    scenario,
                    current_topic,
                    opening=(current_count < total_participants)
                )

            record = ConversationRecord(
                scenario_id=id,
                participant_id=participant.id,
                content=response,
                is_host=False,
                message_type='normal',
                interaction_intent=interaction_type
            )
            db.session.add(record)
            update_participant_state_v2(participant, response, round_num)
            db.session.commit()

            conversation_history.append({"participant": participant.persona_name, "content": response})
            yield f"data: {json.dumps({'type': 'message', 'participant_name': participant.persona_name, 'content': response, 'is_host': False, 'message_type': 'normal', 'participant_id': participant.id, 'timestamp': record.timestamp.isoformat()}, ensure_ascii=False)}\n\n"

            current_count += 1
            participant_index += 1

            if current_count in host_insert_points and len(conversation_history) >= 2:
                step += 1
                yield f"data: {json.dumps({'type': 'progress', 'current': step, 'total': total_steps, 'phase': '主持人正在追问关键分歧', 'pct': int(step / total_steps * 100)}, ensure_ascii=False)}\n\n"
                host_prompt, host_system_prompt = build_quality_host_prompt(scenario, current_topic, conversation_history)
                host_response = call_deepseek(host_prompt, host_system_prompt)
                if is_ai_error_text(host_response):
                    host_response = build_quality_host_fallback(scenario, current_topic)
                guide_msg = ConversationRecord(
                    scenario_id=id,
                    participant_id=0,
                    content=host_response,
                    is_host=True,
                    message_type='guide',
                    interaction_intent='host_probe'
                )
                db.session.add(guide_msg)
                db.session.commit()
                conversation_history.append({"participant": "AI主持人", "content": host_response})
                yield f"data: {json.dumps({'type': 'message', 'participant_name': 'AI主持人', 'content': host_response, 'is_host': True, 'message_type': 'guide', 'participant_id': 0, 'timestamp': guide_msg.timestamp.isoformat()}, ensure_ascii=False)}\n\n"

        conclusion = (
            "今天这轮讨论先到这里。现在我们已经收集到了一些支持推进的理由，也暴露了需要继续验证的风险点。"
            "下一步我会基于这些分歧和证据，整理成更适合企业决策阅读的分析报告。"
        )
        conclusion_msg = ConversationRecord(
            scenario_id=id,
            participant_id=0,
            content=conclusion,
            is_host=True,
            message_type='conclusion',
            interaction_intent='host_conclusion'
        )
        db.session.add(conclusion_msg)
        db.session.commit()
        yield f"data: {json.dumps({'type': 'message', 'participant_name': 'AI主持人', 'content': conclusion, 'is_host': True, 'message_type': 'conclusion', 'participant_id': 0, 'timestamp': conclusion_msg.timestamp.isoformat()}, ensure_ascii=False)}\n\n"

        total_msgs = ConversationRecord.query.filter_by(scenario_id=id).count()
        yield f"data: {json.dumps({'type': 'done', 'total_messages': total_msgs}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(event_stream()),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive'
        }
    )

@app.route('/api/scenarios/<int:id>/conversation', methods=['GET'])
def get_conversation(id):
    scenario, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response
    
    records = ConversationRecord.query.filter_by(scenario_id=id).order_by(ConversationRecord.timestamp).all()

    participants = get_room_managed_participants(id) if not is_standalone_scenario(scenario) else get_standalone_participants(id)
    participant_map = {p.id: p.persona_name for p in participants}
    participant_map[0] = "AI主持人"

    result = []
    for record in records:
        if record.participant_id == -1:
            participant_name = "我"
        else:
            participant_name = participant_map.get(record.participant_id, "未知")

        result.append({
            "id": record.id,
            "participant_id": record.participant_id,
            "participant_name": participant_name,
            "content": record.content,
            "timestamp": record.timestamp.isoformat(),
            "is_host": record.is_host,
            "message_type": record.message_type
        })

    return app.response_class(
        response=json.dumps(result, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/conversation', methods=['POST'])
def add_message(id):
    scenario, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response
    standalone_error = reject_if_not_standalone(scenario, "该场景已绑定会议室，请在会议室页面继续正式讨论。")
    if standalone_error:
        return standalone_error

    meeting_room = MeetingRoom.query.filter_by(scenario_id=id).first()
    if meeting_room and (scenario.meeting_status or 'standalone') != 'standalone':
        return json_api_error("该场景已绑定会议室，请在会议室页面继续正式讨论", 400)
    
    data = request.get_json()
    content = data.get('content', '').strip()
    participant_name = data.get('participant_name', '我')

    if not content:
        return app.response_class(
            response=json.dumps({"error": "消息内容不能为空"}, ensure_ascii=False),
            status=400,
            mimetype='application/json; charset=utf-8'
        )

    participants = VirtualParticipant.query.filter_by(scenario_id=id).all()
    
    if not participants:
        return app.response_class(
            response=json.dumps({"error": "该场景没有参与者"}, ensure_ascii=False),
            status=400,
            mimetype='application/json; charset=utf-8'
        )

    record = ConversationRecord(
        scenario_id=id,
        participant_id=-1,
        content=content,
        is_host=False,
        message_type='user'
    )
    db.session.add(record)
    db.session.commit()

    selling_points = json.loads(scenario.core_selling_points)
    topics = json.loads(scenario.discussion_topics)

    records = ConversationRecord.query.filter_by(scenario_id=id).order_by(ConversationRecord.timestamp).all()
    conversation_history = []
    for r in records[-10:]:
        if r.is_host:
            conversation_history.append({"participant": "主持人", "content": r.content})
        elif r.participant_id == -1:
            conversation_history.append({"participant": "我", "content": r.content})
        else:
            participant = next((p for p in participants if p.id == r.participant_id), None)
            if participant:
                conversation_history.append({"participant": participant.persona_name, "content": r.content})

    import random
    num_responses = min(random.randint(1, 3), len(participants))
    selected_participants = random.sample(participants, num_responses)

    for participant in selected_participants:
        history_text = "\n".join([f"{h['participant']}：{h['content']}" for h in conversation_history[-8:]])
        interaction_types = ["回应", "赞同", "补充", "提问", "建议"]
        interaction_type = random.choice(interaction_types)

        prompt = f"""你"{participant.persona_name}"正在参加一个关于产品的讨论会，刚才有个用户发表了自己的看法。

你的身份背景：{participant.persona_name}，{participant.background}
你的性格特点：{participant.personality}
你的标签：{','.join(json.loads(participant.persona_tags))}

我们今天讨论的产品是「{scenario.product_name}」：
- 产品概念：{scenario.product_concept}
- 核心卖点：{','.join(selling_points)}

之前的讨论：
{history_text}

用户刚才说："{content}"

听到用户的话后，你作为"{participant.persona_name}"想回应一下。你可以说：
- 对用户观点的认同或补充
- 结合自己经历的一些想法
- 提出一个相关的问题

请用口语化的方式自然地回应，60-150字，不用太长。"""

        system_prompt = "你是一个真实的消费者，说话有自己的风格。不会刻意迎合，但也不会无礼。会因为同意而认可，也会因为不同意见而表达疑虑。"

        response = call_deepseek(prompt, system_prompt)

        if response.startswith(("API Error", "Request Error", "请求失败")):
            response = random.choice([
                f"{participant.persona_name}点头表示赞同：\"嗯，你说得有道理。\"",
                f"{participant.persona_name}若有所思地说：\"这个观点很有意思...\"",
                f"{participant.persona_name}：\"我也有类似的看法。\""
            ])
            print(f"AI请求失败，使用备用回复: {response}")

        ai_record = ConversationRecord(
            scenario_id=id,
            participant_id=participant.id,
            content=response,
            is_host=False,
            message_type='normal'
        )
        db.session.add(ai_record)
        db.session.commit()

        conversation_history.append({"participant": participant.persona_name, "content": response})

    return app.response_class(
        response=json.dumps({"status": "success", "message": "消息已添加", "ai_responses": num_responses}, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

# ---------------------------
# 报告生成API
# ---------------------------

@app.route('/api/scenarios/<int:id>/generate-report', methods=['POST'])
def generate_report_with_ai(id):
    log_report_event('report_generation_started', scenario_id=id, mode='sync')
    report_inputs, error_response = build_report_inputs(id)
    if error_response:
        log_report_event(
            'report_generation_blocked',
            scenario_id=id,
            mode='sync',
            http_status=getattr(error_response, 'status_code', None),
            error=extract_error_message_from_response(error_response),
            level='WARNING'
        )
        return error_response

    report, report_error = generate_analysis_report_bundle(report_inputs)
    if report_error:
        log_report_event(
            'report_generation_failed',
            scenario_id=id,
            mode='sync',
            http_status=502,
            error=report_error,
            level='ERROR'
        )
        return app.response_class(
            response=json.dumps({"error": report_error}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    log_report_event(
        'report_generation_completed',
        scenario_id=id,
        mode='sync',
        report_id=getattr(report, 'id', None)
    )
    return app.response_class(
        response=json.dumps(report.to_dict(), ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/generate-report/stream', methods=['GET'])
def generate_report_stream(id):
    log_report_event('report_generation_started', scenario_id=id, mode='stream')
    report_inputs, error_response = build_report_inputs(id)
    if error_response:
        log_report_event(
            'report_generation_blocked',
            scenario_id=id,
            mode='stream',
            http_status=getattr(error_response, 'status_code', None),
            error=extract_error_message_from_response(error_response),
            level='WARNING'
        )
        return error_response

    def event_stream():
        event_queue = Queue()
        result_holder = {'report': None, 'error': None}

        @copy_current_request_context
        def worker():
            def capture_progress(stage_code, phase, pct):
                event_queue.put({
                    'type': 'progress',
                    'stage_code': stage_code,
                    'phase': phase,
                    'pct': pct,
                })

            try:
                report, report_error = generate_analysis_report_bundle(report_inputs, progress_callback=capture_progress)
                result_holder['report'] = report
                result_holder['error'] = report_error
            except Exception as exc:
                result_holder['error'] = str(exc)
                log_report_event(
                    'report_generation_worker_exception',
                    scenario_id=id,
                    mode='stream',
                    error=str(exc),
                    level='ERROR'
                )
            finally:
                event_queue.put(None)

        worker_thread = threading.Thread(target=worker, daemon=True)
        worker_thread.start()

        while True:
            event = event_queue.get()
            if event is None:
                break
            yield f"data: {json.dumps(event, ensure_ascii=False)}\n\n"

        report = result_holder.get('report')
        report_error = result_holder.get('error')
        if report_error:
            log_report_event(
                'report_generation_failed',
                scenario_id=id,
                mode='stream',
                error=report_error,
                level='ERROR'
            )
            yield f"data: {json.dumps({'type': 'error', 'error': report_error}, ensure_ascii=False)}\n\n"
            return

        log_report_event(
            'report_generation_completed',
            scenario_id=id,
            mode='stream',
            report_id=getattr(report, 'id', None)
        )
        yield f"data: {json.dumps({'type': 'done', 'report': report.to_dict()}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(event_stream()),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive'
        }
    )

@app.route('/api/scenarios/<int:id>/report', methods=['GET'])
def get_scenario_report(id):
    _, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response
    
    report = AnalysisReport.query.filter_by(scenario_id=id).first()
    if not report:
        return app.response_class(
            response=json.dumps({"error": "暂无报告，请先生成"}, ensure_ascii=False),
            status=404,
            mimetype='application/json; charset=utf-8'
        )
    return app.response_class(
        response=json.dumps(report.to_dict(), ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/reports/<int:id>', methods=['GET'])
def get_report(id):
    report = AnalysisReport.query.get_or_404(id)
    return app.response_class(
        response=json.dumps(report.to_dict(), ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/research-rounds', methods=['GET'])
def get_research_rounds(id):
    _, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response

    rounds = ResearchRound.query.filter_by(scenario_id=id).order_by(ResearchRound.round_index.desc()).all()
    return app.response_class(
        response=json.dumps([item.to_dict() for item in rounds], ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/report-quality', methods=['GET'])
def get_report_quality(id):
    _, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response

    report = AnalysisReport.query.filter_by(scenario_id=id).first()
    if not report:
        return app.response_class(
            response=json.dumps({"error": "暂无报告，请先生成"}, ensure_ascii=False),
            status=404,
            mimetype='application/json; charset=utf-8'
        )

    data = report.to_dict()
    result = {
        "confidence_level": data.get("confidence_level") or "中",
        "evidence_count": len(data.get("evidence_items") or []),
        "risk_count": len(data.get("decision_risks") or []),
        "action_count": len(data.get("recommended_actions") or []),
        "assumption_count": len(data.get("key_assumptions") or []),
        "has_external_evidence": bool(ExternalEvidence.query.filter_by(scenario_id=id).first())
    }
    return app.response_class(
        response=json.dumps(result, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

def get_report_or_404(scenario_id):
    report = AnalysisReport.query.filter_by(scenario_id=scenario_id).first()
    if report:
        return report, None
    return None, app.response_class(
        response=json.dumps({"error": "暂无报告，请先生成报告"}, ensure_ascii=False),
        status=404,
        mimetype='application/json; charset=utf-8'
    )


def build_report_export_filename(prefix, extension):
    return f"{prefix}_{datetime.now().strftime('%Y%m%d')}.{extension}"


def build_report_export_text(report):
    lines = [
        report.report_title or "分析报告",
        "",
        f"生成时间：{report.generated_at.strftime('%Y-%m-%d %H:%M:%S') if report.generated_at else ''}",
        "",
        "执行摘要",
        report.executive_summary or "暂无执行摘要",
        "",
        "正文",
        report.content or "",
        "",
        "关键假设",
    ]
    lines.extend([f"- {item}" for item in safe_json_loads(report.key_assumptions, [])] or ["- 暂无"])
    lines.extend(["", "风险"])
    lines.extend([f"- {item}" for item in safe_json_loads(report.decision_risks, [])] or ["- 暂无"])
    lines.extend(["", "建议动作"])
    lines.extend([f"- {item}" for item in safe_json_loads(report.recommended_actions, [])] or ["- 暂无"])
    lines.extend(["", "证据"])
    lines.extend([f"- {item}" for item in safe_json_loads(report.evidence_items, [])] or ["- 暂无"])
    lines.extend(["", "来源拆解"])
    source_breakdown = safe_json_loads(report.source_breakdown, [])
    if source_breakdown:
        for item in source_breakdown:
            if isinstance(item, dict):
                lines.append(
                    f"- {item.get('section', '未命名来源')} | {item.get('source_type', 'discussion')} | {item.get('summary', '')}"
                )
            else:
                lines.append(f"- {item}")
    else:
        lines.append("- 暂无")
    return "\n".join(lines).strip()


def build_report_docx_bytes(report):
    doc = Document()
    doc.add_heading(report.report_title or "分析报告", level=0)
    if report.generated_at:
        doc.add_paragraph(f"生成时间：{report.generated_at.strftime('%Y-%m-%d %H:%M:%S')}")
    doc.add_heading("执行摘要", level=1)
    doc.add_paragraph(report.executive_summary or "暂无执行摘要")
    doc.add_heading("报告正文", level=1)
    for line in (report.content or '').splitlines():
        text = line.strip()
        if not text:
            continue
        if text.startswith('## '):
            doc.add_heading(text[3:], level=2)
        elif text.startswith('# '):
            doc.add_heading(text[2:], level=1)
        elif text.startswith('- ') or text.startswith('* '):
            doc.add_paragraph(text[2:], style='List Bullet')
        else:
            doc.add_paragraph(text)

    sections = [
        ("关键假设", safe_json_loads(report.key_assumptions, [])),
        ("风险", safe_json_loads(report.decision_risks, [])),
        ("建议动作", safe_json_loads(report.recommended_actions, [])),
        ("证据", safe_json_loads(report.evidence_items, [])),
    ]
    for title, items in sections:
        doc.add_heading(title, level=1)
        if items:
            for item in items:
                doc.add_paragraph(str(item), style='List Bullet')
        else:
            doc.add_paragraph("暂无")

    doc.add_heading("来源拆解", level=1)
    source_breakdown = safe_json_loads(report.source_breakdown, [])
    if source_breakdown:
        for item in source_breakdown:
            if isinstance(item, dict):
                line = f"{item.get('section', '未命名来源')} | {item.get('source_type', 'discussion')} | {item.get('summary', '')}"
                doc.add_paragraph(line, style='List Bullet')
            else:
                doc.add_paragraph(str(item), style='List Bullet')
    else:
        doc.add_paragraph("暂无")

    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def split_presentation_sentences(text, limit=4):
    parts = [item.strip(" -•\t") for item in re.split(r'[。！？；\n]+', str(text or '')) if item.strip(" -•\t")]
    return parts[:limit]


def split_report_markdown_sections(content):
    sections = []
    current_title = None
    current_lines = []

    for raw_line in str(content or '').splitlines():
        line = raw_line.strip()
        if line.startswith('## '):
            if current_title:
                sections.append({"title": current_title, "lines": current_lines})
            current_title = line[3:].strip()
            current_lines = []
        elif current_title:
            current_lines.append(line)

    if current_title:
        sections.append({"title": current_title, "lines": current_lines})
    return sections


def find_report_section(sections, keyword, fallback_index=None):
    for section in sections:
        if keyword in section.get("title", ""):
            return section
    if fallback_index is not None and len(sections) > fallback_index:
        return sections[fallback_index]
    return {"title": "", "lines": []}


def extract_clean_bullets(lines, limit=6):
    cleaned = []
    for line in lines or []:
        text = str(line or '').strip()
        if not text:
            continue
        text = re.sub(r'^\*\*(.+?)\*\*$', r'\1', text)
        text = re.sub(r'^[-*]\s*', '', text)
        text = re.sub(r'^\d+\.\s*', '', text)
        if text and not re.fullmatch(r'[-|:\s]+', text):
            cleaned.append(text)
        if len(cleaned) >= limit:
            break
    return cleaned


def extract_insight_blocks(section):
    blocks = []
    current_block = None

    for raw_line in section.get("lines", []):
        line = str(raw_line or '').strip()
        if not line:
            continue

        title_match = re.match(r'^\*\*(.+?)\*\*$', line)
        if title_match:
            if current_block:
                blocks.append(current_block)
            current_block = {
                "title": title_match.group(1).strip(),
                "lines": []
            }
            continue

        if current_block is None:
            current_block = {"title": "关键讨论点", "lines": []}
        current_block["lines"].append(line)

    if current_block:
        blocks.append(current_block)

    return blocks


def build_default_presentation_payload(report):
    report_dict = report.to_dict()
    sections = split_report_markdown_sections(report_dict.get("content", ""))
    summary_sentences = split_presentation_sentences(report_dict.get("executive_summary", ""), limit=4)
    overview_section = find_report_section(sections, "核心观点", fallback_index=0)
    discussion_section = find_report_section(sections, "关键讨论点", fallback_index=1)
    divergence_section = find_report_section(sections, "主要分歧", fallback_index=2)
    conclusion_section = find_report_section(sections, "结论", fallback_index=3)

    insight_blocks = extract_insight_blocks(discussion_section)
    if len(insight_blocks) < 4:
        existing_titles = {item["title"] for item in insight_blocks}
        fallback_topics = normalize_text_list(
            report_dict.get("discussion_topics")
            or report_dict.get("key_assumptions")
            or ["需求强度", "噪音控制", "维护复杂度", "定价策略"]
        )
        for topic in fallback_topics:
            normalized_title = f"关键讨论点：{topic}" if "关键讨论点" not in topic else topic
            if normalized_title in existing_titles:
                continue
            insight_blocks.append({
                "title": normalized_title,
                "lines": [
                    f"围绕“{topic}”需要进一步收集真实使用反馈。",
                    "建议结合试点数据判断是否具备持续推进价值。"
                ]
            })
            existing_titles.add(normalized_title)
            if len(insight_blocks) >= 4:
                break

    while len(insight_blocks) < 4:
        index = len(insight_blocks) + 1
        insight_blocks.append({
            "title": f"关键讨论点：补充议题 {index}",
            "lines": ["当前报告未提供足够材料，建议在下一轮研究中补充。"]
        })

    slides = [
        {
            "type": "cover",
            "title": report_dict.get("report_title") or "分析演示稿",
            "subtitle": "基于当前分析报告自动生成的详细汇报稿",
            "bullets": []
        },
        {
            "type": "summary",
            "title": "执行摘要",
            "subtitle": report_dict.get("executive_summary") or "",
            "bullets": summary_sentences or ["当前暂无执行摘要，请先生成完整报告。"]
        },
        {
            "type": "overview",
            "title": overview_section.get("title") or "核心观点总结",
            "subtitle": "先给出管理层可快速判断的总体结论。",
            "bullets": extract_clean_bullets(overview_section.get("lines"), limit=6) or summary_sentences[:3]
        },
    ]

    for block in insight_blocks[:4]:
        title = block.get("title") or "关键讨论点"
        if "关键讨论点" not in title:
            title = f"关键讨论点：{title}"
        slides.append({
            "type": "insight",
            "title": title,
            "subtitle": "",
            "bullets": extract_clean_bullets(block.get("lines"), limit=6) or ["当前议题需要补充更多讨论材料。"]
        })

    slides.extend([
        {
            "type": "divergence",
            "title": divergence_section.get("title") or "主要分歧",
            "subtitle": "列出当前争议点与偏向判断。",
            "bullets": extract_clean_bullets(divergence_section.get("lines"), limit=6)
                or ["当前报告未明确列出主要分歧，需要补充更细的对立观点。"]
        },
        {
            "type": "risks",
            "title": "关键风险",
            "subtitle": "推进前需要重点盯住的风险项。",
            "bullets": normalize_text_list(report_dict.get("decision_risks"))[:6] or ["当前暂无结构化风险数据。"]
        },
        {
            "type": "actions",
            "title": "建议动作",
            "subtitle": "优先执行的验证动作与推进步骤。",
            "bullets": normalize_text_list(report_dict.get("recommended_actions"))[:6] or ["当前暂无结构化行动建议。"]
        },
        {
            "type": "conclusion",
            "title": "结论与下一步",
            "subtitle": "给出是否建议推进，以及对应的停止条件。",
            "bullets": extract_clean_bullets(conclusion_section.get("lines"), limit=6)
                or summary_sentences
                or ["建议结合下一轮试点结果再做最终判断。"]
        },
    ])

    return {
        "scenario_id": report.scenario_id,
        "report_title": report.report_title,
        "generated_at": datetime.now().isoformat(),
        "slide_count": len(slides),
        "slides": slides,
    }


def build_presentation_prompt(report):
    report_dict = report.to_dict()
    payload = {
        "report_title": report_dict.get("report_title", ""),
        "executive_summary": report_dict.get("executive_summary", ""),
        "content": report_dict.get("content", ""),
        "decision_risks": report_dict.get("decision_risks", []),
        "recommended_actions": report_dict.get("recommended_actions", []),
        "evidence_items": report_dict.get("evidence_items", []),
        "source_breakdown": report_dict.get("source_breakdown", []),
    }
    return f"""请基于下面的分析报告，输出一个 JSON 对象用于网页演示稿与 PPT 导出。

要求：
1. 只输出 JSON，不要解释。
2. 顶层结构必须是 {{ "slides": [...] }}。
3. slides 总页数控制在 10-14 页，必须章节单独拆页，不要再压缩成 6 页摘要版。
4. 每页必须包含 type、title，正文使用 bullets 数组，可选 subtitle。
5. type 允许使用：cover、summary、overview、insight、divergence、risks、actions、conclusion。
6. 关键讨论点至少拆成 4 页，每个关键议题单独一页。
7. 必须包含“结论与下一步”这一页，并明确是否建议继续推进、下一步动作和停止条件。
8. 每页 bullet 数量建议 3-6 条，内容要具体，不要只写关键词。

报告数据：
{json.dumps(payload, ensure_ascii=False)}"""


def normalize_presentation_payload(raw_payload, report):
    parsed = safe_json_loads(raw_payload, {}) if isinstance(raw_payload, str) else raw_payload
    if not isinstance(parsed, dict):
        return build_default_presentation_payload(report)
    slides = parsed.get("slides")
    if not isinstance(slides, list):
        return build_default_presentation_payload(report)

    normalized = []
    for index, item in enumerate(slides):
        if not isinstance(item, dict):
            continue
        slide_type = str(item.get("type") or '').strip()
        title = str(item.get("title") or f"第 {index + 1} 页").strip()
        subtitle = str(item.get("subtitle") or '').strip()
        bullets = item.get("bullets") if isinstance(item.get("bullets"), list) else []
        normalized_bullets = [str(bullet).strip() for bullet in bullets if str(bullet).strip()][:6]
        if slide_type and title:
            normalized.append({
                "type": slide_type,
                "title": title,
                "subtitle": subtitle,
                "bullets": normalized_bullets,
            })

    if len(normalized) < 10:
        return build_default_presentation_payload(report)

    return {
        "scenario_id": report.scenario_id,
        "report_title": report.report_title,
        "generated_at": datetime.now().isoformat(),
        "slide_count": len(normalized),
        "slides": normalized,
    }


def build_presentation_pptx_bytes(payload):
    from pptx import Presentation
    from pptx.util import Pt

    presentation = Presentation()
    presentation.core_properties.title = payload.get("report_title") or "分析演示稿"

    for index, slide_data in enumerate(payload.get("slides", [])):
        slide_type = str(slide_data.get("type") or '').strip()
        title = str(slide_data.get("title") or f"第 {index + 1} 页").strip()
        subtitle = str(slide_data.get("subtitle") or '').strip()
        bullets = [str(item).strip() for item in slide_data.get("bullets", []) if str(item).strip()]

        if slide_type == "cover":
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle or "基于分析报告自动生成"
            continue

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        first_written = False
        if subtitle:
            paragraph = text_frame.paragraphs[0]
            paragraph.text = subtitle
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            first_written = True

        for bullet in bullets:
            paragraph = text_frame.add_paragraph() if first_written else text_frame.paragraphs[0]
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(22)
            first_written = True

        if not first_written:
            text_frame.paragraphs[0].text = "暂无内容"
            text_frame.paragraphs[0].font.size = Pt(20)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def save_presentation_payload(scenario_id, payload):
    serialized = json.dumps(payload, ensure_ascii=False)
    with db.engine.begin() as connection:
        connection.exec_driver_sql(
            "UPDATE analysis_report SET presentation_payload = ?, presentation_generated_at = ? WHERE scenario_id = ?",
            (serialized, datetime.now().isoformat(sep=' '), scenario_id)
        )


def load_presentation_payload(scenario_id):
    with db.engine.begin() as connection:
        row = connection.exec_driver_sql(
            "SELECT presentation_payload FROM analysis_report WHERE scenario_id = ?",
            (scenario_id,)
        ).fetchone()
    if not row or not row[0]:
        return None
    return safe_json_loads(row[0], None)


@app.route('/api/scenarios/<int:id>/report/pdf', methods=['GET'])
def get_report_pdf(id):
    _, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response
    
    report = AnalysisReport.query.filter_by(scenario_id=id).first()
    if not report:
        return app.response_class(
            response=json.dumps({"error": "暂无报告，请先生成"}, ensure_ascii=False),
            status=404,
            mimetype='application/json; charset=utf-8'
        )

    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import inch
    from reportlab.lib.enums import TA_CENTER

    result = BytesIO()
    doc = SimpleDocTemplate(result, pagesize=A4)
    styles = getSampleStyleSheet()
    
    try:
        if platform.system() == 'Windows':
            font_paths = [
                ('MicrosoftYaHei', 'C:\\Windows\\Fonts\\msyh.ttc'),
                ('SimSun', 'C:\\Windows\\Fonts\\simsun.ttc'),
                ('SimHei', 'C:\\Windows\\Fonts\\simhei.ttf')
            ]
            font_name = None
            for name, path in font_paths:
                if os.path.exists(path):
                    try:
                        pdfmetrics.registerFont(TTFont(name, path))
                        font_name = name
                        break
                    except:
                        continue
        else:
            font_name = 'Helvetica'
    except:
        font_name = 'Helvetica'
    
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontName=font_name,
        fontSize=18,
        alignment=TA_CENTER,
        spaceAfter=30
    )
    
    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=11,
        leading=18
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontName=font_name,
        fontSize=14,
        spaceBefore=15,
        spaceAfter=8
    )
    
    story = []
    
    story.append(Paragraph(report.report_title, title_style))
    
    meta_text = f"报告生成时间：{report.generated_at.strftime('%Y年%m月%d日 %H:%M:%S')}"
    meta_style = ParagraphStyle(
        'MetaStyle',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=10,
        textColor='#666666',
        alignment=TA_CENTER
    )
    story.append(Paragraph(meta_text, meta_style))
    story.append(Spacer(1, 20))
    
    lines = report.content.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith('# '):
            story.append(Paragraph(line[2:], heading_style))
        elif line.startswith('## '):
            story.append(Paragraph(line[3:], heading_style))
        elif line.startswith('### '):
            story.append(Paragraph(line[4:], heading_style))
        elif line.startswith('- ') or line.startswith('* '):
            story.append(Paragraph('• ' + line[2:], normal_style))
        elif line.startswith('**') and line.endswith('**'):
            story.append(Paragraph(f'<b>{line[2:-2]}</b>', normal_style))
        else:
            story.append(Paragraph(line, normal_style))
        story.append(Spacer(1, 6))
    
    story.append(Spacer(1, 30))
    footer_style = ParagraphStyle(
        'FooterStyle',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=9,
        textColor='#999999',
        alignment=TA_CENTER
    )
    story.append(Paragraph('焦点小组讨论系统 - 分析报告', footer_style))
    
    doc.build(story)
    pdf_bytes = result.getvalue()

    response = make_response(pdf_bytes)
    response.headers['Content-Type'] = 'application/pdf'
    filename = f"analysis_report_{datetime.now().strftime('%Y%m%d')}.pdf"
    response.headers['Content-Disposition'] = f'attachment; filename="{quote(filename)}"'
    
    return response


@app.route('/api/scenarios/<int:id>/report/txt', methods=['GET'])
def get_report_txt(id):
    _, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response

    report, report_error = get_report_or_404(id)
    if report_error:
        return report_error

    payload = build_report_export_text(report).encode('utf-8')
    response = make_response(payload)
    response.headers['Content-Type'] = 'text/plain; charset=utf-8'
    filename = build_report_export_filename('analysis_report', 'txt')
    response.headers['Content-Disposition'] = f'attachment; filename="{quote(filename)}"'
    return response


@app.route('/api/scenarios/<int:id>/report/docx', methods=['GET'])
def get_report_docx(id):
    _, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response

    report, report_error = get_report_or_404(id)
    if report_error:
        return report_error

    docx_bytes = build_report_docx_bytes(report)
    response = make_response(docx_bytes)
    response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    filename = build_report_export_filename('analysis_report', 'docx')
    response.headers['Content-Disposition'] = f'attachment; filename="{quote(filename)}"'
    return response


@app.route('/api/scenarios/<int:id>/presentation', methods=['GET'])
def get_report_presentation(id):
    _, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response

    report, report_error = get_report_or_404(id)
    if report_error:
        return report_error

    payload = load_presentation_payload(id)
    if not payload:
        return app.response_class(
            response=json.dumps({"error": "暂无演示稿，请先生成演示稿"}, ensure_ascii=False),
            status=404,
            mimetype='application/json; charset=utf-8'
        )

    result = {
        "scenario_id": id,
        "report_title": report.report_title,
        "slide_count": len(payload.get("slides", [])),
        "slides": payload.get("slides", []),
        "generated_at": payload.get("generated_at", ""),
    }
    return app.response_class(
        response=json.dumps(result, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/meeting-rooms', methods=['GET'])
def list_meeting_rooms():
    user = get_current_user_model()
    if not user:
        return json_api_error("请先登录", 401)

    joined_room_ids = [
        item[0]
        for item in db.session.query(MeetingRoomMember.room_id)
        .filter_by(user_id=user.id)
        .all()
    ]
    room_query = MeetingRoom.query.filter(
        or_(
            MeetingRoom.owner_user_id == user.id,
            MeetingRoom.id.in_(joined_room_ids or [-1]),
        )
    ).order_by(MeetingRoom.updated_at.desc(), MeetingRoom.created_at.desc(), MeetingRoom.id.desc())
    rooms = room_query.all()
    return json_api_response({
        'rooms': [serialize_meeting_room_bundle(room) for room in rooms]
    })


@app.route('/api/meeting-rooms', methods=['POST'])
def create_meeting_room():
    user = get_current_user_model()
    if not user:
        return json_api_error("请先登录", 401)

    data = request.get_json() or {}
    room_name = str(data.get('room_name') or '').strip()
    topic_title = str(data.get('topic_title') or '').strip()
    product_name = str(data.get('product_name') or '').strip()
    product_concept = str(data.get('product_concept') or '').strip()

    if not room_name:
        return json_api_error("请填写会议室名称", 400)
    if not topic_title:
        return json_api_error("请填写会议主题", 400)
    if not product_name:
        return json_api_error("请填写产品名称", 400)
    if not product_concept:
        return json_api_error("请填写产品概念", 400)

    try:
        target_count = max(1, int(data.get('target_count', 4)))
    except (TypeError, ValueError):
        target_count = 4

    scenario = build_meeting_scenario_from_payload(user.id, data)
    db.session.add(scenario)
    db.session.flush()

    room = MeetingRoom(
        scenario_id=scenario.id,
        owner_user_id=user.id,
        room_name=room_name,
        room_code=generate_meeting_room_code(),
        invite_token=generate_meeting_invite_token(),
        status='waiting',
        topic_title=topic_title,
        topic_notes=str(data.get('topic_notes') or '').strip(),
        target_count=target_count,
        created_at=datetime.now(),
        updated_at=datetime.now()
    )
    db.session.add(room)
    db.session.flush()

    upsert_room_member(room, user, member_role='owner')
    db.session.commit()
    return json_api_response(serialize_meeting_room_bundle(room, scenario), status=201)


@app.route('/api/meeting-rooms/join', methods=['POST'])
def join_meeting_room():
    user = get_current_user_model()
    if not user:
        return json_api_error("请先登录", 401)

    data = request.get_json() or {}
    room_code = str(data.get('room_code') or '').strip().upper()
    invite_token = str(data.get('invite_token') or '').strip()

    room = None
    if room_code:
        room = MeetingRoom.query.filter_by(room_code=room_code).first()
    elif invite_token:
        room = MeetingRoom.query.filter_by(invite_token=invite_token).first()

    if not room:
        return json_api_error("会议室不存在或邀请码无效", 404)

    existing_member = MeetingRoomMember.query.filter_by(room_id=room.id, user_id=user.id).first()
    room_status = room.status or 'waiting'
    if room_status != 'waiting' and not existing_member and room.owner_user_id != user.id:
        if room_status == 'ended':
            return json_api_error("会议已结束，不支持新成员加入", 409)
        return json_api_error("会议已经开始，暂不支持新成员加入", 409)

    upsert_room_member(room, user, member_role='member')
    room.updated_at = datetime.now()
    db.session.commit()
    return json_api_response(serialize_meeting_room_bundle(room), status=200)


@app.route('/api/meeting-rooms/<int:room_id>', methods=['GET'])
def get_meeting_room(room_id):
    room, _, error_response = get_meeting_room_or_error(room_id)
    if error_response:
        return error_response
    scenario = db.session.get(ProductScenario, room.scenario_id)
    if scenario:
        ensure_current_human_turn(room, scenario)
        db.session.commit()
    return json_api_response(serialize_meeting_room_bundle(room, scenario, include_messages=True))


@app.route('/api/meeting-rooms/<int:room_id>', methods=['PATCH'])
def update_meeting_room(room_id):
    room, membership, error_response = get_meeting_room_or_error(room_id)
    if error_response:
        return error_response
    if room.owner_user_id != session.get('user_id'):
        return json_api_error("只有房主可以修改会议主题", 403)
    if (room.status or 'waiting') != 'waiting':
        return json_api_error("会议开始后不能再修改主题", 400)

    data = request.get_json() or {}
    scenario = db.session.get(ProductScenario, room.scenario_id)
    if not scenario:
        return json_api_error("会议绑定的场景不存在", 404)

    topic_title = str(data.get('topic_title') or room.topic_title or '').strip()
    topic_notes = str(data.get('topic_notes') or '').strip()
    discussion_topics = normalize_text_list(data.get('discussion_topics', scenario.to_dict().get('discussion_topics', [])))

    room.topic_title = topic_title or room.topic_title
    room.topic_notes = topic_notes
    room.updated_at = datetime.now()

    scenario.decision_problem = room.topic_title
    if discussion_topics:
        scenario.discussion_topics = json_dumps(discussion_topics, [])
    scenario.research_plan = build_research_plan(scenario)

    create_meeting_room_system_message(room.id, f"房主已更新主题：{room.topic_title}")
    if membership:
        membership.last_seen_at = datetime.now()
    db.session.commit()
    return json_api_response(
        serialize_meeting_room_bundle(room, scenario, include_messages=True),
        status=200
    )


@app.route('/api/meeting-rooms/<int:room_id>/turn-order', methods=['PATCH'])
def update_meeting_room_turn_order(room_id):
    room, membership, error_response = get_meeting_room_or_error(room_id)
    if error_response:
        return error_response
    if room.owner_user_id != session.get('user_id'):
        return json_api_error("只有房主可以调整发言顺序", 403)
    if (room.status or 'waiting') != 'waiting':
        return json_api_error("会议开始后不能再调整发言顺序", 400)

    data = request.get_json() or {}
    member_ids = data.get('member_ids', [])
    if not isinstance(member_ids, list) or not member_ids:
        return json_api_error("请提交完整的成员顺序", 400)

    members = get_meeting_room_members(room)
    valid_ids = [member.id for member in members]
    normalized_ids = []
    seen = set()
    for member_id in member_ids:
        try:
            normalized_id = int(member_id)
        except (TypeError, ValueError):
            return json_api_error("成员顺序格式无效", 400)
        if normalized_id not in valid_ids or normalized_id in seen:
            return json_api_error("成员顺序不完整或包含无效成员", 400)
        normalized_ids.append(normalized_id)
        seen.add(normalized_id)

    if set(normalized_ids) != set(valid_ids):
        return json_api_error("成员顺序必须覆盖全部当前成员", 400)

    room.turn_order = json_dumps(normalized_ids, [])
    room.updated_at = datetime.now()
    if membership:
        membership.last_seen_at = datetime.now()
    db.session.commit()
    return json_api_response(serialize_meeting_room_bundle(room, include_messages=True), status=200)


@app.route('/api/meeting-rooms/<int:room_id>/messages', methods=['GET'])
def get_meeting_room_messages(room_id):
    room, membership, error_response = get_meeting_room_or_error(room_id)
    if error_response:
        return error_response
    if membership:
        membership.last_seen_at = datetime.now()
        db.session.commit()
    messages = MeetingRoomMessage.query.filter_by(room_id=room.id).order_by(MeetingRoomMessage.created_at.asc()).all()
    return json_api_response({'messages': [message.to_dict() for message in messages]})


@app.route('/api/meeting-rooms/<int:room_id>/messages', methods=['POST'])
def post_meeting_room_message(room_id):
    room, membership, error_response = get_meeting_room_or_error(room_id)
    if error_response:
        return error_response

    data = request.get_json() or {}
    content = str(data.get('content') or '').strip()
    if not content:
        return json_api_error("消息内容不能为空", 400)

    user = get_current_user_model()
    message = MeetingRoomMessage(
        room_id=room.id,
        user_id=user.id if user else None,
        sender_name=membership.display_name if membership and membership.display_name else get_room_member_display_name(user),
        content=content,
        message_type='member',
        created_at=datetime.now()
    )
    db.session.add(message)
    if membership:
        membership.last_seen_at = datetime.now()
    room.updated_at = datetime.now()
    db.session.commit()
    scenario = db.session.get(ProductScenario, room.scenario_id)
    return json_api_response(
        serialize_meeting_room_bundle(room, scenario, include_messages=True),
        status=201
    )


@app.route('/api/meeting-rooms/<int:room_id>/start', methods=['POST'])
def start_meeting_room(room_id):
    room, membership, error_response = get_meeting_room_or_error(room_id)
    if error_response:
        return error_response
    if room.owner_user_id != session.get('user_id'):
        return json_api_error("只有房主可以开始会议", 403)
    if (room.status or 'waiting') == 'ended':
        return json_api_error("会议已结束，不能再次开始", 400)
    if (room.status or 'waiting') == 'active' or (room.discussion_phase or 'waiting') == 'live':
        return json_api_error("会议已经开始", 400)

    scenario = db.session.get(ProductScenario, room.scenario_id)
    if not scenario:
        return json_api_error("会议绑定的场景不存在", 404)

    current_count = MeetingRoomMember.query.filter_by(room_id=room.id).count()
    target_count = max(room.target_count or 0, 0)
    gap_count = max(target_count - current_count, 0)
    generated_participants = []
    if gap_count > 0:
        generated_participants = generate_participants_for_room_gap(scenario, gap_count)
        create_meeting_room_system_message(room.id, f"当前人数不足，系统已补入 {gap_count} 位 AI 参与者。")

    initialize_room_discussion(room, scenario, generated_participants)
    create_meeting_room_system_message(room.id, "会议已开始，系统将按既定顺序轮流发言。")

    now = datetime.now()
    room.started_at = now
    room.updated_at = now
    scenario.meeting_status = 'meeting_active'
    scenario.research_plan = build_research_plan(scenario)
    if membership:
        membership.last_seen_at = now
    turn_state = ensure_current_human_turn(room, scenario, async_ai=True)
    db.session.commit()

    payload = serialize_meeting_room_bundle(room, scenario, include_messages=True)
    payload['generated_participants'] = generated_participants
    payload['next_url'] = f"/meeting-room?room_id={room.id}"
    if turn_state == 'ai_pending':
        enqueue_room_ai_turn_worker(room.id)
    return json_api_response(payload, status=200)

@app.route('/api/meeting-rooms/<int:room_id>/discussion/messages', methods=['POST'])
def post_meeting_room_discussion_message(room_id):
    room, membership, error_response = get_meeting_room_or_error(room_id)
    if error_response:
        return error_response

    scenario = db.session.get(ProductScenario, room.scenario_id)
    if not scenario:
        return json_api_error("会议绑定的场景不存在", 404)
    if (room.discussion_phase or 'waiting') != 'live' or (room.status or 'waiting') == 'ended':
        return json_api_error("会议尚未开始或已经结束", 400)

    ensure_current_human_turn(room, scenario)
    active_participant = db.session.get(VirtualParticipant, room.active_speaker_participant_id) if room.active_speaker_participant_id else None
    if not active_participant:
        return json_api_error("当前没有可发言席位", 400)
    if (active_participant.speaker_origin or 'ai') != 'human':
        return json_api_error("当前轮到 AI 席位，系统会自动推进", 400)
    if active_participant.linked_user_id != session.get('user_id'):
        return json_api_error("当前还没有轮到你发言", 403)

    data = request.get_json() or {}
    content = str(data.get('content') or '').strip()
    if not content:
        return json_api_error("发言内容不能为空", 400)

    participants = get_room_managed_participants(room.scenario_id)
    participant_count = max(len(participants), 1)
    round_index = max((room.turn_number or 1) - 1, 0) // participant_count

    record = ConversationRecord(
        scenario_id=scenario.id,
        participant_id=active_participant.id,
        content=content,
        is_host=False,
        message_type='turn'
    )
    db.session.add(record)
    update_participant_state_v2(active_participant, content, round_index)

    now = datetime.now()
    room.updated_at = now
    if membership:
        membership.last_seen_at = now

    turn_state = continue_room_discussion_after_turn(room, scenario, participants=participants, async_ai=True)
    db.session.commit()
    payload = serialize_meeting_room_bundle(room, scenario, include_messages=True)
    if turn_state == 'ai_pending':
        enqueue_room_ai_turn_worker(room.id)
    return json_api_response(payload, status=200)


@app.route('/api/meeting-rooms/<int:room_id>/discussion/skip', methods=['POST'])
def skip_meeting_room_discussion_turn(room_id):
    room, membership, error_response = get_meeting_room_or_error(room_id)
    if error_response:
        return error_response
    if room.owner_user_id != session.get('user_id'):
        return json_api_error("只有房主可以跳过当前发言人", 403)

    scenario = db.session.get(ProductScenario, room.scenario_id)
    if not scenario:
        return json_api_error("会议绑定的场景不存在", 404)
    if (room.discussion_phase or 'waiting') != 'live' or (room.status or 'waiting') == 'ended':
        return json_api_error("会议尚未开始或已经结束", 400)

    ensure_current_human_turn(room, scenario)
    active_participant = db.session.get(VirtualParticipant, room.active_speaker_participant_id) if room.active_speaker_participant_id else None
    if not active_participant:
        return json_api_error("当前没有可跳过的发言人", 400)
    if (active_participant.speaker_origin or 'ai') != 'human':
        return json_api_error("当前席位不是人工成员", 400)

    create_meeting_room_system_message(room.id, f"房主已跳过当前发言人：{active_participant.persona_name}")
    now = datetime.now()
    room.updated_at = now
    if membership:
        membership.last_seen_at = now
    participants = get_room_managed_participants(room.scenario_id)
    turn_state = continue_room_discussion_after_turn(room, scenario, participants=participants, async_ai=True)
    db.session.commit()
    payload = serialize_meeting_room_bundle(room, scenario, include_messages=True)
    if turn_state == 'ai_pending':
        enqueue_room_ai_turn_worker(room.id)
    return json_api_response(payload, status=200)


@app.route('/api/meeting-rooms/<int:room_id>/discussion/end', methods=['POST'])
def end_meeting_room_discussion(room_id):
    room, membership, error_response = get_meeting_room_or_error(room_id)
    if error_response:
        return error_response
    if room.owner_user_id != session.get('user_id'):
        return json_api_error("只有房主可以结束会议", 403)

    scenario = db.session.get(ProductScenario, room.scenario_id)
    if not scenario:
        return json_api_error("会议绑定的场景不存在", 404)

    now = datetime.now()
    room.status = 'ended'
    room.discussion_phase = 'ended'
    room.active_speaker_participant_id = None
    room.updated_at = now
    scenario.meeting_status = 'meeting_ended'
    if membership:
        membership.last_seen_at = now
    create_meeting_room_system_message(room.id, "房主已结束本次会议，正式讨论已关闭。")
    db.session.commit()

    payload = serialize_meeting_room_bundle(room, scenario, include_messages=True)
    payload['next_url'] = f"/report?scenario_id={scenario.id}"
    return json_api_response(payload, status=200)


@app.route('/api/scenarios/<int:id>/presentation/pptx', methods=['GET'])
def export_report_presentation_pptx(id):
    _, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response

    _, report_error = get_report_or_404(id)
    if report_error:
        return report_error

    payload = load_presentation_payload(id)
    if not payload:
        return app.response_class(
            response=json.dumps({"error": "暂无演示稿，请先生成演示稿"}, ensure_ascii=False),
            status=404,
            mimetype='application/json; charset=utf-8'
        )

    try:
        pptx_bytes = build_presentation_pptx_bytes(payload)
    except ModuleNotFoundError:
        return app.response_class(
            response=json.dumps({"error": "服务器缺少 PPTX 导出依赖，请安装 python-pptx"}, ensure_ascii=False),
            status=500,
            mimetype='application/json; charset=utf-8'
        )

    response = make_response(pptx_bytes)
    response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    filename = build_report_export_filename('analysis_presentation', 'pptx')
    response.headers['Content-Disposition'] = f'attachment; filename="{quote(filename)}"'
    return response


@app.route('/api/scenarios/<int:id>/presentation/generate', methods=['POST'])
def generate_report_presentation(id):
    _, error_response = get_scenario_or_403(id, allow_meeting_member=True)
    if error_response:
        return error_response

    report, report_error = get_report_or_404(id)
    if report_error:
        return report_error

    prompt = build_presentation_prompt(report)
    result = call_deepseek(
        prompt,
        "你是一位企业汇报顾问，擅长把分析报告整理成适合网页演示稿和 PPT 汇报的详细结构化内容。"
    )
    if is_ai_error_text(result):
        return app.response_class(
            response=json.dumps({"error": result}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    normalized = normalize_presentation_payload(result, report)
    if not normalized:
        return app.response_class(
            response=json.dumps({"error": "演示稿生成失败，请稍后重试"}, ensure_ascii=False),
            status=502,
            mimetype='application/json; charset=utf-8'
        )

    save_presentation_payload(id, normalized)
    return app.response_class(
        response=json.dumps({"status": "success", "presentation": normalized}, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

def format_report_content(content):
    if not content:
        return '<p>报告内容为空</p>'
    
    import re
    text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', content)
    
    lines = text.split('\n')
    result = []
    in_table = False
    table_rows = []
    
    for line in lines:
        line = line.strip()
        if not line:
            if in_table:
                result.append('</tbody></table>')
                result.append('<br>')
                in_table = False
                table_rows = []
            continue
        
        heading_match = re.match(r'^(#{1,6})\s+(.+)$', line)
        if heading_match:
            if in_table:
                result.append('</tbody></table>')
                result.append('<br>')
                in_table = False
                table_rows = []
            level = len(heading_match.group(1))
            result.append(f'<h{level}>{heading_match.group(2)}</h{level}>')
            continue
        
        if re.match(r'^\|.+\|$', line):
            cells = [c.strip() for c in line.split('|') if c.strip() and not re.match(r'^-+$', c)]
            if len(cells) > 0:
                if not in_table:
                    in_table = True
                    table_rows = []
                    result.append('<table class="report-table">')
                
                is_separator = re.match(r'^\|[-:\s]+\|[-:\s]+\|$', line)
                if is_separator:
                    continue
                
                table_rows.append(cells)
                
                if len(table_rows) == 1:
                    result.append('<thead><tr>')
                    for cell in cells:
                        result.append(f'<th>{cell}</th>')
                    result.append('</tr></thead><tbody>')
                else:
                    result.append('<tr>')
                    for cell in cells:
                        result.append(f'<td>{cell}</td>')
                    result.append('</tr>')
            continue
        else:
            if in_table:
                result.append('</tbody></table>')
                result.append('<br>')
                in_table = False
                table_rows = []
        
        if line.startswith('- ') or line.startswith('* '):
            result.append(f'<li>{line[2:]}</li>')
        elif line == '---':
            result.append('<hr>')
        else:
            result.append(f'<p>{line}</p>')
    
    if in_table:
        result.append('</tbody></table>')
    
    html = '\n'.join(result)
    
    if '<thead><tbody>' in html:
        html = html.replace('<thead><tbody>', '<thead></thead><tbody>')
    if html.endswith('</tbody></table>'):
        html = html.replace('<table class="report-table"></thead></tbody></table>', '<table class="report-table"></table>')
    
    return html

@app.route('/api/participants', methods=['GET'])
def get_participants():
    participants = VirtualParticipant.query.all()
    return app.response_class(
        response=json.dumps([p.to_dict() for p in participants], ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/participants', methods=['GET'])
def get_scenario_participants(id):
    scenario, error_response = get_scenario_or_403(id)
    if error_response:
        return error_response
    standalone_error = reject_if_not_standalone(scenario, "会议室内部场景不能使用普通参与者接口，请前往会议室系统继续。")
    if standalone_error:
        return standalone_error
    participants = get_standalone_participants(id)
    return app.response_class(
        response=json.dumps([p.to_dict() for p in participants], ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>/participants', methods=['POST'])
def create_participant(id):
    """
    创建自定义参与者
    
    POST /api/scenarios/{id}/participants
    请求体：{ persona_name, persona_tags, personality, background, usage_goal, budget_sensitivity, brand_preference, risk_aversion, decision_style, deal_breakers, stance_summary }
    返回：创建的参与者对象
    """
    scenario, error_response = get_scenario_or_403(id)
    if error_response:
        return error_response
    standalone_error = reject_if_not_standalone(scenario, "会议室内部场景不能使用普通参与者接口，请前往会议室系统继续。")
    if standalone_error:
        return standalone_error
    
    data = request.get_json() or {}
    payload, error_message = parse_participant_payload(data, require_name=True)
    if error_message:
        return app.response_class(
            response=json.dumps({"error": error_message}, ensure_ascii=False),
            status=400,
            mimetype='application/json; charset=utf-8'
        )

    participant = VirtualParticipant(
        scenario_id=id,
        persona_name=payload['persona_name'],
        persona_tags=json.dumps(payload['persona_tags'], ensure_ascii=False),
        personality=payload['personality'],
        background=payload['background'],
        usage_goal=payload['usage_goal'],
        budget_sensitivity=payload['budget_sensitivity'],
        brand_preference=payload['brand_preference'],
        risk_aversion=payload['risk_aversion'],
        decision_style=payload['decision_style'],
        deal_breakers=json_dumps(payload['deal_breakers'], []),
        stance_summary=payload['stance_summary'],
        stance_state=build_participant_stance_state(payload['stance_summary']),
        is_custom=True,
        is_ai_generated=False,
        custom_params=json.dumps(payload['custom_params'], ensure_ascii=False)
    )
    db.session.add(participant)
    db.session.commit()
    return app.response_class(
        response=json.dumps({"status": "success", "participant": participant.to_dict()}, ensure_ascii=False),
        status=201,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/participants/<int:id>', methods=['PATCH'])
def update_participant(id):
    """
    更新参与者信息
    
    PATCH /api/participants/{id}
    请求体：{ persona_name, persona_tags, personality, background, usage_goal, budget_sensitivity, brand_preference, risk_aversion, decision_style, deal_breakers, stance_summary }
    返回：更新后的参与者对象
    """
    participant = VirtualParticipant.query.get_or_404(id)
    if participant.room_managed:
        return json_api_error("会议室正式席位不能通过普通参与者接口修改。", 409)
    
    # 检查权限
    scenario, error_response = get_scenario_or_403(participant.scenario_id)
    if error_response:
        return error_response

    data = request.get_json() or {}
    payload, error_message = parse_participant_payload(
        data,
        require_name='persona_name' in data,
        existing_custom_params=safe_json_loads(participant.custom_params, {})
    )
    if error_message:
        return app.response_class(
            response=json.dumps({"error": error_message}, ensure_ascii=False),
            status=400,
            mimetype='application/json; charset=utf-8'
        )

    if 'persona_name' in data:
        participant.persona_name = payload['persona_name']
    if 'persona_tags' in data or 'tags' in data:
        participant.persona_tags = json.dumps(payload['persona_tags'], ensure_ascii=False)
    if 'personality' in data:
        participant.personality = payload['personality']
    if 'background' in data:
        participant.background = payload['background']
    if 'usage_goal' in data:
        participant.usage_goal = payload['usage_goal']
    if 'budget_sensitivity' in data:
        participant.budget_sensitivity = payload['budget_sensitivity']
    if 'brand_preference' in data:
        participant.brand_preference = payload['brand_preference']
    if 'risk_aversion' in data:
        participant.risk_aversion = payload['risk_aversion']
    if 'decision_style' in data:
        participant.decision_style = payload['decision_style']
    if 'deal_breakers' in data:
        participant.deal_breakers = json_dumps(payload['deal_breakers'], [])
    if 'stance_summary' in data:
        participant.stance_summary = payload['stance_summary']
        state = safe_json_loads(participant.stance_state, {})
        state['current_position'] = payload['stance_summary']
        participant.stance_state = json.dumps(state, ensure_ascii=False)
    if 'custom_params' in data:
        participant.custom_params = json.dumps(payload['custom_params'], ensure_ascii=False)

    db.session.commit()
    return app.response_class(
        response=json.dumps({"status": "success", "participant": participant.to_dict()}, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/participants/<int:id>', methods=['DELETE'])
def delete_participant(id):
    participant = VirtualParticipant.query.get_or_404(id)
    if participant.room_managed:
        return json_api_error("会议室正式席位不能通过普通参与者接口删除。", 409)
    scenario_id = participant.scenario_id
    
    # 检查权限
    scenario, error_response = get_scenario_or_403(scenario_id)
    if error_response:
        return error_response
    
    db.session.delete(participant)
    db.session.commit()
    return app.response_class(
        response=json.dumps({"status": "success", "message": f"参与者 {id} 已删除", "scenario_id": scenario_id}, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

@app.route('/api/scenarios/<int:id>', methods=['DELETE'])
def delete_scenario(id):
    scenario, error_response = get_scenario_or_403(id)
    if error_response:
        return error_response
    standalone_error = reject_if_not_standalone(scenario, "会议室内部场景不能在普通产品讨论中删除，请前往会议室系统处理。")
    if standalone_error:
        return standalone_error

    ConversationRecord.query.filter_by(scenario_id=id).delete()
    VirtualParticipant.query.filter(VirtualParticipant.scenario_id == id).filter(
        or_(VirtualParticipant.room_managed.is_(False), VirtualParticipant.room_managed.is_(None))
    ).delete()
    AnalysisReport.query.filter_by(scenario_id=id).delete()
    db.session.delete(scenario)
    db.session.commit()

    return app.response_class(
        response=json.dumps({"status": "deleted", "message": f"场景 {id} 已删除"}, ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

# ---------------------------
# 用户认证API
# ---------------------------

@app.route('/api/register', methods=['POST'])
def register():
    """
    用户注册
    
    POST /api/register
    请求体：{ username, email, phone, code, password }
    返回：用户对象
    
    功能说明：
        - 创建新用户账户
        - 验证验证码
        - 验证用户名和邮箱唯一性
        - 使用PBKDF2加密存储密码
        - 注册成功后自动登录
    """
    data = request.json or {}
    username = (data.get('username') or '').strip()
    email = normalize_auth_target(data.get('email'))
    password = data.get('password', '')

    user, error_message, status = register_user_account(username, email, password)
    if error_message:
        return json_api_error(error_message, status)

    apply_login_session(user, update_stats=False)
    return json_api_response({"status": "success", "user": user.to_dict()}, status=201)

@app.route('/api/login', methods=['POST'])
def login():
    """
    用户登录
    
    POST /api/login
    请求体：{ target, password }
    返回：用户对象
    
    功能说明：
        - 通过邮箱或手机号登录
        - 使用PBKDF2验证密码（兼容旧SHA256格式）
        - 设置登录会话
        - 登录速率限制
    """
    data = request.json or {}
    email = normalize_auth_target(data.get('email') or data.get('target'))
    password = data.get('password', '').strip()

    if not email or not password or not validate_email(email):
        return json_api_error("邮箱或密码错误", 401)

    # 检查登录速率限制
    ip = _get_client_ip()
    if not check_login_rate_limit(f"{email}:{ip}"):
        return json_api_error("登录过于频繁，请稍后再试", 429)

    # 当前版本仅支持邮箱密码登录
    user = find_user_by_email(email)
    
    if not user or not verify_password(password, user.password_hash):
        return json_api_error("邮箱或密码错误", 401)

    apply_login_session(user)
    return json_api_response({"status": "success", "user": user.to_dict()})

@app.route('/api/logout', methods=['POST'])
def logout():
    """
    用户登出
    
    POST /api/logout
    返回：成功消息
    
    功能说明：清除用户登录会话
    """
    session.clear()
    return json_api_response({"message": "已退出登录"})

@app.route('/api/user', methods=['GET', 'PATCH'])
def get_user():
    """
    获取或更新当前登录用户信息
    
    GET /api/user
    返回：用户对象
    
    PATCH /api/user
    请求体：{ nickname, avatar, company, bio, locale }
    返回：更新后的用户对象
    
    功能说明：验证登录状态并返回/更新用户信息
    """
    user_id = session.get('user_id')
    if not user_id:
        return json_api_error("未登录", 401)

    user = db.session.get(User, user_id)
    if not user:
        return json_api_error("无效的登录状态", 401)

    if request.method == 'GET':
        return json_api_response(user.to_dict())
    else:  # PATCH
        data = request.get_json() or {}
        if 'nickname' in data:
            user.nickname = str(data['nickname']).strip()[:80]
        if 'avatar' in data:
            user.avatar = str(data['avatar']).strip()[:200]
        if 'company' in data:
            user.company = str(data['company']).strip()[:120]
        if 'bio' in data:
            user.bio = str(data['bio']).strip()
        if 'locale' in data:
            user.locale = str(data['locale']).strip()[:10]
        db.session.commit()
        return json_api_response({"status": "success", "user": user.to_dict()})

@app.route('/api/user/ai-config', methods=['GET', 'PATCH'])
def user_ai_config():
    user_id = session.get('user_id')
    if not user_id:
        return json_api_error("未登录", 401)

    user = db.session.get(User, user_id)
    if not user:
        return json_api_error("无效的登录状态", 401)

    if request.method == 'GET':
        return json_api_response(get_user_ai_config_summary(user))

    data = request.get_json() or {}
    success, payload, valid_keys = test_user_ai_config(
        data.get('endpoint_url'),
        data.get('model_name'),
        data.get('api_keys')
    )
    user.ai_last_test_status = 'success' if success else 'error'
    user.ai_last_test_message = payload.get('message', '')
    user.ai_last_tested_at = datetime.now()

    if not success:
        db.session.commit()
        return json_api_error(payload.get('message', '用户 AI 配置不可用'), 400)

    endpoint_url, model_name, _, _ = validate_user_ai_config_payload(
        data.get('endpoint_url'),
        data.get('model_name'),
        valid_keys
    )
    user.ai_endpoint_url = endpoint_url
    user.ai_model_name = model_name
    user.ai_api_keys_encrypted = encrypt_user_api_keys(valid_keys, user=user)
    user.ai_config_enabled = True
    db.session.commit()
    return json_api_response({
        "status": "success",
        "message": payload.get('message', '用户 AI 配置已保存'),
        "config": get_user_ai_config_summary(user)
    })


@app.route('/api/user/ai-config/test', methods=['POST'])
def user_ai_config_test():
    user_id = session.get('user_id')
    if not user_id:
        return json_api_error("未登录", 401)

    user = db.session.get(User, user_id)
    if not user:
        return json_api_error("无效的登录状态", 401)

    data = request.get_json() or {}
    success, payload, _ = test_user_ai_config(
        data.get('endpoint_url'),
        data.get('model_name'),
        data.get('api_keys')
    )
    return json_api_response(payload, status=200 if success else 400)


@app.route('/api/user/password', methods=['POST'])
def change_password():
    """
    修改密码
    
    POST /api/user/password
    请求体：{ old_password, new_password }
    返回：成功消息
    
    功能说明：验证旧密码后设置新密码
    """
    user_id = session.get('user_id')
    if not user_id:
        return json_api_error("未登录", 401)

    user = db.session.get(User, user_id)
    if not user:
        return json_api_error("无效的登录状态", 401)

    data = request.get_json() or {}
    old_pwd = data.get('old_password', '')
    new_pwd = data.get('new_password', '')
    code = data.get('code')

    if not verify_password(old_pwd, user.password_hash):
        return json_api_error("旧密码不正确", 403)

    if not code:
        return json_api_error("请填写邮箱验证码", 400)

    verified, message = verify_verification_code(user.email, code, 'change_password')
    if not verified:
        return json_api_error(message, 400)

    pwd_err = validate_password(new_pwd)
    if pwd_err:
        return json_api_error(pwd_err, 400)

    user.password_hash = hash_password(new_pwd)
    db.session.commit()
    return json_api_response({"status": "success", "message": "密码已修改"})

@app.route('/api/auth/send-code', methods=['POST'])
def api_send_code():
    """
    发送验证码
    
    POST /api/auth/send-code
    请求体：{ target, purpose }
    返回：{ status, message, channel, test_code }
    
    功能说明：向指定邮箱/手机号发送验证码
    """
    data = request.get_json() or {}
    raw_purpose = data.get('purpose', 'login')
    if raw_purpose == 'change_password':
        user_id = session.get('user_id')
        if not user_id:
            return json_api_error("未登录", 401)
        user = db.session.get(User, user_id)
        if not user:
            return json_api_error("无效的登录状态", 401)
        target = normalize_auth_target(data.get('target'))
        if target != normalize_auth_target(user.email):
            return json_api_error("只能向当前登录邮箱发送验证码", 403)

    target, purpose, error_message, status = validate_verification_request(
        data.get('target'),
        raw_purpose
    )
    if error_message:
        return json_api_error(error_message, status)

    success, message, channel, test_code = send_verification_code(target, purpose)
    if not success:
        error_status = 429 if '频繁' in (message or '') else 400
        return json_api_error(message or "验证码发送失败", error_status)

    payload = {
        "status": "success",
        "message": message,
        "channel": channel,
        "cooldown_seconds": VERIFICATION_COOLDOWN
    }
    if test_code:
        payload["test_code"] = test_code
    return json_api_response(payload)

@app.route('/api/auth/verify-code', methods=['POST'])
def api_verify_code():
    """
    验证验证码
    
    POST /api/auth/verify-code
    请求体：{ target, code, purpose }
    返回：{ status, message }
    
    功能说明：验证验证码是否有效
    """
    data = request.get_json() or {}
    target = data.get('target')
    code = data.get('code')
    purpose = data.get('purpose', 'login')

    verified, message = verify_verification_code(target, code, purpose)
    if not verified:
        return json_api_error(message, 400)
    return json_api_response({"status": "success", "message": message})

@app.route('/api/auth/login-with-code', methods=['POST'])
def api_login_with_code():
    """
    验证码登录
    
    POST /api/auth/login-with-code
    请求体：{ target, code }
    返回：{ status, user, message }
    
    功能说明：使用验证码登录（首次登录自动创建账号）
    """
    data = request.get_json() or {}
    target = normalize_auth_target(data.get('target'))
    code = data.get('code')

    verified, message = verify_verification_code(target, code, 'login')
    if not verified:
        return json_api_error(message, 400)

    user = find_user_by_email(target)
    created = False
    if not user:
        user = create_code_login_user(target)
        created = True

    apply_login_session(user)
    response_message = "验证码登录成功"
    if created:
        response_message = "验证码登录成功，已为你创建账号"
    return json_api_response({"status": "success", "user": user.to_dict(), "message": response_message})

@app.route('/api/auth/register-with-code', methods=['POST'])
def api_register_with_code():
    """
    验证码注册
    
    POST /api/auth/register-with-code
    请求体：{ target, code, username, password }
    返回：{ status, user }
    
    功能说明：使用验证码验证后注册账号
    """
    data = request.get_json() or {}
    target, purpose, error_message, status = validate_verification_request(
        data.get('target'),
        'register'
    )
    if error_message:
        return json_api_error(error_message, status)

    username = (data.get('username') or '').strip()
    password = data.get('password', '')
    code = data.get('code')

    if User.query.filter_by(username=username).first():
        return json_api_error("用户名已存在", 400)

    verified, message = verify_verification_code(target, code, purpose)
    if not verified:
        return json_api_error(message, 400)

    user, register_error, register_status = register_user_account(username, target, password)
    if register_error:
        return json_api_error(register_error, register_status)

    apply_login_session(user, update_stats=False)
    return json_api_response({"status": "success", "user": user.to_dict()}, status=201)

@app.route('/api/auth/reset-password', methods=['POST'])
def api_reset_password():
    """
    重置密码
    
    POST /api/auth/reset-password
    请求体：{ target, code, new_password }
    返回：{ status, message }
    
    功能说明：使用验证码重置密码
    """
    data = request.get_json() or {}
    target, _, error_message, status = validate_verification_request(
        data.get('target'),
        'reset_password'
    )
    if error_message:
        return json_api_error(error_message, status)

    code = data.get('code')
    new_password = data.get('new_password', '')
    pwd_err = validate_password(new_password)
    if pwd_err:
        return json_api_error(pwd_err, 400)

    verified, message = verify_verification_code(target, code, 'reset_password')
    if not verified:
        return json_api_error(message, 400)

    user = find_user_by_email(target)
    if not user:
        return json_api_error("该邮箱尚未注册", 404)

    user.password_hash = hash_password(new_password)
    db.session.commit()
    return json_api_response({"status": "success", "message": "密码已重置"})

@app.route('/api/current-user', methods=['GET'])
def get_current_user():
    """
    获取当前登录用户信息（别名）
    """
    return get_user()

@app.route('/api/occasion-types', methods=['GET'])
def get_occasion_types():
    return app.response_class(
        response=json.dumps([
            {"value": "focus_group", "label": "焦点小组讨论"},
            {"value": "user_interview", "label": "用户深度访谈"},
            {"value": "brainstorming", "label": "头脑风暴会议"},
            {"value": "product_team", "label": "产品团队评审"},
            {"value": "sales_conversation", "label": "销售对话模拟"}
        ], ensure_ascii=False),
        status=200,
        mimetype='application/json; charset=utf-8'
    )

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
        ensure_extended_schema()
    app.run(debug=True, host='0.0.0.0', port=5000)
